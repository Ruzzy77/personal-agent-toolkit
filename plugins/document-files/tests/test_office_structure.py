from __future__ import annotations

import tempfile
import unittest
from pathlib import Path

from document_files.extraction_protocol import run_builtin_extraction


class OfficeStructureTest(unittest.TestCase):
    def test_word_body_order_nested_tables_and_merged_cells(self):
        from docx import Document

        with tempfile.TemporaryDirectory() as temporary:
            path = Path(temporary) / "source.docx"
            document = Document()
            document.add_heading("Title", level=1)
            document.add_paragraph("Before the table")
            table = document.add_table(rows=3, cols=2)
            table.cell(0, 0).merge(table.cell(0, 1)).text = "Horizontal merge"
            table.cell(1, 0).merge(table.cell(2, 0)).text = "Vertical merge"
            table.cell(1, 1).text = "Before nested"
            nested = table.cell(1, 1).add_table(rows=1, cols=1)
            nested.cell(0, 0).text = "Nested cell"
            table.cell(1, 1).add_paragraph("After nested")
            table.cell(2, 1).text = "Last cell"
            document.add_paragraph("After the table")
            document.sections[0].header.paragraphs[0].text = "Header text"
            document.sections[0].footer.paragraphs[0].text = "Footer text"
            document.save(path)
            original = path.read_bytes()
            result = run_builtin_extraction(path, "docx")
            self.assertEqual(original, path.read_bytes())
        texts = [unit.content for unit in result.units if unit.content]
        self.assertEqual(
            texts,
            [
                "Title",
                "Before the table",
                "Horizontal merge",
                "Vertical merge",
                "Before nested",
                "Nested cell",
                "After nested",
                "Last cell",
                "After the table",
                "Header text",
                "Footer text",
            ],
        )
        vertical = next(unit for unit in result.units if unit.content == "Vertical merge")
        horizontal = next(unit for unit in result.units if unit.content == "Horizontal merge")
        nested = next(unit for unit in result.units if unit.content == "Nested cell")
        self.assertEqual(vertical.structure_path["row_span"], 2)
        self.assertEqual(horizontal.structure_path["col_span"], 2)
        self.assertTrue(nested.structure_path["container_path"])
        self.assertEqual(result.units[0].unit_type, "heading")

    def test_word_textbox_and_field_do_not_duplicate_or_evaluate(self):
        from docx import Document
        from docx.oxml import parse_xml

        with tempfile.TemporaryDirectory() as temporary:
            path = Path(temporary) / "source.docx"
            document = Document()
            paragraph = document.add_paragraph("Before ")
            paragraph._p.append(
                parse_xml("""
              <w:r xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"
                   xmlns:mc="http://schemas.openxmlformats.org/markup-compatibility/2006">
                <mc:AlternateContent>
                  <mc:Choice Requires="w"><w:drawing><w:txbxContent><w:p><w:r>
                    <w:t>Textbox</w:t></w:r></w:p></w:txbxContent></w:drawing></mc:Choice>
                  <mc:Fallback><w:pict><w:txbxContent><w:p><w:r>
                    <w:t>Textbox</w:t></w:r></w:p></w:txbxContent></w:pict></mc:Fallback>
                </mc:AlternateContent>
              </w:r>""")
            )
            paragraph.add_run("After ")
            paragraph._p.append(
                parse_xml("""
              <w:fldSimple xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"
                           w:instr="DATE"><w:r><w:t>Stored date</w:t></w:r></w:fldSimple>""")
            )
            document.save(path)
            result = run_builtin_extraction(path, "docx")
        self.assertEqual(
            [unit.content for unit in result.units if unit.content],
            ["Before", "Textbox", "After", "Stored date"],
        )
        field = next(unit for unit in result.units if unit.unit_type == "field")
        self.assertEqual(field.structure_path["instruction"], "DATE")
        self.assertEqual(field.structure_path["evaluation"], "stored_result_only")

    def test_slides_recursive_shapes_tables_chart_and_notes(self):
        from pptx import Presentation
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE
        from pptx.util import Inches

        with tempfile.TemporaryDirectory() as temporary:
            path = Path(temporary) / "source.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            slide.shapes.add_textbox(0, 0, Inches(2), Inches(1)).text = "Top-level"
            group = slide.shapes.add_group_shape()
            group.shapes.add_textbox(0, 0, Inches(2), Inches(1)).text = "Group text"
            nested = group.shapes.add_group_shape()
            nested.shapes.add_textbox(0, 0, Inches(2), Inches(1)).text = "Nested group"
            table = slide.shapes.add_table(2, 2, 0, 0, Inches(4), Inches(2)).table
            table.cell(0, 0).merge(table.cell(0, 1))
            table.cell(0, 0).text = "Merged heading"
            table.cell(1, 0).text = "Cell A"
            table.cell(1, 1).text = "Cell B"
            data = CategoryChartData()
            data.categories = ["North", "South"]
            data.add_series("Revenue", [12, 18])
            slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, 0, 0, Inches(4), Inches(2), data)
            slide.notes_slide.notes_text_frame.text = "Speaker notes"
            presentation.save(path)
            first = run_builtin_extraction(path, "pptx")
            second = run_builtin_extraction(path, "pptx")
        self.assertEqual(first.manifest_hash, second.manifest_hash)
        texts = [unit.content for unit in first.units if unit.content]
        self.assertEqual(
            texts[:6],
            [
                "Top-level",
                "Group text",
                "Nested group",
                "Merged heading",
                "Cell A",
                "Cell B",
            ],
        )
        self.assertEqual(texts[-1], "Speaker notes")
        self.assertTrue({"Revenue", "North", "South", "12", "18"}.issubset(texts))
        cell = next(unit for unit in first.units if unit.content == "Merged heading")
        self.assertEqual(cell.structure_path["col_span"], 2)


if __name__ == "__main__":
    unittest.main()
