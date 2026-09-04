"""Small reopen checks for the four bundled document-authoring Skills."""

from pathlib import Path

from docx import Document
from openpyxl import Workbook, load_workbook
from pptx import Presentation
from pypdf import PdfReader
from reportlab.pdfgen.canvas import Canvas


def test_docx_xlsx_pptx_and_pdf_reopen(tmp_path: Path) -> None:
    docx_path = tmp_path / "sample.docx"
    document = Document()
    document.add_paragraph("한국어 문서")
    document.save(docx_path)
    assert Document(docx_path).paragraphs[0].text == "한국어 문서"

    xlsx_path = tmp_path / "sample.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet["A1"] = "값"
    sheet["B1"] = "=1+1"
    workbook.save(xlsx_path)
    reopened = load_workbook(xlsx_path, data_only=False)
    assert (reopened.active["A1"].value, reopened.active["B1"].value) == ("값", "=1+1")

    pptx_path = tmp_path / "sample.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "한국어 발표"
    presentation.save(pptx_path)
    assert Presentation(pptx_path).slides[0].shapes.title.text == "한국어 발표"

    pdf_path = tmp_path / "sample.pdf"
    canvas = Canvas(str(pdf_path))
    canvas.drawString(72, 720, "PDF document")
    canvas.save()
    assert len(PdfReader(pdf_path).pages) == 1
