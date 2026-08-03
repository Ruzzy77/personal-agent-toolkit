"""Application service used by the CLI and MCP surfaces."""

from __future__ import annotations

import errno
import hashlib
import hmac
import json
import math
import os
import stat
import time
import unicodedata
import uuid
from contextlib import suppress
from pathlib import Path

from .adapter_registry import AdapterRegistry, build_default_registry
from .adapters import AdapterDescriptor, ExtractionEnvelope
from .capture import (
    COPY_CHUNK_BYTES,
    CapturedSource,
    capture_to_staging,
    cleanup_abandoned_staging,
    cleanup_source_copies,
    current_source_path_identity,
    discard_staged_capture,
    native_source_path,
    observe_staging,
    source_identity_from_stat,
    validate_file_boundary,
)
from .config import (
    RuntimePaths,
    default_data_root,
    is_within,
    normalize_corpus_id,
)
from .contexts import CONTEXT_MAX_LIMIT, ContextService
from .database import (
    configure_corpus_source_scope,
    corpus_connection,
    corpus_read_connection,
    corpus_schema_status,
    encode_json,
    get_corpus,
    list_corpora,
    migrate_corpus,
    register_corpus,
    utc_now,
)
from .errors import (
    BudgetExceededError,
    ConfigurationError,
    CorpusError,
    ExtractionError,
    InvalidRequestError,
    SemanticCommitError,
    SnapshotConflictError,
    SourceBoundaryError,
    SourceUnavailableError,
)
from .golden import (
    GoldenEvaluationError,
    evaluate_projection_observation,
    projection_observation_sha256,
    validate_golden_annotation,
)
from .locking import writer_lock
from .scanner import scan_corpus
from .schema import EXTRACTION_SCHEMA_VERSION
from .session_sources import SESSION_SOURCE_FETCH_DEFAULT_CHARS
from .source_access import (
    opened_source_file,
    opened_source_root,
    relative_source_parts,
    source_root_identity,
)

_MATERIALIZATION_RECEIPT_KEY_NAME = "materialization-receipt-hmac.key"
_MATERIALIZATION_RECEIPT_KEY_BYTES = 32
# These domains are persisted identity/proof inputs. Renaming them would break
# existing source links and receipts, so they remain stable across branding.
_MATERIALIZATION_RECEIPT_DOMAIN = b"work-corpus-materialization-receipt-v1\0"
_MATERIALIZATION_RECEIPT_INSTALL_RETRIES = 100
_MATERIALIZATION_RECEIPT_INSTALL_RETRY_SECONDS = 0.005
_MAX_SEARCH_RESULTS = 200
CORPUS_INVENTORY_DEFAULT_LIMIT = 100
CORPUS_INVENTORY_MAX_LIMIT = 200
CORPUS_INVENTORY_MAX_OFFSET = 100_000
CORPUS_INVENTORY_MAX_PATH_FILTER_CHARS = 1_000
CORPUS_INVENTORY_MAX_EXTENSION_CHARS = 20
CORPUS_INVENTORY_MAX_LOGICAL_BYTES = (1 << 63) - 1
CORPUS_INVENTORY_MAX_SERIALIZED_BYTES = 1024 * 1024
CORPUS_SEARCH_EXCERPT_MAX_CHARS = 2_000
CORPUS_SEARCH_EXCERPT_CONTEXT_BEFORE_CHARS = 400
CORPUS_SEARCH_MAX_SERIALIZED_BYTES = 1024 * 1024
_MAX_READ_UNITS = 200
_MAX_NEIGHBOR_SPAN = 10
CORPUS_READ_MIN_CHARS = 1_000
CORPUS_READ_DEFAULT_CHARS = 30_000
CORPUS_READ_MAX_CHARS = 200_000
CORPUS_READ_MAX_SELECTED_UNITS = 500
CORPUS_READ_MAX_SERIALIZED_BYTES = 2 * 1024 * 1024
CORPUS_OVERVIEW_DEFAULT_ITEMS_PER_CONTEXT = 5
CORPUS_OVERVIEW_MAX_ITEMS_PER_CONTEXT = 20
CORPUS_OVERVIEW_MAX_BODY_CHARS = 1_200
CORPUS_OVERVIEW_MAX_SERIALIZED_BYTES = 512 * 1024
INTERPRETATION_MATERIAL_MAX_SERIALIZED_BYTES = 2 * 1024 * 1024
SEMANTIC_CONTEXT_MAX_SERIALIZED_BYTES = 2 * 1024 * 1024
_MAX_INGEST_FILES = 50
_MAX_INGEST_BYTES = 500 * 1024 * 1024
_MAX_INGEST_FILE_BYTES = 250 * 1024 * 1024
_MAX_INGEST_DOCUMENT_IDS = 100
_MAX_INGEST_TIMEOUT_SECONDS = 600
SEMANTIC_COMMIT_MAX_CLAIMS = 50
SEMANTIC_COMMIT_MAX_EVIDENCE_PER_CLAIM = 50
SEMANTIC_COMMIT_MAX_EVIDENCE_LINKS = 500
SEMANTIC_COMMIT_MAX_COMPLETED_REVISIONS = 50
SEMANTIC_COMMIT_MAX_PROGRESS_UPDATES = 50
SEMANTIC_COMMIT_MAX_IDENTIFIER_CHARS = 200
SEMANTIC_COMMIT_MAX_BODY_CHARS = 10_000
SEMANTIC_COMMIT_MAX_SUBJECT_CHARS = 2_000
SEMANTIC_COMMIT_MAX_QUALIFIER_CHARS = 2_000
SEMANTIC_COMMIT_MAX_STATUS_CHARS = 1_000
SEMANTIC_COMMIT_MAX_MATERIALIZER_VERSION_CHARS = 200
SEMANTIC_COMMIT_MAX_COLLECTION_ITEMS = 100
SEMANTIC_COMMIT_MAX_NESTING_DEPTH = 8
SEMANTIC_COMMIT_MAX_NESTED_VALUES = 10_000
SEMANTIC_COMMIT_MAX_SERIALIZED_BYTES = 1024 * 1024
SEMANTIC_COMMIT_MAX_INTEGER = (1 << 63) - 1

_INVENTORY_ELIGIBILITY_STATES = {
    "all",
    "supported",
    "unsupported",
    "ignored",
}
_INVENTORY_RESIDENCY_STATES = {
    "all",
    "resident",
    "remote_only",
    "unknown",
}
_INVENTORY_INDEX_STATES = {
    "all",
    "current",
    "refresh_required",
    "unindexed",
    "not_applicable",
}


def _validate_ingest_budgets(
    *,
    max_files: int,
    max_bytes: int,
    max_file_bytes: int,
    timeout_seconds: float,
) -> None:
    if (
        not 1 <= max_files <= _MAX_INGEST_FILES
        or not 1 <= max_bytes <= _MAX_INGEST_BYTES
        or not 1 <= max_file_bytes <= _MAX_INGEST_FILE_BYTES
        or not 0 < timeout_seconds <= _MAX_INGEST_TIMEOUT_SECONDS
    ):
        raise BudgetExceededError(
            "ingest budgets exceed the supported request bounds",
            details={
                "max_files": max_files,
                "max_bytes": max_bytes,
                "max_file_bytes": max_file_bytes,
                "timeout_seconds": timeout_seconds,
                "allowed": {
                    "max_files": [1, _MAX_INGEST_FILES],
                    "max_bytes": [1, _MAX_INGEST_BYTES],
                    "max_file_bytes": [1, _MAX_INGEST_FILE_BYTES],
                    "timeout_seconds": [">0", _MAX_INGEST_TIMEOUT_SECONDS],
                },
            },
        )


def _ephemeral_capture_ref(sha256: str) -> str:
    return f"ephemeral:sha256:{sha256}"


def _source_is_missing(error: SourceBoundaryError) -> bool:
    reason = error.details.get("reason")
    return isinstance(reason, str) and reason.endswith(f":{errno.ENOENT}")


def _canonical_legacy_blob_ref(sha256: str) -> str:
    return f"blobs/{sha256[:2]}/{sha256}.blob"


def _safe_relative_inventory_path(value: object) -> bool:
    if not isinstance(value, str) or not value or "\x00" in value:
        return False
    if value.startswith(("/", "\\")):
        return False
    if (
        len(value) >= 3
        and value[0].isalpha()
        and value[1] == ":"
        and value[2] in {"/", "\\"}
    ):
        return False
    return all(
        part not in {"", ".", ".."}
        for part in value.replace("\\", "/").split("/")
    )


def _normalize_inventory_path_filter(value: str | None) -> str | None:
    if value is None:
        return None
    normalized = unicodedata.normalize("NFC", value.strip())
    if not normalized:
        return None
    if len(normalized) > CORPUS_INVENTORY_MAX_PATH_FILTER_CHARS:
        raise BudgetExceededError(
            "inventory path filter is too long",
            details={
                "path_filter_chars": len(normalized),
                "maximum": CORPUS_INVENTORY_MAX_PATH_FILTER_CHARS,
            },
        )
    if not _safe_relative_inventory_path(normalized):
        raise ConfigurationError(
            "inventory path filter must be a safe relative literal",
        )
    return normalized


def _normalize_inventory_extension(value: str | None) -> str | None:
    if value is None:
        return None
    normalized = value.strip().lower().removeprefix(".")
    if not normalized:
        return None
    if (
        len(normalized) > CORPUS_INVENTORY_MAX_EXTENSION_CHARS
        or not all(character.isascii() and character.isalnum() for character in normalized)
    ):
        raise ConfigurationError(
            "inventory extension must be a lowercase alphanumeric suffix",
            details={
                "maximum_extension_chars": CORPUS_INVENTORY_MAX_EXTENSION_CHARS,
            },
        )
    return normalized


def _revision_id(document_id: str, sha256: str) -> str:
    value = hashlib.sha256(
        f"work-corpus-revision-v1\0{document_id}\0{sha256}".encode()
    ).hexdigest()
    return f"rev_{value[:32]}"


def _unit_id(projection_id: str, ordinal: int, content_sha256: str) -> str:
    value = hashlib.sha256(
        f"work-corpus-unit-v2\0{projection_id}\0{ordinal}\0{content_sha256}".encode()
    ).hexdigest()
    return f"unit_{value[:32]}"


def _projection_id(
    revision_id: str,
    adapter_id: str,
    adapter_version: str,
    config_hash: str,
    result_manifest_hash: str,
) -> str:
    value = hashlib.sha256(
        (
            "work-corpus-projection-v1\0"
            f"{revision_id}\0{adapter_id}\0{adapter_version}\0"
            f"{config_hash}\0{result_manifest_hash}"
        ).encode()
    ).hexdigest()
    return f"projection_{value[:32]}"


def _issue_locator(issue: dict) -> tuple[dict, str]:
    details = issue.get("details", {})
    structural_locator = issue.get("structural_locator")
    if not isinstance(structural_locator, dict) and isinstance(details, dict):
        structural_locator = details.get("structural_locator")
    if not isinstance(structural_locator, dict):
        structural_locator = {
            key: details.get(key, issue.get(key))
            for key in ("page", "slide", "sheet", "range", "section", "paragraph")
            if key in issue or (isinstance(details, dict) and key in details)
        }
    locator_payload = {
        "code": issue.get("code", "extractor_issue"),
        "structural_locator": structural_locator,
    }
    locator_key = hashlib.sha256(encode_json(locator_payload).encode()).hexdigest()
    return structural_locator, locator_key


def _source_span(structure: dict) -> dict | None:
    if "line_start" in structure:
        return {
            "line_start": structure["line_start"],
            "line_end": structure.get("line_end", structure["line_start"]),
        }
    if "page" in structure:
        return {"page": structure["page"]}
    if "slide" in structure:
        return {"slide": structure["slide"]}
    if "sheet" in structure:
        return {"sheet": structure["sheet"], "range": structure.get("range")}
    if "section" in structure:
        return {
            "section": structure["section"],
            "paragraph": structure.get("paragraph"),
        }
    return None


def _semantic_state(connection) -> dict:
    state_hash = hashlib.sha256()
    state_hash.update(b'{"claims":[')
    first_row = True
    for row in connection.execute(
        """
        SELECT c.claim_id, c.dependency_state, c.validation_state,
               c.claim_assessment, e.evidence_link_id, e.source_unit_id,
               e.stance
        FROM atomic_claims c
        LEFT JOIN evidence_links e ON e.claim_id = c.claim_id
        ORDER BY c.claim_id, e.evidence_link_id
        """
    ):
        if not first_row:
            state_hash.update(b",")
        state_hash.update(encode_json(dict(row)).encode())
        first_row = False
    state_hash.update(b'],"commits":[')
    commit_count = 0
    for row in connection.execute(
        """
        SELECT commit_id, base_snapshot_id, input_sha256, materializer_version
        FROM semantic_commits ORDER BY commit_id
        """
    ):
        if commit_count:
            state_hash.update(b",")
        state_hash.update(encode_json(dict(row)).encode())
        commit_count += 1
    state_hash.update(b"]}")
    claim_count = connection.execute("SELECT COUNT(*) FROM atomic_claims").fetchone()[0]
    valid_claim_count = connection.execute(
        "SELECT COUNT(*) FROM atomic_claims WHERE dependency_state = 'valid'"
    ).fetchone()[0]
    return {
        "semantic_state_hash": state_hash.hexdigest(),
        "commit_count": commit_count,
        "claim_count": claim_count,
        "valid_claim_count": valid_claim_count,
    }


def _receipt_key_error(key_path: Path, reason: str) -> ConfigurationError:
    return ConfigurationError(
        "materialization receipt key is missing, unsafe, or invalid",
        details={
            "key_path_sha256": hashlib.sha256(str(key_path).encode()).hexdigest(),
            "reason": reason,
        },
    )


def _validate_receipt_key_stat(key_path: Path, key_stat: os.stat_result) -> None:
    if not stat.S_ISREG(key_stat.st_mode):
        raise _receipt_key_error(key_path, "not_regular")
    if key_stat.st_uid != os.geteuid():
        raise _receipt_key_error(key_path, "wrong_owner")
    if stat.S_IMODE(key_stat.st_mode) != 0o600:
        raise _receipt_key_error(key_path, "unsafe_permissions")
    if key_stat.st_nlink != 1:
        reason = (
            "installation_in_progress"
            if key_stat.st_nlink == 2
            else "unexpected_link_count"
        )
        raise _receipt_key_error(key_path, reason)
    if key_stat.st_size != _MATERIALIZATION_RECEIPT_KEY_BYTES:
        raise _receipt_key_error(key_path, "invalid_length")


def _read_receipt_key(
    runtime_descriptor: int,
    *,
    key_path: Path,
) -> bytes | None:
    try:
        path_stat = os.stat(
            _MATERIALIZATION_RECEIPT_KEY_NAME,
            dir_fd=runtime_descriptor,
            follow_symlinks=False,
        )
    except FileNotFoundError:
        return None
    except OSError as exc:
        raise _receipt_key_error(key_path, f"stat_failed:{exc.errno}") from exc
    _validate_receipt_key_stat(key_path, path_stat)

    flags = os.O_RDONLY | getattr(os, "O_CLOEXEC", 0) | getattr(os, "O_NOFOLLOW", 0)
    try:
        key_descriptor = os.open(
            _MATERIALIZATION_RECEIPT_KEY_NAME,
            flags,
            dir_fd=runtime_descriptor,
        )
    except OSError as exc:
        raise _receipt_key_error(key_path, f"open_failed:{exc.errno}") from exc
    try:
        opened_stat = os.fstat(key_descriptor)
        _validate_receipt_key_stat(key_path, opened_stat)
        if (
            opened_stat.st_dev != path_stat.st_dev
            or opened_stat.st_ino != path_stat.st_ino
        ):
            raise _receipt_key_error(key_path, "path_changed_during_open")
        key = os.read(key_descriptor, _MATERIALIZATION_RECEIPT_KEY_BYTES + 1)
        closed_stat = os.fstat(key_descriptor)
        try:
            final_path_stat = os.stat(
                _MATERIALIZATION_RECEIPT_KEY_NAME,
                dir_fd=runtime_descriptor,
                follow_symlinks=False,
            )
        except OSError as exc:
            raise _receipt_key_error(
                key_path,
                f"final_stat_failed:{exc.errno}",
            ) from exc
        _validate_receipt_key_stat(key_path, final_path_stat)
        if (
            closed_stat.st_dev != opened_stat.st_dev
            or closed_stat.st_ino != opened_stat.st_ino
            or closed_stat.st_size != opened_stat.st_size
            or closed_stat.st_mtime_ns != opened_stat.st_mtime_ns
            or closed_stat.st_ctime_ns != opened_stat.st_ctime_ns
            or final_path_stat.st_dev != opened_stat.st_dev
            or final_path_stat.st_ino != opened_stat.st_ino
        ):
            raise _receipt_key_error(key_path, "key_changed_during_read")
    finally:
        os.close(key_descriptor)
    if len(key) != _MATERIALIZATION_RECEIPT_KEY_BYTES:
        raise _receipt_key_error(key_path, "invalid_length")
    return key


def _read_receipt_key_with_install_retry(
    runtime_descriptor: int,
    *,
    key_path: Path,
) -> bytes | None:
    for attempt in range(_MATERIALIZATION_RECEIPT_INSTALL_RETRIES):
        try:
            return _read_receipt_key(runtime_descriptor, key_path=key_path)
        except ConfigurationError as exc:
            if (
                exc.details.get("reason") != "installation_in_progress"
                or attempt == _MATERIALIZATION_RECEIPT_INSTALL_RETRIES - 1
            ):
                raise
            time.sleep(_MATERIALIZATION_RECEIPT_INSTALL_RETRY_SECONDS)
    raise AssertionError("unreachable receipt-key retry state")


def _materialization_receipt_key(data_root: Path, *, create: bool) -> bytes:
    runtime_path = data_root / "runtime"
    key_path = runtime_path / _MATERIALIZATION_RECEIPT_KEY_NAME
    try:
        runtime_path.mkdir(parents=True, exist_ok=True, mode=0o700)
        runtime_stat = os.lstat(runtime_path)
    except OSError as exc:
        raise _receipt_key_error(key_path, f"runtime_unavailable:{exc.errno}") from exc
    if not stat.S_ISDIR(runtime_stat.st_mode):
        raise _receipt_key_error(key_path, "runtime_not_directory")
    if runtime_stat.st_uid != os.geteuid():
        raise _receipt_key_error(key_path, "runtime_wrong_owner")
    if stat.S_IMODE(runtime_stat.st_mode) & 0o077:
        raise _receipt_key_error(key_path, "runtime_unsafe_permissions")

    runtime_flags = (
        os.O_RDONLY
        | getattr(os, "O_CLOEXEC", 0)
        | getattr(os, "O_DIRECTORY", 0)
        | getattr(os, "O_NOFOLLOW", 0)
    )
    try:
        runtime_descriptor = os.open(runtime_path, runtime_flags)
    except OSError as exc:
        raise _receipt_key_error(key_path, f"runtime_open_failed:{exc.errno}") from exc
    try:
        opened_runtime_stat = os.fstat(runtime_descriptor)
        if (
            not stat.S_ISDIR(opened_runtime_stat.st_mode)
            or opened_runtime_stat.st_dev != runtime_stat.st_dev
            or opened_runtime_stat.st_ino != runtime_stat.st_ino
        ):
            raise _receipt_key_error(key_path, "runtime_changed_during_open")
        existing_key = _read_receipt_key_with_install_retry(
            runtime_descriptor,
            key_path=key_path,
        )
        if existing_key is not None:
            return existing_key
        if not create:
            raise _receipt_key_error(key_path, "missing")

        key = os.urandom(_MATERIALIZATION_RECEIPT_KEY_BYTES)
        temporary_name = (
            f".{_MATERIALIZATION_RECEIPT_KEY_NAME}.{uuid.uuid4().hex}.tmp"
        )
        temporary_descriptor: int | None = None
        try:
            temporary_descriptor = os.open(
                temporary_name,
                os.O_WRONLY
                | os.O_CREAT
                | os.O_EXCL
                | getattr(os, "O_CLOEXEC", 0)
                | getattr(os, "O_NOFOLLOW", 0),
                0o600,
                dir_fd=runtime_descriptor,
            )
            os.fchmod(temporary_descriptor, 0o600)
            offset = 0
            while offset < len(key):
                written = os.write(temporary_descriptor, key[offset:])
                if written <= 0:
                    raise OSError("receipt key write made no progress")
                offset += written
            os.fsync(temporary_descriptor)
            temporary_stat = os.fstat(temporary_descriptor)
            if (
                not stat.S_ISREG(temporary_stat.st_mode)
                or stat.S_IMODE(temporary_stat.st_mode) != 0o600
                or temporary_stat.st_size != _MATERIALIZATION_RECEIPT_KEY_BYTES
            ):
                raise _receipt_key_error(key_path, "temporary_key_invalid")
            os.close(temporary_descriptor)
            temporary_descriptor = None
            with suppress(FileExistsError):
                os.link(
                    temporary_name,
                    _MATERIALIZATION_RECEIPT_KEY_NAME,
                    src_dir_fd=runtime_descriptor,
                    dst_dir_fd=runtime_descriptor,
                    follow_symlinks=False,
                )
        except ConfigurationError:
            raise
        except OSError as exc:
            raise _receipt_key_error(key_path, f"creation_failed:{exc.errno}") from exc
        finally:
            if temporary_descriptor is not None:
                os.close(temporary_descriptor)
            try:
                os.unlink(temporary_name, dir_fd=runtime_descriptor)
            except FileNotFoundError:
                pass
            except OSError as exc:
                raise _receipt_key_error(
                    key_path,
                    f"temporary_cleanup_failed:{exc.errno}",
                ) from exc
        try:
            os.fsync(runtime_descriptor)
        except OSError as exc:
            raise _receipt_key_error(key_path, f"runtime_sync_failed:{exc.errno}") from exc

        installed_key = _read_receipt_key_with_install_retry(
            runtime_descriptor,
            key_path=key_path,
        )
        if installed_key is None:
            raise _receipt_key_error(key_path, "atomic_install_failed")
        return installed_key
    finally:
        os.close(runtime_descriptor)


def _materialization_batch_receipt(
    *,
    receipt_key: bytes,
    base_snapshot_id: str,
    base_semantic_state_hash: str,
    queue_id: str,
    revision_id: str,
    projection_id: str,
    processed_from_ordinal: int,
    processed_through_ordinal: int,
    units: list[dict],
) -> str:
    manifest = {
        "schema_version": 1,
        "base_snapshot_id": base_snapshot_id,
        "base_semantic_state_hash": base_semantic_state_hash,
        "queue_id": queue_id,
        "revision_id": revision_id,
        "projection_id": projection_id,
        "processed_from_ordinal": processed_from_ordinal,
        "processed_through_ordinal": processed_through_ordinal,
        "units": [
            {
                "unit_id": unit["unit_id"],
                "ordinal": int(unit["ordinal"]),
                "content_sha256": unit["content_sha256"],
            }
            for unit in units
        ],
    }
    return hmac.new(
        receipt_key,
        _MATERIALIZATION_RECEIPT_DOMAIN + encode_json(manifest).encode(),
        hashlib.sha256,
    ).hexdigest()


def _resident_source_matches_revision(
    *,
    source_path: Path,
    source_root: Path,
    staging_root: Path,
    document: dict,
    revision_sha256: str,
) -> bool:
    try:
        validate_file_boundary(source_path, source_root, staging_root)
        relative_parts = relative_source_parts(source_path, source_root)
        with (
            opened_source_root(source_root) as root_descriptor,
            opened_source_file(root_descriptor, relative_parts) as descriptor,
        ):
            initial_source_root_identity = source_root_identity(root_descriptor)
            before = source_identity_from_stat(os.fstat(descriptor))
            if before.dataless or document["residency_state"] != "resident":
                return False
            if before.stable_key() != (
                document["logical_size"],
                document["modified_ns"],
                document["changed_ns"],
                document["device"],
                document["inode"],
            ):
                return False
            digest = hashlib.sha256()
            os.lseek(descriptor, 0, os.SEEK_SET)
            while chunk := os.read(descriptor, COPY_CHUNK_BYTES):
                digest.update(chunk)
            after = source_identity_from_stat(os.fstat(descriptor))
            current_path = current_source_path_identity(
                source_root,
                initial_source_root_identity,
                relative_parts,
            )
    except (OSError, CorpusError):
        return False
    return bool(
        before.stable_key() == after.stable_key()
        and current_path.stable_key() == after.stable_key()
        and digest.hexdigest() == revision_sha256
    )


def _require_semantic_commit_string(
    value: object,
    *,
    field: str,
    maximum: int,
    nonempty: bool = False,
) -> str:
    if not isinstance(value, str):
        raise SemanticCommitError(
            f"{field} must be a string",
            details={"field": field},
        )
    if len(value) > maximum:
        raise SemanticCommitError(
            f"{field} may contain at most {maximum} characters",
            details={"field": field, "maximum": maximum, "actual": len(value)},
        )
    if nonempty and not value.strip():
        raise SemanticCommitError(
            f"{field} must be non-empty",
            details={"field": field},
        )
    return value


def _validate_semantic_commit_json_bounds(value: object) -> None:
    state = {
        "values": 0,
        "string_bytes": 0,
    }
    active_containers: set[int] = set()

    def visit(item: object, *, depth: int, path: str) -> None:
        if depth > SEMANTIC_COMMIT_MAX_NESTING_DEPTH:
            raise SemanticCommitError(
                "semantic commit JSON nesting exceeds the request limit",
                details={
                    "field": path,
                    "maximum_depth": SEMANTIC_COMMIT_MAX_NESTING_DEPTH,
                },
            )
        state["values"] += 1
        if state["values"] > SEMANTIC_COMMIT_MAX_NESTED_VALUES:
            raise SemanticCommitError(
                "semantic commit contains too many nested values",
                details={
                    "maximum": SEMANTIC_COMMIT_MAX_NESTED_VALUES,
                },
            )
        if isinstance(item, str):
            if len(item) > SEMANTIC_COMMIT_MAX_BODY_CHARS:
                raise SemanticCommitError(
                    "semantic commit contains an oversized nested string",
                    details={
                        "field": path,
                        "maximum": SEMANTIC_COMMIT_MAX_BODY_CHARS,
                        "actual": len(item),
                    },
                )
            state["string_bytes"] += len(item.encode("utf-8"))
            if state["string_bytes"] > SEMANTIC_COMMIT_MAX_SERIALIZED_BYTES:
                raise SemanticCommitError(
                    "semantic commit string content exceeds the aggregate request budget",
                    details={
                        "maximum_bytes": SEMANTIC_COMMIT_MAX_SERIALIZED_BYTES,
                    },
                )
            return
        if item is None or isinstance(item, bool):
            return
        if isinstance(item, int):
            if abs(item) > SEMANTIC_COMMIT_MAX_INTEGER:
                raise SemanticCommitError(
                    "semantic commit integer exceeds the request limit",
                    details={
                        "field": path,
                        "maximum": SEMANTIC_COMMIT_MAX_INTEGER,
                    },
                )
            return
        if isinstance(item, float):
            if not math.isfinite(item):
                raise SemanticCommitError(
                    "semantic commit numbers must be finite",
                    details={"field": path},
                )
            return
        if not isinstance(item, (dict, list)):
            raise SemanticCommitError(
                "semantic commit nested values must be JSON-compatible",
                details={"field": path},
            )
        if len(item) > SEMANTIC_COMMIT_MAX_COLLECTION_ITEMS:
            raise SemanticCommitError(
                "semantic commit nested collection exceeds the request limit",
                details={
                    "field": path,
                    "maximum": SEMANTIC_COMMIT_MAX_COLLECTION_ITEMS,
                    "actual": len(item),
                },
            )
        container_id = id(item)
        if container_id in active_containers:
            raise SemanticCommitError(
                "semantic commit JSON payload may not contain cycles",
                details={"field": path},
            )
        active_containers.add(container_id)
        try:
            if isinstance(item, list):
                for index, child in enumerate(item):
                    visit(child, depth=depth + 1, path=f"{path}[{index}]")
                return
            for key, child in item.items():
                if not isinstance(key, str):
                    raise SemanticCommitError(
                        "semantic commit object keys must be strings",
                        details={"field": path},
                    )
                if len(key) > SEMANTIC_COMMIT_MAX_IDENTIFIER_CHARS:
                    raise SemanticCommitError(
                        "semantic commit object key exceeds the request limit",
                        details={
                            "field": path,
                            "maximum": SEMANTIC_COMMIT_MAX_IDENTIFIER_CHARS,
                            "actual": len(key),
                        },
                    )
                state["string_bytes"] += len(key.encode("utf-8"))
                if state["string_bytes"] > SEMANTIC_COMMIT_MAX_SERIALIZED_BYTES:
                    raise SemanticCommitError(
                        "semantic commit string content exceeds the aggregate request budget",
                        details={
                            "maximum_bytes": SEMANTIC_COMMIT_MAX_SERIALIZED_BYTES,
                        },
                    )
                visit(child, depth=depth + 1, path=f"{path}.value")
        finally:
            active_containers.remove(container_id)

    visit(value, depth=0, path="request")


def _bounded_semantic_commit_input(
    *,
    base_snapshot_id: object,
    base_semantic_state_hash: object,
    idempotency_key: object,
    claims: object,
    completed_revision_ids: object,
    progress_updates: object,
    materializer_version: object,
) -> tuple[list[dict], list[str], list[dict], bytes]:
    _require_semantic_commit_string(
        base_snapshot_id,
        field="base_snapshot_id",
        maximum=SEMANTIC_COMMIT_MAX_IDENTIFIER_CHARS,
        nonempty=True,
    )
    semantic_state_hash = _require_semantic_commit_string(
        base_semantic_state_hash,
        field="base_semantic_state_hash",
        maximum=64,
    )
    if (
        len(semantic_state_hash) != 64
        or any(character not in "0123456789abcdef" for character in semantic_state_hash)
    ):
        raise SemanticCommitError(
            "base_semantic_state_hash must be a lowercase SHA-256 digest"
        )
    _require_semantic_commit_string(
        idempotency_key,
        field="idempotency_key",
        maximum=SEMANTIC_COMMIT_MAX_IDENTIFIER_CHARS,
        nonempty=True,
    )
    _require_semantic_commit_string(
        materializer_version,
        field="materializer_version",
        maximum=SEMANTIC_COMMIT_MAX_MATERIALIZER_VERSION_CHARS,
        nonempty=True,
    )
    if not isinstance(claims, list):
        raise SemanticCommitError("claims must be a list")
    if len(claims) > SEMANTIC_COMMIT_MAX_CLAIMS:
        raise SemanticCommitError(
            f"a semantic commit may contain at most {SEMANTIC_COMMIT_MAX_CLAIMS} "
            "atomic claims",
            details={"claim_count": len(claims)},
        )
    if completed_revision_ids is None:
        completed_revision_values: list[object] = []
    elif isinstance(completed_revision_ids, list):
        completed_revision_values = completed_revision_ids
    else:
        raise SemanticCommitError("completed_revision_ids must be a list")
    if len(completed_revision_values) > SEMANTIC_COMMIT_MAX_COMPLETED_REVISIONS:
        raise SemanticCommitError(
            "a semantic commit may complete at most "
            f"{SEMANTIC_COMMIT_MAX_COMPLETED_REVISIONS} revisions",
            details={"completed_revision_count": len(completed_revision_values)},
        )
    if progress_updates is None:
        progress_values: list[object] = []
    elif isinstance(progress_updates, list):
        progress_values = progress_updates
    else:
        raise SemanticCommitError("progress_updates must be a list")
    if len(progress_values) > SEMANTIC_COMMIT_MAX_PROGRESS_UPDATES:
        raise SemanticCommitError(
            "a semantic commit may update at most "
            f"{SEMANTIC_COMMIT_MAX_PROGRESS_UPDATES} materialization checkpoints",
            details={"progress_update_count": len(progress_values)},
        )

    bounded_completed: list[str] = []
    for index, revision_id in enumerate(completed_revision_values, start=1):
        bounded_completed.append(
            _require_semantic_commit_string(
                revision_id,
                field=f"completed_revision_ids[{index}]",
                maximum=SEMANTIC_COMMIT_MAX_IDENTIFIER_CHARS,
                nonempty=True,
            ).strip()
        )
    bounded_completed = list(dict.fromkeys(bounded_completed))

    bounded_claims: list[dict] = []
    evidence_link_count = 0
    for claim_index, claim in enumerate(claims, start=1):
        if not isinstance(claim, dict):
            raise SemanticCommitError(
                "each claim must be an object",
                details={"claim_index": claim_index},
            )
        if "client_ref" in claim:
            _require_semantic_commit_string(
                claim["client_ref"],
                field=f"claims[{claim_index}].client_ref",
                maximum=SEMANTIC_COMMIT_MAX_IDENTIFIER_CHARS,
                nonempty=True,
            )
        body = _require_semantic_commit_string(
            claim.get("body"),
            field=f"claims[{claim_index}].body",
            maximum=SEMANTIC_COMMIT_MAX_BODY_CHARS,
            nonempty=True,
        )
        if not body.strip():
            raise SemanticCommitError(
                "claim body must contain 1-10000 characters",
                details={"claim_index": claim_index},
            )
        for field, maximum in (
            ("subject", SEMANTIC_COMMIT_MAX_SUBJECT_CHARS),
            ("apparent_status", SEMANTIC_COMMIT_MAX_STATUS_CHARS),
        ):
            if claim.get(field) is not None:
                _require_semantic_commit_string(
                    claim[field],
                    field=f"claims[{claim_index}].{field}",
                    maximum=maximum,
                )
        for field in (
            "modality",
            "claim_assessment",
            "temporal_applicability",
            "contest_state",
        ):
            if field in claim:
                _require_semantic_commit_string(
                    claim[field],
                    field=f"claims[{claim_index}].{field}",
                    maximum=SEMANTIC_COMMIT_MAX_IDENTIFIER_CHARS,
                    nonempty=True,
                )
        for field in ("scope_and_conditions", "time_window"):
            nested = claim.get(field, {})
            if not isinstance(nested, dict):
                raise SemanticCommitError(
                    f"claim {field} must be an object",
                    details={"claim_index": claim_index, "field": field},
                )
        evidence = claim.get("evidence", [])
        if (
            not isinstance(evidence, list)
            or not 1 <= len(evidence) <= SEMANTIC_COMMIT_MAX_EVIDENCE_PER_CLAIM
        ):
            raise SemanticCommitError(
                "each claim requires 1-50 evidence links",
                details={"claim_index": claim_index},
            )
        evidence_link_count += len(evidence)
        if evidence_link_count > SEMANTIC_COMMIT_MAX_EVIDENCE_LINKS:
            raise SemanticCommitError(
                "a semantic commit contains too many evidence links",
                details={
                    "maximum": SEMANTIC_COMMIT_MAX_EVIDENCE_LINKS,
                    "actual": evidence_link_count,
                },
            )
        for evidence_index, link in enumerate(evidence, start=1):
            if not isinstance(link, dict):
                raise SemanticCommitError(
                    "each evidence link must be an object",
                    details={
                        "claim_index": claim_index,
                        "evidence_index": evidence_index,
                    },
                )
            _require_semantic_commit_string(
                link.get("source_unit_id"),
                field=(
                    f"claims[{claim_index}].evidence[{evidence_index}].source_unit_id"
                ),
                maximum=SEMANTIC_COMMIT_MAX_IDENTIFIER_CHARS,
                nonempty=True,
            )
            if "stance" in link:
                _require_semantic_commit_string(
                    link["stance"],
                    field=f"claims[{claim_index}].evidence[{evidence_index}].stance",
                    maximum=SEMANTIC_COMMIT_MAX_IDENTIFIER_CHARS,
                    nonempty=True,
                )
            if link.get("qualifier") is not None:
                _require_semantic_commit_string(
                    link["qualifier"],
                    field=f"claims[{claim_index}].evidence[{evidence_index}].qualifier",
                    maximum=SEMANTIC_COMMIT_MAX_QUALIFIER_CHARS,
                )
            if not isinstance(link.get("applicability", {}), dict):
                raise SemanticCommitError(
                    "evidence applicability must be an object",
                    details={
                        "claim_index": claim_index,
                        "evidence_index": evidence_index,
                    },
                )
        bounded_claims.append(claim)

    bounded_progress: list[dict] = []
    for progress_index, update in enumerate(progress_values, start=1):
        if not isinstance(update, dict):
            raise SemanticCommitError(
                "each progress update must be an object",
                details={"progress_index": progress_index},
            )
        _require_semantic_commit_string(
            update.get("revision_id"),
            field=f"progress_updates[{progress_index}].revision_id",
            maximum=SEMANTIC_COMMIT_MAX_IDENTIFIER_CHARS,
            nonempty=True,
        )
        batch_receipt = _require_semantic_commit_string(
            update.get("batch_receipt"),
            field=f"progress_updates[{progress_index}].batch_receipt",
            maximum=64,
        )
        if (
            len(batch_receipt) != 64
            or any(character not in "0123456789abcdef" for character in batch_receipt)
        ):
            raise SemanticCommitError(
                "progress requires a valid server-issued batch_receipt",
                details={"progress_index": progress_index},
            )
        for field, minimum in (
            ("processed_from_ordinal", 1),
            ("processed_through_ordinal", 1),
            ("next_ordinal", 2),
        ):
            ordinal = update.get(field)
            if (
                isinstance(ordinal, bool)
                or not isinstance(ordinal, int)
                or not minimum <= ordinal <= SEMANTIC_COMMIT_MAX_INTEGER
            ):
                raise SemanticCommitError(
                    "progress requires bounded integer ordinals",
                    details={
                        "progress_index": progress_index,
                        "field": field,
                        "minimum": minimum,
                        "maximum": SEMANTIC_COMMIT_MAX_INTEGER,
                    },
                )
        if update.get("note") is not None:
            _require_semantic_commit_string(
                update["note"],
                field=f"progress_updates[{progress_index}].note",
                maximum=SEMANTIC_COMMIT_MAX_STATUS_CHARS,
            )
        bounded_progress.append(update)

    if not bounded_claims and not bounded_progress and not bounded_completed:
        raise SemanticCommitError(
            "a semantic commit requires claims, progress, or completed revisions"
        )
    canonical_input = {
        "base_snapshot_id": base_snapshot_id,
        "base_semantic_state_hash": semantic_state_hash,
        "claims": bounded_claims,
        "completed_revision_ids": bounded_completed,
        "progress_updates": bounded_progress,
        "materializer_version": materializer_version,
    }
    _validate_semantic_commit_json_bounds(canonical_input)
    canonical_bytes = encode_json(canonical_input).encode()
    if len(canonical_bytes) > SEMANTIC_COMMIT_MAX_SERIALIZED_BYTES:
        raise SemanticCommitError(
            "semantic commit serialized request exceeds the aggregate byte limit",
            details={
                "maximum_bytes": SEMANTIC_COMMIT_MAX_SERIALIZED_BYTES,
                "actual_bytes": len(canonical_bytes),
            },
        )
    return bounded_claims, bounded_completed, bounded_progress, canonical_bytes


class CorpusService:
    def __init__(
        self,
        data_root: Path | None = None,
        *,
        adapter_registry: AdapterRegistry | None = None,
        maintain_legacy_semantic_cache: bool = False,
    ) -> None:
        self.data_root = (data_root or default_data_root()).expanduser().resolve()
        self.adapter_registry = adapter_registry or build_default_registry(
            self.data_root / "runtime"
        )
        self.maintain_legacy_semantic_cache = maintain_legacy_semantic_cache
        self.contexts = ContextService(
            self.data_root,
            adapter_registry=self.adapter_registry,
        )

    def _projection_uses_current_adapter(
        self,
        extension: str,
        adapter_id: str,
        adapter_version: str,
        config_hash: str,
    ) -> bool:
        try:
            descriptor = self.adapter_registry.resolve(extension).descriptor
        except ExtractionError:
            return False
        return (
            adapter_id,
            adapter_version,
            config_hash,
        ) == (
            descriptor.adapter_id,
            descriptor.adapter_version,
            descriptor.config_hash,
        )

    def register(
        self,
        *,
        corpus_id: str,
        source_root: Path,
        execution_policy: str,
        provider_kind: str = "filesystem",
        source_scope: dict | None = None,
    ) -> dict:
        return register_corpus(
            data_root=self.data_root,
            corpus_id=corpus_id,
            source_root=source_root,
            execution_policy=execution_policy,
            provider_kind=provider_kind,
            source_scope=source_scope,
        )

    def configure_source_scope(
        self,
        *,
        corpus_id: str,
        exclude_directory_names: object = (),
        exclude_path_prefixes: object = (),
    ) -> dict:
        return configure_corpus_source_scope(
            data_root=self.data_root,
            corpus_id=corpus_id,
            exclude_directory_names=exclude_directory_names,
            exclude_path_prefixes=exclude_path_prefixes,
        )

    def corpora(self) -> list[dict]:
        return list_corpora(self.data_root)

    def overview(
        self,
        *,
        audience: str = "local_cli",
        max_items_per_context: int = CORPUS_OVERVIEW_DEFAULT_ITEMS_PER_CONTEXT,
    ) -> dict:
        if audience not in {"local_cli", "external_mcp"}:
            raise ConfigurationError(
                "unsupported corpus overview audience",
                details={"audience": audience},
            )
        if not 1 <= max_items_per_context <= CORPUS_OVERVIEW_MAX_ITEMS_PER_CONTEXT:
            raise BudgetExceededError(
                "corpus overview item limit is outside the supported range",
                details={
                    "max_items_per_context": max_items_per_context,
                    "maximum": CORPUS_OVERVIEW_MAX_ITEMS_PER_CONTEXT,
                },
            )

        corpora = self.corpora()
        if audience == "external_mcp":
            corpora = [
                corpus
                for corpus in corpora
                if corpus["execution_policy"] == "external_host_allowed"
            ]

        context_listing = self.context_read(
            context_id=None,
            state="active",
            limit=100,
            offset=0,
            audience=audience,
            view="restricted",
        )
        archived_context_listing = self.context_read(
            context_id=None,
            state="archived",
            limit=100,
            offset=0,
            audience=audience,
            view="restricted",
        )
        contexts_by_corpus: dict[str, list[dict]] = {
            corpus["corpus_id"]: [] for corpus in corpora
        }
        archived_contexts_by_corpus: dict[str, list[dict]] = {
            corpus["corpus_id"]: [] for corpus in corpora
        }
        active_contexts = []
        for context in context_listing["contexts"]:
            detail = self.context_read(
                context_id=context["context_id"],
                state="active",
                include_history=False,
                limit=CONTEXT_MAX_LIMIT,
                offset=0,
                audience=audience,
                view="restricted",
            )
            history = self.context_read(
                context_id=context["context_id"],
                state="active",
                include_history=True,
                limit=CONTEXT_MAX_LIMIT,
                offset=0,
                audience=audience,
                view="restricted",
            )
            stale_item_count, stale_source_link_count = (
                self._overview_stale_counts(detail["items"])
            )
            context_view = {
                **detail["context"],
                "item_count": detail["total_matching"],
                "items_truncated": (
                    detail["total_matching"] > max_items_per_context
                ),
                "items": [
                    self._overview_context_item(item)
                    for item in detail["items"][:max_items_per_context]
                ],
                "stale_item_count": stale_item_count,
                "stale_source_link_count": stale_source_link_count,
                "stale_counts_truncated": detail["has_more"],
                "superseded_item_count": max(
                    0,
                    history["total_matching"] - detail["total_matching"],
                ),
                "history_truncated": history["has_more"],
            }
            active_contexts.append(context_view)
            for corpus_id in context["corpus_ids"]:
                if corpus_id in contexts_by_corpus:
                    contexts_by_corpus[corpus_id].append(context_view)
        for context in archived_context_listing["contexts"]:
            context_view = {
                "context_id": context["context_id"],
                "title": context["title"],
                "purpose": context["purpose"],
                "scope": context["scope"],
                "state": context["state"],
                "version": context["version"],
                "updated_at": context["updated_at"],
            }
            for corpus_id in context["corpus_ids"]:
                if corpus_id in archived_contexts_by_corpus:
                    archived_contexts_by_corpus[corpus_id].append(context_view)

        corpus_views = []
        for corpus in corpora:
            corpus_id = corpus["corpus_id"]
            status = self.status(corpus_id)
            linked_sources = self.corpus_source_read(
                corpus_id=corpus_id,
                record_state="active",
                limit=1,
                offset=0,
                audience=audience,
            )
            contexts = contexts_by_corpus[corpus_id]
            archived_contexts = archived_contexts_by_corpus[corpus_id]

            source_root = corpus["source_root_nfc"]
            corpus_views.append(
                {
                    "corpus_id": corpus_id,
                    "display_name": (
                        Path(source_root).name
                        if audience == "local_cli"
                        else corpus_id
                    ),
                    "execution_policy": corpus["execution_policy"],
                    "provider_kind": corpus["provider_kind"],
                    "source_scope": corpus["source_scope"],
                    "source_root": source_root if audience == "local_cli" else None,
                    "source_index": {
                        "inventory_complete": bool(
                            status["latest_scan"]
                            and status["latest_scan"]["status"] == "complete"
                        ),
                        "scan_completed_at": (
                            status["latest_scan"]["completed_at"]
                            if status["latest_scan"]
                            else None
                        ),
                        "documents": int(status["totals"]["documents"] or 0),
                        "supported_documents": int(
                            status["totals"]["supported_documents"] or 0
                        ),
                        "indexed_documents": int(
                            status["totals"]["indexed_documents"] or 0
                        ),
                        "remote_supported_documents": int(
                            status["document_states"].get(
                                "supported:remote_only",
                                0,
                            )
                        ),
                        "active_source_units": status["active_source_units"],
                        "coverage_gaps": status["coverage_gaps"],
                    },
                    "linked_sources": [
                        {
                            "binding_id": binding["binding_id"],
                            "provider_kind": binding["provider_kind"],
                            "selector": binding["selector"],
                            "state": binding["state"],
                            "active_record_count": binding["active_record_count"],
                            "last_complete_at": binding["last_complete_at"],
                        }
                        for binding in linked_sources["bindings"]
                    ],
                    "contexts": contexts,
                    "archived_contexts": archived_contexts,
                    "context_lifecycle": {
                        "active_context_count": len(contexts),
                        "archived_context_count": len(archived_contexts),
                        "active_item_count": sum(
                            context["item_count"] for context in contexts
                        ),
                        "stale_item_count": sum(
                            context["stale_item_count"] for context in contexts
                        ),
                        "stale_source_link_count": sum(
                            context["stale_source_link_count"]
                            for context in contexts
                        ),
                        "superseded_item_count": sum(
                            context["superseded_item_count"]
                            for context in contexts
                        ),
                    },
                }
            )

        response = {
            "view": "personal",
            "corpus_count": len(corpus_views),
            "corpora": corpus_views,
            "context_lifecycle": {
                "active_context_count": len(active_contexts),
                "archived_context_count": archived_context_listing[
                    "total_matching"
                ],
                "active_item_count": sum(
                    context["item_count"] for context in active_contexts
                ),
                "stale_item_count": sum(
                    context["stale_item_count"] for context in active_contexts
                ),
                "stale_source_link_count": sum(
                    context["stale_source_link_count"]
                    for context in active_contexts
                ),
                "superseded_item_count": sum(
                    context["superseded_item_count"]
                    for context in active_contexts
                ),
            },
        }
        serialized_bytes = len(encode_json(response).encode())
        if serialized_bytes > CORPUS_OVERVIEW_MAX_SERIALIZED_BYTES:
            raise BudgetExceededError(
                "corpus overview exceeds the serialized response budget",
                details={
                    "serialized_bytes": serialized_bytes,
                    "maximum_bytes": CORPUS_OVERVIEW_MAX_SERIALIZED_BYTES,
                    "suggestion": "reduce max_items_per_context",
                },
            )
        return response

    @staticmethod
    def _overview_context_item(item: dict) -> dict:
        body_text = item["body_text"]
        truncated = len(body_text) > CORPUS_OVERVIEW_MAX_BODY_CHARS
        if truncated:
            body_text = body_text[:CORPUS_OVERVIEW_MAX_BODY_CHARS].rstrip() + "…"
        return {
            "kind": item["kind"],
            "body_text": body_text,
            "body_truncated": truncated,
            "attributes": item["attributes"],
            "created_at": item["created_at"],
            "source_count": len(item["sources"]),
            "linked_source_count": len(item["external_sources"]),
        }

    @staticmethod
    def _overview_stale_counts(items: list[dict]) -> tuple[int, int]:
        stale_item_count = 0
        stale_source_link_count = 0
        for item in items:
            stale_links = sum(
                source.get("dependency_state") != "valid"
                for source in [*item["sources"], *item["external_sources"]]
            )
            if stale_links:
                stale_item_count += 1
                stale_source_link_count += stale_links
        return stale_item_count, stale_source_link_count

    def context_read(
        self,
        *,
        context_id: str | None = None,
        state: str = "active",
        include_history: bool = False,
        limit: int = 100,
        offset: int = 0,
        audience: str = "local_cli",
        view: str = "restricted",
    ) -> dict:
        return self.contexts.read(
            context_id=context_id,
            state=state,
            include_history=include_history,
            limit=limit,
            offset=offset,
            audience=audience,
            view=view,
        )

    def context_update(
        self,
        *,
        action: str,
        context_id: str,
        expected_version: int,
        payload: dict,
        confirm_persistent_context_write: bool,
        confirm_general_release_approval: bool = False,
        audience: str = "local_cli",
    ) -> dict:
        return self.contexts.update(
            action=action,
            context_id=context_id,
            expected_version=expected_version,
            payload=payload,
            confirm_persistent_context_write=confirm_persistent_context_write,
            confirm_general_release_approval=confirm_general_release_approval,
            audience=audience,
        )

    def context_migrate(self) -> dict:
        return self.contexts.migrate()

    def corpus_source_read(
        self,
        *,
        corpus_id: str,
        binding_id: str | None = None,
        record_state: str = "active",
        occurred_after: str | None = None,
        limit: int = 100,
        offset: int = 0,
        audience: str = "local_cli",
    ) -> dict:
        return self.contexts.source_read(
            corpus_id=corpus_id,
            binding_id=binding_id,
            record_state=record_state,
            occurred_after=occurred_after,
            limit=limit,
            offset=offset,
            audience=audience,
        )

    def corpus_source_update(
        self,
        *,
        action: str,
        corpus_id: str,
        binding_id: str,
        payload: dict,
        confirm_persistent_context_write: bool,
        audience: str = "local_cli",
    ) -> dict:
        return self.contexts.source_update(
            action=action,
            corpus_id=corpus_id,
            binding_id=binding_id,
            payload=payload,
            confirm_persistent_context_write=confirm_persistent_context_write,
            audience=audience,
        )

    def corpus_source_fetch(
        self,
        *,
        corpus_id: str,
        binding_id: str,
        external_id: str,
        max_chars: int = SESSION_SOURCE_FETCH_DEFAULT_CHARS,
        audience: str = "local_cli",
    ) -> dict:
        return self.contexts.source_fetch(
            corpus_id=corpus_id,
            binding_id=binding_id,
            external_id=external_id,
            max_chars=max_chars,
            audience=audience,
        )

    def migrate(self, corpus_id: str) -> dict:
        corpus_id = normalize_corpus_id(corpus_id)
        paths = self._paths(corpus_id)
        paths.ensure()
        with writer_lock(paths.corpus_root / "writer.lock"):
            return migrate_corpus(self.data_root, corpus_id)

    def migration_status(self, corpus_id: str) -> dict:
        return corpus_schema_status(self.data_root, corpus_id)

    def scan(self, corpus_id: str) -> dict:
        corpus_id = normalize_corpus_id(corpus_id)
        paths = self._paths(corpus_id)
        paths.ensure()
        with writer_lock(paths.corpus_root / "writer.lock"):
            result = scan_corpus(self.data_root, corpus_id)
            result["snapshot"] = self._publish_snapshot(corpus_id)
            return result

    def cleanup_source_copies(
        self,
        corpus_id: str,
        *,
        confirm_delete: bool = False,
    ) -> dict:
        """Plan or remove retained source bytes without changing extracted units."""

        corpus_id = normalize_corpus_id(corpus_id)
        get_corpus(self.data_root, corpus_id)
        paths = self._paths(corpus_id)
        paths.ensure()
        with writer_lock(paths.corpus_root / "writer.lock"):
            cleanup = cleanup_source_copies(paths, delete=False)
            if confirm_delete:
                cleanup = cleanup_source_copies(paths, delete=True)
            references_marked_ephemeral = 0
            with corpus_connection(self.data_root, corpus_id) as connection:
                source_units_preserved = int(
                    connection.execute("SELECT COUNT(*) FROM source_units").fetchone()[0]
                )
                legacy_reference_rows = connection.execute(
                    """
                    SELECT revision_id, sha256, immutable_blob_ref
                    FROM revisions
                    WHERE immutable_blob_ref LIKE 'blobs/%'
                    """
                ).fetchall()
                canonical_reference_rows = [
                    row
                    for row in legacy_reference_rows
                    if row["immutable_blob_ref"]
                    == _canonical_legacy_blob_ref(row["sha256"])
                ]
                present_digests = set(cleanup.canonical_blob_digests)
                missing_legacy_references = sum(
                    1
                    for row in canonical_reference_rows
                    if row["sha256"] not in present_digests
                )
                if confirm_delete:
                    for row in canonical_reference_rows:
                        cursor = connection.execute(
                            """
                            UPDATE revisions
                            SET immutable_blob_ref = ?
                            WHERE revision_id = ?
                              AND immutable_blob_ref = ?
                            """,
                            (
                                _ephemeral_capture_ref(row["sha256"]),
                                row["revision_id"],
                                row["immutable_blob_ref"],
                            ),
                        )
                        references_marked_ephemeral += cursor.rowcount
            result = cleanup.as_dict(deleted=confirm_delete)
            result.update(
                {
                    "corpus_id": corpus_id,
                    "confirmation_required": not confirm_delete,
                    "references_marked_ephemeral": references_marked_ephemeral,
                    "canonical_legacy_references": len(canonical_reference_rows),
                    "missing_legacy_references": missing_legacy_references,
                    "noncanonical_legacy_references_skipped": (
                        len(legacy_reference_rows) - len(canonical_reference_rows)
                    ),
                    "source_units_preserved": source_units_preserved,
                    "search_index_preserved": True,
                }
            )
            return result

    def _paths(self, corpus_id: str) -> RuntimePaths:
        corpus_id = normalize_corpus_id(corpus_id)
        return RuntimePaths(data_root=self.data_root, corpus_id=corpus_id)

    def status(
        self,
        corpus_id: str,
        *,
        include_derived: bool = False,
    ) -> dict:
        corpus = get_corpus(self.data_root, corpus_id)
        with corpus_read_connection(self.data_root, corpus_id) as connection:
            scan = connection.execute(
                "SELECT * FROM scan_runs ORDER BY rowid DESC LIMIT 1"
            ).fetchone()
            document_counts = {
                row["state"]: row["count"]
                for row in connection.execute(
                    """
                    SELECT eligibility_state || ':' || residency_state AS state, COUNT(*) AS count
                    FROM documents WHERE deleted_at IS NULL GROUP BY state
                    """
                )
            }
            totals = dict(
                connection.execute(
                    """
                    SELECT
                        COUNT(*) AS documents,
                        COALESCE(SUM(logical_size), 0) AS logical_bytes,
                        COALESCE(SUM(allocated_size), 0) AS allocated_bytes,
                        SUM(CASE WHEN is_dataless = 1 THEN 1 ELSE 0 END) AS dataless_documents,
                        SUM(CASE WHEN eligibility_state = 'supported' THEN 1 ELSE 0 END)
                            AS supported_documents,
                        SUM(
                            CASE WHEN eligibility_state = 'supported' AND EXISTS (
                                SELECT 1 FROM extraction_projections p
                                JOIN revisions r ON r.revision_id = p.revision_id
                                WHERE p.revision_id = documents.current_revision_id
                                  AND p.is_active = 1
                                  AND r.source_size = documents.logical_size
                                  AND r.source_modified_ns = documents.modified_ns
                                  AND r.source_changed_ns = documents.changed_ns
                                  AND r.source_device = documents.device
                                  AND r.source_inode = documents.inode
                            ) THEN 1 ELSE 0 END
                        ) AS indexed_documents
                    FROM documents
                    WHERE deleted_at IS NULL
                    """
                ).fetchone()
            )
            revisions = connection.execute("SELECT COUNT(*) FROM revisions").fetchone()[0]
            units = connection.execute("SELECT COUNT(*) FROM source_units").fetchone()[0]
            active_units = connection.execute(
                """
                SELECT COUNT(*)
                FROM source_units u
                JOIN extraction_projections p ON p.projection_id = u.projection_id
                JOIN revisions r ON r.revision_id = u.revision_id
                JOIN documents d ON d.current_revision_id = u.revision_id
                WHERE p.is_active = 1 AND d.deleted_at IS NULL
                  AND d.eligibility_state = 'supported'
                  AND r.source_size = d.logical_size
                  AND r.source_modified_ns = d.modified_ns
                  AND r.source_changed_ns = d.changed_ns
                  AND r.source_device = d.device
                  AND r.source_inode = d.inode
                """
            ).fetchone()[0]
            projections = {
                row["state"]: row["count"]
                for row in connection.execute(
                    """
                    SELECT
                        CASE
                            WHEN is_active = 1
                            THEN 'active:' || completeness_state || ':' || assurance_state
                            ELSE 'historical:' || completeness_state || ':' || assurance_state
                        END AS state,
                        COUNT(*) AS count
                    FROM extraction_projections
                    GROUP BY state
                    """
                )
            }
            attempts = {
                row["state"]: row["count"]
                for row in connection.execute(
                    """
                    SELECT state, COUNT(*) AS count
                    FROM extraction_attempts
                    GROUP BY state
                    """
                )
            }
            active_projection_rows = connection.execute(
                """
                SELECT d.extension, p.adapter_id, p.adapter_version, p.config_hash,
                       p.completeness_state
                FROM documents d
                JOIN revisions r ON r.revision_id = d.current_revision_id
                JOIN extraction_projections p
                  ON p.revision_id = d.current_revision_id AND p.is_active = 1
                WHERE d.deleted_at IS NULL
                  AND d.eligibility_state = 'supported'
                  AND r.source_size = d.logical_size
                  AND r.source_modified_ns = d.modified_ns
                  AND r.source_changed_ns = d.changed_ns
                  AND r.source_device = d.device
                  AND r.source_inode = d.inode
                """
            ).fetchall()
            issues = connection.execute("SELECT COUNT(*) FROM extraction_issues").fetchone()[0]
            issue_lifecycle = {
                row["lifecycle_state"]: row["count"]
                for row in connection.execute(
                    """
                    SELECT lifecycle_state, COUNT(*) AS count
                    FROM extraction_issues
                    GROUP BY lifecycle_state
                    """
                )
            }
            queue = None
            semantic_claims = None
            if include_derived:
                queue = {
                    row["state"]: row["count"]
                    for row in connection.execute(
                        "SELECT state, COUNT(*) AS count "
                        "FROM interpretation_queue GROUP BY state"
                    )
                }
                semantic_claims = {
                    row["dependency_state"]: row["count"]
                    for row in connection.execute(
                        """
                        SELECT dependency_state, COUNT(*) AS count
                        FROM atomic_claims GROUP BY dependency_state
                        """
                    )
                }
            snapshot = connection.execute(
                """
                SELECT * FROM snapshots
                WHERE state = 'complete'
                ORDER BY rowid DESC
                LIMIT 1
                """
            ).fetchone()
        outdated_projections = 0
        partial_projections = 0
        for row in active_projection_rows:
            descriptor = self.adapter_registry.resolve(row["extension"]).descriptor
            if row["completeness_state"] != "complete":
                partial_projections += 1
            if (
                row["adapter_id"],
                row["adapter_version"],
                row["config_hash"],
            ) != (
                descriptor.adapter_id,
                descriptor.adapter_version,
                descriptor.config_hash,
            ):
                outdated_projections += 1
        supported_documents = int(totals["supported_documents"] or 0)
        indexed_documents = int(totals["indexed_documents"] or 0)
        coverage_gaps = {
            "supported_documents_without_usable_projection": max(
                0, supported_documents - indexed_documents
            ),
            "partial_active_projections": partial_projections,
            "outdated_active_projections": outdated_projections,
        }
        staging_observation = observe_staging(self._paths(corpus_id))
        response = {
            "corpus": corpus,
            "data_root": str(self.data_root),
            "latest_scan": dict(scan) if scan else None,
            "totals": totals,
            "document_states": document_counts,
            "revisions": revisions,
            "source_units": units,
            "active_source_units": active_units,
            "extraction_projections": projections,
            "extraction_attempts": attempts,
            "current_snapshot": dict(snapshot) if snapshot else None,
            "coverage_gaps": coverage_gaps,
            "issues": issues,
            "issue_lifecycle": issue_lifecycle,
            "extraction_adapters": [
                descriptor.to_dict() for descriptor in self.adapter_registry.descriptors
            ],
            "authority": {
                "source": "registered source bytes",
                "extracted_projection": "source_units",
                "request_time_interpretation": "host agent ephemeral task context",
            },
            "source_copy_retention": {
                "default": "ephemeral",
                "persistent_source_bytes_required_for_search": False,
                "intentional_absence_marker": "ephemeral:sha256:<digest>",
                "staging_observation": staging_observation,
            },
        }
        if include_derived:
            response["interpretation_queue"] = queue
            response["semantic_claims"] = semantic_claims
            response["authority"]["optional_legacy_semantic_cache"] = (
                "atomic_claims + evidence_links"
            )
        return response

    def _document_index_state(self, document: dict) -> tuple[str, list[str]]:
        if document["eligibility_state"] != "supported":
            return "not_applicable", [f"eligibility:{document['eligibility_state']}"]
        if document["current_revision_id"] is None:
            return "unindexed", ["no_current_revision"]

        reasons: list[str] = []
        if document["active_projection_id"] is None:
            reasons.append("active_projection_missing")
        source_observation_current = (
            document["revision_source_size"] == document["logical_size"]
            and document["revision_source_modified_ns"] == document["modified_ns"]
            and document["revision_source_changed_ns"] == document["changed_ns"]
            and document["revision_source_device"] == document["device"]
            and document["revision_source_inode"] == document["inode"]
        )
        if not source_observation_current:
            reasons.append("source_observation_changed")

        descriptor = self.adapter_registry.resolve(document["extension"]).descriptor
        adapter_current = (
            document["active_projection_id"] is not None
            and document["projection_adapter_id"] == descriptor.adapter_id
            and document["projection_adapter_version"] == descriptor.adapter_version
            and document["projection_config_hash"] == descriptor.config_hash
        )
        if document["active_projection_id"] is not None and not adapter_current:
            reasons.append("outdated_adapter")
        return ("refresh_required", reasons) if reasons else ("current", [])

    def inventory(
        self,
        corpus_id: str,
        *,
        path_contains: str | None = None,
        eligibility_state: str = "supported",
        residency_state: str = "all",
        index_state: str = "all",
        extension: str | None = None,
        max_logical_bytes: int | None = None,
        limit: int = CORPUS_INVENTORY_DEFAULT_LIMIT,
        offset: int = 0,
    ) -> dict:
        if eligibility_state not in _INVENTORY_ELIGIBILITY_STATES:
            raise ConfigurationError(
                "unsupported inventory eligibility state",
                details={
                    "eligibility_state": eligibility_state,
                    "allowed": sorted(_INVENTORY_ELIGIBILITY_STATES),
                },
            )
        if residency_state not in _INVENTORY_RESIDENCY_STATES:
            raise ConfigurationError(
                "unsupported inventory residency state",
                details={
                    "residency_state": residency_state,
                    "allowed": sorted(_INVENTORY_RESIDENCY_STATES),
                },
            )
        if index_state not in _INVENTORY_INDEX_STATES:
            raise ConfigurationError(
                "unsupported inventory index state",
                details={
                    "index_state": index_state,
                    "allowed": sorted(_INVENTORY_INDEX_STATES),
                },
            )
        if not 1 <= limit <= CORPUS_INVENTORY_MAX_LIMIT:
            raise BudgetExceededError(
                "inventory limit must be between 1 and 200",
                details={
                    "limit": limit,
                    "maximum": CORPUS_INVENTORY_MAX_LIMIT,
                },
            )
        if not 0 <= offset <= CORPUS_INVENTORY_MAX_OFFSET:
            raise BudgetExceededError(
                "inventory offset is outside the supported range",
                details={
                    "offset": offset,
                    "maximum": CORPUS_INVENTORY_MAX_OFFSET,
                },
            )
        if max_logical_bytes is not None and not (
            1 <= max_logical_bytes <= CORPUS_INVENTORY_MAX_LOGICAL_BYTES
        ):
            raise BudgetExceededError(
                "inventory maximum logical size is outside the supported range",
                details={
                    "max_logical_bytes": max_logical_bytes,
                    "maximum": CORPUS_INVENTORY_MAX_LOGICAL_BYTES,
                },
            )

        normalized_path_filter = _normalize_inventory_path_filter(path_contains)
        normalized_extension = _normalize_inventory_extension(extension)
        with corpus_read_connection(self.data_root, corpus_id) as connection:
            latest_scan = connection.execute(
                """
                SELECT scan_id, status, completed_at
                FROM scan_runs
                ORDER BY rowid DESC
                LIMIT 1
                """
            ).fetchone()
            snapshot = connection.execute(
                """
                SELECT snapshot_id, coverage_state
                FROM snapshots
                WHERE state = 'complete'
                ORDER BY rowid DESC
                LIMIT 1
                """
            ).fetchone()
            rows = connection.execute(
                """
                SELECT d.document_id, d.relative_path, d.relative_path_nfc,
                       d.extension, d.media_type, d.logical_size, d.modified_ns,
                       d.residency_state, d.eligibility_state,
                       d.current_revision_id, d.device, d.inode, d.changed_ns,
                       r.source_size AS revision_source_size,
                       r.source_modified_ns AS revision_source_modified_ns,
                       r.source_changed_ns AS revision_source_changed_ns,
                       r.source_device AS revision_source_device,
                       r.source_inode AS revision_source_inode,
                       p.projection_id AS active_projection_id,
                       p.adapter_id AS projection_adapter_id,
                       p.adapter_version AS projection_adapter_version,
                       p.config_hash AS projection_config_hash,
                       p.completeness_state AS projection_completeness
                FROM documents d
                LEFT JOIN revisions r ON r.revision_id = d.current_revision_id
                LEFT JOIN extraction_projections p
                  ON p.revision_id = d.current_revision_id AND p.is_active = 1
                WHERE d.deleted_at IS NULL
                ORDER BY d.relative_path_nfc COLLATE BINARY, d.document_id
                """
            ).fetchall()

        path_filter_key = (
            normalized_path_filter.casefold()
            if normalized_path_filter is not None
            else None
        )
        documents = []
        for row in rows:
            document = dict(row)
            if not _safe_relative_inventory_path(document["relative_path"]):
                raise ConfigurationError(
                    "inventory contains an unsafe relative document locator",
                    details={"document_id": document["document_id"]},
                )
            if (
                eligibility_state != "all"
                and document["eligibility_state"] != eligibility_state
            ):
                continue
            if (
                residency_state != "all"
                and document["residency_state"] != residency_state
            ):
                continue
            if normalized_extension is not None and (
                document["extension"] != normalized_extension
            ):
                continue
            if (
                max_logical_bytes is not None
                and document["logical_size"] > max_logical_bytes
            ):
                continue
            if path_filter_key is not None and path_filter_key not in (
                unicodedata.normalize(
                    "NFC",
                    document["relative_path_nfc"],
                ).casefold()
            ):
                continue

            document_index_state, refresh_reasons = self._document_index_state(
                document
            )
            if index_state != "all" and document_index_state != index_state:
                continue
            documents.append(
                {
                    "document_id": document["document_id"],
                    "relative_path": document["relative_path"],
                    "extension": document["extension"],
                    "media_type": document["media_type"],
                    "logical_size": document["logical_size"],
                    "modified_ns": document["modified_ns"],
                    "residency_state": document["residency_state"],
                    "eligibility_state": document["eligibility_state"],
                    "current_revision_id": document["current_revision_id"],
                    "active_projection_id": document["active_projection_id"],
                    "projection_completeness": document[
                        "projection_completeness"
                    ],
                    "index_state": document_index_state,
                    "refresh_reasons": refresh_reasons,
                }
            )

        total_matching = len(documents)
        page = documents[offset : offset + limit]
        next_offset = offset + len(page)
        has_more = next_offset < total_matching
        response = {
            "corpus_id": normalize_corpus_id(corpus_id),
            "observation": {
                "latest_scan_id": latest_scan["scan_id"] if latest_scan else None,
                "latest_scan_status": latest_scan["status"] if latest_scan else None,
                "scan_completed_at": latest_scan["completed_at"] if latest_scan else None,
                "inventory_complete": bool(
                    latest_scan is not None and latest_scan["status"] == "complete"
                ),
                "current_snapshot_id": (
                    snapshot["snapshot_id"] if snapshot else None
                ),
                "snapshot_coverage_state": (
                    snapshot["coverage_state"] if snapshot else None
                ),
            },
            "filters": {
                "path_contains": normalized_path_filter,
                "eligibility_state": eligibility_state,
                "residency_state": residency_state,
                "index_state": index_state,
                "extension": normalized_extension,
                "max_logical_bytes": max_logical_bytes,
            },
            "offset": offset,
            "limit": limit,
            "returned_count": len(page),
            "total_matching": total_matching,
            "has_more": has_more,
            "next_offset": next_offset if has_more else None,
            "documents": page,
            "metadata_only": True,
            "relevance_assessed": False,
            "notice": (
                "Inventory metadata supports exact document selection only. "
                "Filenames are untrusted metadata, not evidence of document content. "
                "A complete inventory does not establish content absence."
            ),
        }
        serialized_bytes = len(encode_json(response).encode())
        if serialized_bytes > CORPUS_INVENTORY_MAX_SERIALIZED_BYTES:
            raise BudgetExceededError(
                "inventory response exceeds the serialized response budget",
                details={
                    "requested_limit": limit,
                    "offset": offset,
                    "returned_document_count": len(page),
                    "serialized_bytes": serialized_bytes,
                    "maximum_serialized_bytes": (
                        CORPUS_INVENTORY_MAX_SERIALIZED_BYTES
                    ),
                    "retry_with_lower": ["limit"],
                    "retry_with_narrower": [
                        "path_contains",
                        "eligibility_state",
                        "residency_state",
                        "index_state",
                        "extension",
                        "max_logical_bytes",
                    ],
                },
            )
        return response

    def _pending_documents(
        self,
        corpus_id: str,
        *,
        include_remote: bool,
        max_file_bytes: int,
        remote_only: bool = False,
    ) -> tuple[list[dict], dict]:
        with corpus_read_connection(self.data_root, corpus_id) as connection:
            rows = connection.execute(
                """
                SELECT d.*, r.source_size AS revision_source_size,
                       r.source_modified_ns AS revision_source_modified_ns,
                       r.source_changed_ns AS revision_source_changed_ns,
                       r.source_device AS revision_source_device,
                       r.source_inode AS revision_source_inode,
                       p.projection_id AS active_projection_id,
                       p.adapter_id AS projection_adapter_id,
                       p.adapter_version AS projection_adapter_version,
                       p.config_hash AS projection_config_hash
                FROM documents d
                LEFT JOIN revisions r ON r.revision_id = d.current_revision_id
                LEFT JOIN extraction_projections p
                  ON p.revision_id = d.current_revision_id AND p.is_active = 1
                WHERE d.deleted_at IS NULL
                  AND d.eligibility_state = 'supported'
                ORDER BY d.is_dataless ASC, d.logical_size ASC, d.relative_path_nfc ASC
                """
            ).fetchall()
        pending: list[dict] = []
        skipped = {
            "current": 0,
            "remote": 0,
            "local": 0,
            "too_large": 0,
            "not_selected": 0,
        }
        for row in rows:
            document = dict(row)
            document_index_state, _ = self._document_index_state(document)
            if document_index_state == "current":
                skipped["current"] += 1
                continue
            if document["logical_size"] > max_file_bytes:
                skipped["too_large"] += 1
                continue
            if remote_only and not document["is_dataless"]:
                skipped["local"] += 1
                continue
            if document["is_dataless"] and not include_remote:
                skipped["remote"] += 1
                continue
            pending.append(document)
        return pending, skipped

    @staticmethod
    def _safe_selection_extension(document: dict) -> str | None:
        extension = document.get("extension")
        if not isinstance(extension, str):
            return None
        normalized = extension.strip().lower().removeprefix(".")
        if not normalized:
            return ""
        if (
            len(normalized) > CORPUS_INVENTORY_MAX_EXTENSION_CHARS
            or not all(
                character.isascii() and character.isalnum()
                for character in normalized
            )
        ):
            return None
        return normalized

    def _exact_document_candidates(
        self,
        corpus_id: str,
        *,
        document_ids: list[str],
        include_remote: bool,
        max_file_bytes: int,
        remote_only: bool,
    ) -> tuple[list[dict], dict[str, dict], dict]:
        placeholders = ",".join("?" for _ in document_ids)
        with corpus_read_connection(self.data_root, corpus_id) as connection:
            rows = connection.execute(
                f"""
                SELECT d.*, r.source_size AS revision_source_size,
                       r.source_modified_ns AS revision_source_modified_ns,
                       r.source_changed_ns AS revision_source_changed_ns,
                       r.source_device AS revision_source_device,
                       r.source_inode AS revision_source_inode,
                       p.projection_id AS active_projection_id,
                       p.adapter_id AS projection_adapter_id,
                       p.adapter_version AS projection_adapter_version,
                       p.config_hash AS projection_config_hash
                FROM documents d
                LEFT JOIN revisions r ON r.revision_id = d.current_revision_id
                LEFT JOIN extraction_projections p
                  ON p.revision_id = d.current_revision_id AND p.is_active = 1
                WHERE d.document_id IN ({placeholders})
                """,
                document_ids,
            ).fetchall()
        documents = {row["document_id"]: dict(row) for row in rows}
        candidates: list[dict] = []
        outcomes: dict[str, dict] = {}
        skipped = {
            "current": 0,
            "remote": 0,
            "local": 0,
            "too_large": 0,
            "not_selected": 0,
            "unknown": 0,
            "deleted": 0,
            "unsupported": 0,
            "max_files_deferred": 0,
            "max_bytes_deferred": 0,
        }

        def record_outcome(
            document_id: str,
            outcome: str,
            document: dict | None = None,
        ) -> None:
            item: dict[str, object] = {
                "document_id": document_id,
                "outcome": outcome,
            }
            if document is not None:
                extension = self._safe_selection_extension(document)
                if extension is not None:
                    item["extension"] = extension
            outcomes[document_id] = item

        for document_id in document_ids:
            document = documents.get(document_id)
            if document is None:
                skipped["unknown"] += 1
                record_outcome(document_id, "unknown")
                continue
            if document["deleted_at"] is not None:
                skipped["deleted"] += 1
                record_outcome(document_id, "deleted", document)
                continue
            if document["eligibility_state"] != "supported":
                skipped["unsupported"] += 1
                record_outcome(document_id, "unsupported", document)
                continue
            document_index_state, _ = self._document_index_state(document)
            if document_index_state == "current":
                skipped["current"] += 1
                record_outcome(document_id, "current", document)
                continue
            if document["logical_size"] > max_file_bytes:
                skipped["too_large"] += 1
                record_outcome(document_id, "too_large", document)
                continue
            if remote_only and not document["is_dataless"]:
                skipped["local"] += 1
                record_outcome(document_id, "remote_disallowed", document)
                continue
            if document["is_dataless"] and not include_remote:
                skipped["remote"] += 1
                record_outcome(document_id, "remote_disallowed", document)
                continue
            candidates.append(document)
        return candidates, outcomes, skipped

    def _pending_state_summary(
        self,
        corpus_id: str,
        *,
        max_file_bytes: int,
    ) -> dict:
        with corpus_read_connection(self.data_root, corpus_id) as connection:
            rows = connection.execute(
                """
                SELECT d.*, r.source_size AS revision_source_size,
                       r.source_modified_ns AS revision_source_modified_ns,
                       r.source_changed_ns AS revision_source_changed_ns,
                       r.source_device AS revision_source_device,
                       r.source_inode AS revision_source_inode,
                       p.projection_id AS active_projection_id,
                       p.adapter_id AS projection_adapter_id,
                       p.adapter_version AS projection_adapter_version,
                       p.config_hash AS projection_config_hash
                FROM documents d
                LEFT JOIN revisions r ON r.revision_id = d.current_revision_id
                LEFT JOIN extraction_projections p
                  ON p.revision_id = d.current_revision_id AND p.is_active = 1
                WHERE d.deleted_at IS NULL
                  AND d.eligibility_state = 'supported'
                """
            ).fetchall()
        result = {
            "remaining": 0,
            "refreshable": 0,
            "pending_remote": 0,
            "too_large": 0,
        }
        for row in rows:
            document = dict(row)
            document_index_state, _ = self._document_index_state(document)
            if document_index_state == "current":
                continue
            result["remaining"] += 1
            if document["is_dataless"]:
                result["pending_remote"] += 1
            if document["logical_size"] > max_file_bytes:
                result["too_large"] += 1
            if (
                not document["is_dataless"]
                and document["logical_size"] <= max_file_bytes
            ):
                result["refreshable"] += 1
        return result

    def _ingest_locked(
        self,
        *,
        corpus: dict,
        paths: RuntimePaths,
        max_files: int,
        max_bytes: int,
        max_file_bytes: int,
        include_remote: bool,
        remote_only: bool,
        document_ids: list[str] | None,
        timeout_seconds: float,
    ) -> dict:
        corpus_id = corpus["corpus_id"]
        abandoned_staging_cleanup = cleanup_abandoned_staging(paths)
        exact_selection = document_ids is not None
        if document_ids is None:
            pending, skipped = self._pending_documents(
                corpus_id,
                include_remote=include_remote,
                max_file_bytes=max_file_bytes,
                remote_only=remote_only,
            )
            outcome_by_id: dict[str, dict] = {}
        else:
            pending, outcome_by_id, skipped = self._exact_document_candidates(
                corpus_id,
                document_ids=document_ids,
                include_remote=include_remote,
                max_file_bytes=max_file_bytes,
                remote_only=remote_only,
            )
        selected: list[dict] = []
        selected_bytes = 0
        for document in pending:
            if len(selected) >= max_files:
                if exact_selection:
                    skipped["max_files_deferred"] += 1
                    deferred = {
                        "document_id": document["document_id"],
                        "outcome": "max_files_deferred",
                    }
                    extension = self._safe_selection_extension(document)
                    if extension is not None:
                        deferred["extension"] = extension
                    outcome_by_id[document["document_id"]] = deferred
                    continue
                break
            if selected_bytes + document["logical_size"] > max_bytes:
                if exact_selection:
                    skipped["max_bytes_deferred"] += 1
                    deferred = {
                        "document_id": document["document_id"],
                        "outcome": "max_bytes_deferred",
                    }
                    extension = self._safe_selection_extension(document)
                    if extension is not None:
                        deferred["extension"] = extension
                    outcome_by_id[document["document_id"]] = deferred
                continue
            selected.append(document)
            selected_bytes += document["logical_size"]

        attempted_results: list[dict] = []
        for document in selected:
            try:
                result = self._ingest_document(
                    corpus=corpus,
                    document=document,
                    allow_hydration=include_remote,
                    maximum_bytes=document["logical_size"],
                    timeout_seconds=timeout_seconds,
                )
            except Exception as exc:
                result = {
                    "document_id": document["document_id"],
                    "relative_path": document["relative_path"],
                    "state": "failed",
                    "error_code": getattr(exc, "code", "unexpected_error"),
                    "message": str(exc),
                    "details": getattr(exc, "details", {}),
                }
                cleanup_failure = getattr(exc, "source_copy_cleanup", None)
                if isinstance(cleanup_failure, dict):
                    result["source_copy_retention"] = "cleanup_failed"
                    result["source_copy_cleanup"] = cleanup_failure
            result["outcome"] = (
                "refreshed"
                if result["state"] in {"indexed", "already_indexed"}
                else "selected"
            )
            attempted_results.append(result)
            if exact_selection:
                outcome_by_id[document["document_id"]] = result
        results = (
            [outcome_by_id[document_id] for document_id in document_ids]
            if document_ids is not None
            else attempted_results
        )
        return {
            "policy": {
                "include_remote": include_remote,
                "remote_only": remote_only,
                "selection_mode": "exact" if exact_selection else "automatic",
                "document_ids": document_ids or [],
                "max_files": max_files,
                "max_bytes": max_bytes,
                "max_file_bytes": max_file_bytes,
                "timeout_seconds": timeout_seconds,
                "concurrency": 1,
                "abandoned_staging_cleanup": abandoned_staging_cleanup,
            },
            "selected_files": len(selected),
            "selected_logical_bytes": selected_bytes,
            "skipped": skipped,
            "results": results,
            "summary": {
                **{
                    state: sum(
                        1
                        for result in attempted_results
                        if result["state"] == state
                    )
                    for state in ("indexed", "already_indexed", "failed")
                },
                "source_copy_cleanup_failed": sum(
                    1
                    for result in attempted_results
                    if result.get("source_copy_cleanup", {}).get("state") == "failed"
                ),
            },
        }

    def ingest(
        self,
        corpus_id: str,
        *,
        max_files: int = 10,
        max_bytes: int = 50 * 1024 * 1024,
        max_file_bytes: int = 25 * 1024 * 1024,
        include_remote: bool = False,
        remote_only: bool = False,
        document_ids: list[str] | None = None,
        timeout_seconds: float = 120,
    ) -> dict:
        corpus_id = normalize_corpus_id(corpus_id)
        corpus = get_corpus(self.data_root, corpus_id)
        _validate_ingest_budgets(
            max_files=max_files,
            max_bytes=max_bytes,
            max_file_bytes=max_file_bytes,
            timeout_seconds=timeout_seconds,
        )
        if document_ids is not None and not document_ids:
            raise InvalidRequestError(
                "explicit ingest document selection must not be empty",
                details={"minimum_document_ids": 1},
            )
        if document_ids is not None and (
            len(document_ids) > _MAX_INGEST_DOCUMENT_IDS
            or any(
                not isinstance(document_id, str)
                or not document_id
                or len(document_id) > 200
                for document_id in document_ids
            )
        ):
            raise BudgetExceededError(
                "ingest document selection exceeds the supported request bounds",
                details={
                    "document_id_count": len(document_ids),
                    "max_document_ids": _MAX_INGEST_DOCUMENT_IDS,
                    "max_document_id_chars": 200,
                },
            )
        if document_ids is not None and len(set(document_ids)) != len(document_ids):
            raise InvalidRequestError(
                "ingest document ids must be unique",
                details={"reason": "duplicate_document_ids"},
            )
        if remote_only and not include_remote:
            raise BudgetExceededError(
                "remote-only selection requires explicit --include-remote",
                details={"remote_only": True, "include_remote": False},
            )
        paths = self._paths(corpus_id)
        paths.ensure()
        with writer_lock(paths.corpus_root / "writer.lock"):
            result = self._ingest_locked(
                corpus=corpus,
                paths=paths,
                max_files=max_files,
                max_bytes=max_bytes,
                max_file_bytes=max_file_bytes,
                include_remote=include_remote,
                remote_only=remote_only,
                document_ids=document_ids,
                timeout_seconds=timeout_seconds,
            )
            result["snapshot"] = self._publish_snapshot(corpus_id)
        return {
            "corpus_id": corpus_id,
            **result,
        }

    def _validate_sync_boundary(self, corpus: dict) -> RuntimePaths:
        source_root = Path(corpus["source_root"]).expanduser().resolve(strict=False)
        paths = self._paths(corpus["corpus_id"])
        runtime_roots = {
            "data_root": self.data_root.expanduser().resolve(strict=False),
            "staging_root": paths.staging.expanduser().resolve(strict=False),
        }
        overlaps = [
            name
            for name, runtime_root in runtime_roots.items()
            if is_within(runtime_root, source_root)
            or is_within(source_root, runtime_root)
        ]
        if overlaps:
            raise SourceBoundaryError(
                "runtime data and source roots must not overlap",
                details={
                    "source_root": str(source_root),
                    "runtime_roots": {
                        name: str(runtime_roots[name]) for name in overlaps
                    },
                },
            )
        return paths

    def sync(
        self,
        corpus_id: str,
        *,
        max_files: int = 10,
        max_bytes: int = 50 * 1024 * 1024,
        max_file_bytes: int = 25 * 1024 * 1024,
        include_remote: bool = False,
        timeout_seconds: float = 120,
    ) -> dict:
        """Scan metadata and refresh only documents whose source index is pending."""

        _validate_ingest_budgets(
            max_files=max_files,
            max_bytes=max_bytes,
            max_file_bytes=max_file_bytes,
            timeout_seconds=timeout_seconds,
        )
        corpus_id = normalize_corpus_id(corpus_id)
        corpus = get_corpus(self.data_root, corpus_id)
        paths = self._validate_sync_boundary(corpus)
        paths.ensure()
        with writer_lock(paths.corpus_root / "writer.lock"):
            scan = dict(scan_corpus(self.data_root, corpus_id))
            inventory_complete = bool(scan.get("observation_complete"))
            ingested: dict | None = None
            pending_state: dict | None = None
            if inventory_complete:
                ingested = self._ingest_locked(
                    corpus=corpus,
                    paths=paths,
                    max_files=max_files,
                    max_bytes=max_bytes,
                    max_file_bytes=max_file_bytes,
                    include_remote=include_remote,
                    remote_only=False,
                    document_ids=None,
                    timeout_seconds=timeout_seconds,
                )
                pending_state = self._pending_state_summary(
                    corpus_id,
                    max_file_bytes=max_file_bytes,
                )
            snapshot = self._publish_snapshot(corpus_id)

        scan.pop("source_root", None)
        change_counts = dict(scan.get("change_counts", {}))
        inventory = {
            "scan_id": scan["scan_id"],
            "inventory_complete": inventory_complete,
            "directories": scan["directories"],
            "files": scan["files"],
            "dataless_files": scan["dataless_files"],
            "logical_bytes": scan["logical_bytes"],
            "supported_files": int(
                scan.get("eligibility_counts", {}).get("supported", 0)
            ),
            "completeness_failure_count": scan["completeness_failure_count"],
            "changed_documents": int(scan.get("changed_documents", 0)),
            "change_counts": change_counts,
        }
        policy = {
            "include_remote": include_remote,
            "max_files": max_files,
            "max_bytes": max_bytes,
            "max_file_bytes": max_file_bytes,
            "timeout_seconds": timeout_seconds,
            "concurrency": 1,
        }
        if not inventory_complete:
            summary = {
                "added": int(change_counts.get("added", 0)),
                "changed": int(scan.get("changed_documents", 0)),
                "reappeared": int(change_counts.get("reappeared", 0)),
                "deleted": int(change_counts.get("deleted", 0)),
                "indexed": 0,
                "reused": 0,
                "failed": 0,
                "pending_remote": None,
                "too_large": None,
                "remaining": None,
            }
            return {
                "corpus_id": corpus_id,
                "state": "scan_incomplete",
                "policy": policy,
                "inventory": inventory,
                "refresh": {
                    "state": "skipped",
                    "reason": "incomplete_metadata_scan",
                    "selected_files": 0,
                    "selected_logical_bytes": 0,
                    "results": [],
                },
                "pending": {
                    "remaining": None,
                    "refreshable": None,
                    "pending_remote": None,
                    "too_large": None,
                },
                "summary": summary,
                "snapshot": snapshot,
            }

        if ingested is None or pending_state is None:
            raise AssertionError("complete sync did not produce refresh state")
        refreshable = int(pending_state["refreshable"])
        pending_remote = int(pending_state["pending_remote"])
        too_large = int(pending_state["too_large"])
        remaining = int(pending_state["remaining"])
        refresh_summary = ingested["summary"]
        failed = int(refresh_summary["failed"])
        state = "complete" if remaining == 0 and failed == 0 else "pending"
        summary = {
            "added": int(change_counts.get("added", 0)),
            "changed": int(scan.get("changed_documents", 0)),
            "reappeared": int(change_counts.get("reappeared", 0)),
            "deleted": int(change_counts.get("deleted", 0)),
            "indexed": int(refresh_summary["indexed"]),
            "reused": int(refresh_summary["already_indexed"]),
            "failed": failed,
            "pending_remote": pending_remote,
            "too_large": too_large,
            "remaining": remaining,
        }
        return {
            "corpus_id": corpus_id,
            "state": state,
            "policy": policy,
            "inventory": inventory,
            "refresh": {
                "state": "completed" if failed == 0 else "completed_with_failures",
                "selected_files": ingested["selected_files"],
                "selected_logical_bytes": ingested["selected_logical_bytes"],
                "skipped": ingested["skipped"],
                "results": ingested["results"],
                "source_copy_cleanup_failed": refresh_summary[
                    "source_copy_cleanup_failed"
                ],
                "abandoned_staging_cleanup": ingested["policy"][
                    "abandoned_staging_cleanup"
                ],
            },
            "pending": {
                "remaining": remaining,
                "refreshable": refreshable,
                "pending_remote": pending_remote,
                "too_large": too_large,
            },
            "summary": summary,
            "snapshot": snapshot,
        }

    def _publish_snapshot(self, corpus_id: str) -> dict:
        import unicodedata

        source_root = Path(get_corpus(self.data_root, corpus_id)["source_root"])
        with corpus_connection(self.data_root, corpus_id) as connection:
            rows = connection.execute(
                """
                SELECT d.document_id, d.current_revision_id AS revision_id,
                       d.extension, p.projection_id, p.completeness_state,
                       p.adapter_id, p.adapter_version, p.config_hash
                FROM documents d
                JOIN revisions r ON r.revision_id = d.current_revision_id
                JOIN extraction_projections p
                  ON p.revision_id = d.current_revision_id AND p.is_active = 1
                WHERE d.deleted_at IS NULL
                  AND d.eligibility_state = 'supported'
                  AND r.source_size = d.logical_size
                  AND r.source_modified_ns = d.modified_ns
                  AND r.source_changed_ns = d.changed_ns
                  AND r.source_device = d.device
                  AND r.source_inode = d.inode
                ORDER BY d.document_id
                """
            ).fetchall()
            supported_rows = connection.execute(
                """
                SELECT document_id, relative_path_nfc, logical_size, allocated_size,
                       modified_ns, changed_ns, device, inode, is_dataless,
                       residency_state
                FROM documents
                WHERE deleted_at IS NULL AND eligibility_state = 'supported'
                ORDER BY document_id
                """
            ).fetchall()
            supported_count = len(supported_rows)
            inventory_manifest = [
                dict(supported_document) for supported_document in supported_rows
            ]
            inventory_set_hash = hashlib.sha256(
                encode_json(inventory_manifest).encode()
            ).hexdigest()
            mapping = [
                (row["document_id"], row["revision_id"], row["projection_id"]) for row in rows
            ]
            revision_manifest = "\n".join(
                f"{document_id}={revision_id}"
                for document_id, revision_id, _projection_id_value in mapping
            )
            projection_manifest = "\n".join(
                f"{document_id}={revision_id}@{projection_id}"
                for document_id, revision_id, projection_id in mapping
            )
            revision_set_hash = hashlib.sha256(revision_manifest.encode()).hexdigest()
            projection_set_hash = hashlib.sha256(projection_manifest.encode()).hexdigest()
            has_partial_projection = any(
                row["completeness_state"] != "complete" for row in rows
            )
            latest_scan = connection.execute(
                """
                SELECT scan_id, status
                FROM scan_runs
                ORDER BY rowid DESC
                LIMIT 1
                """
            ).fetchone()
            has_incomplete_scan = bool(
                latest_scan is None
                or latest_scan["status"] != "complete"
            )
            completeness_failure_counts: dict[tuple[str, str], int] = {}
            if latest_scan is not None:
                completeness_issues = connection.execute(
                    """
                    SELECT code, details_json, structural_locator_json
                    FROM extraction_issues
                    WHERE scan_id = ?
                      AND stage = 'scan'
                      AND code IN (
                          'source_root_open_failed',
                          'source_root_changed_during_scan',
                          'directory_scan_failed',
                          'directory_changed_during_scan',
                          'directory_open_failed',
                          'scan_resource_exhausted',
                          'scan_permission_denied',
                          'source_root_revalidation_failed',
                          'stat_failed'
                      )
                    ORDER BY rowid
                    """,
                    (latest_scan["scan_id"],),
                ).fetchall()
                for issue in completeness_issues:
                    try:
                        structural_locator = json.loads(
                            issue["structural_locator_json"]
                        )
                    except (TypeError, json.JSONDecodeError):
                        structural_locator = {}
                    relative_path = (
                        structural_locator.get("relative_path")
                        if isinstance(structural_locator, dict)
                        else None
                    )
                    if not isinstance(relative_path, str) or not relative_path:
                        try:
                            details = json.loads(issue["details_json"])
                        except (TypeError, json.JSONDecodeError):
                            details = {}
                        raw_path = details.get("path") if isinstance(details, dict) else None
                        if isinstance(raw_path, str):
                            try:
                                relative_path = Path(raw_path).relative_to(
                                    source_root
                                ).as_posix()
                            except ValueError:
                                relative_path = None
                    if not isinstance(relative_path, str) or not relative_path:
                        relative_path = "__unlocated__"
                    relative_path = unicodedata.normalize("NFC", relative_path)
                    path_segments = relative_path.split("/")
                    if relative_path.startswith("/") or ".." in path_segments:
                        relative_path = "__unlocated__"
                    else:
                        relative_path = (
                            "/".join(
                                segment
                                for segment in path_segments
                                if segment not in ("", ".")
                            )
                            or "."
                        )
                    normalized_locator = {
                        "relative_path": relative_path,
                    }
                    locator_sha256 = hashlib.sha256(
                        encode_json(normalized_locator).encode()
                    ).hexdigest()
                    failure_key = (issue["code"], locator_sha256)
                    completeness_failure_counts[failure_key] = (
                        completeness_failure_counts.get(failure_key, 0) + 1
                    )
            completeness_failure_scope = [
                {
                    "code": code,
                    "locator_sha256": locator_sha256,
                    "count": count,
                }
                for (code, locator_sha256), count in sorted(
                    completeness_failure_counts.items()
                )
            ]
            completeness_failure_scope_fingerprint = hashlib.sha256(
                encode_json(completeness_failure_scope).encode()
            ).hexdigest()
            has_outdated_projection = False
            adapter_expectations = []
            for row in rows:
                descriptor = self.adapter_registry.resolve(row["extension"]).descriptor
                adapter_expectations.append(
                    {
                        "projection_id": row["projection_id"],
                        "adapter_id": descriptor.adapter_id,
                        "adapter_version": descriptor.adapter_version,
                        "config_hash": descriptor.config_hash,
                    }
                )
                projection_identity = (
                    row["adapter_id"],
                    row["adapter_version"],
                    row["config_hash"],
                )
                current_identity = (
                    descriptor.adapter_id,
                    descriptor.adapter_version,
                    descriptor.config_hash,
                )
                if projection_identity != current_identity:
                    has_outdated_projection = True
            coverage_state = (
                "complete"
                if (
                    len(mapping) == supported_count
                    and not has_partial_projection
                    and not has_outdated_projection
                    and not has_incomplete_scan
                )
                else "partial"
            )
            coverage_gaps = []
            if len(mapping) != supported_count:
                coverage_gaps.append("supported_documents_without_usable_projection")
            if has_partial_projection:
                coverage_gaps.append("partial_extraction_projection")
            if has_outdated_projection:
                coverage_gaps.append("outdated_extraction_projection")
            if has_incomplete_scan:
                coverage_gaps.append("incomplete_metadata_scan")
            publication_manifest = {
                "schema_version": 1,
                "extraction_schema_version": EXTRACTION_SCHEMA_VERSION,
                "inventory_set_hash": inventory_set_hash,
                "document_revision_set_hash": revision_set_hash,
                "extraction_projection_set_hash": projection_set_hash,
                "document_count": len(mapping),
                "supported_document_count": supported_count,
                "coverage_state": coverage_state,
                "coverage_gaps": coverage_gaps,
                "completeness_failure_scope": completeness_failure_scope,
                "completeness_failure_scope_fingerprint": (
                    completeness_failure_scope_fingerprint
                ),
                "projection_completeness": [
                    {
                        "projection_id": row["projection_id"],
                        "completeness_state": row["completeness_state"],
                    }
                    for row in rows
                ],
                "adapter_expectations": adapter_expectations,
            }
            publication_state_hash = hashlib.sha256(
                encode_json(publication_manifest).encode()
            ).hexdigest()
            latest_snapshot = connection.execute(
                """
                SELECT snapshot_id
                FROM snapshots
                WHERE state = 'complete'
                ORDER BY rowid DESC
                LIMIT 1
                """
            ).fetchone()
            latest_publication = connection.execute(
                """
                SELECT payload_json
                FROM events
                WHERE event_type = 'snapshot_published'
                ORDER BY rowid DESC
                LIMIT 1
                """
            ).fetchone()
            reuse_latest_snapshot = False
            if latest_snapshot is not None and latest_publication is not None:
                try:
                    latest_payload = json.loads(latest_publication["payload_json"])
                except (TypeError, json.JSONDecodeError):
                    latest_payload = {}
                reuse_latest_snapshot = (
                    latest_payload.get("snapshot_id") == latest_snapshot["snapshot_id"]
                    and latest_payload.get("publication_state_hash")
                    == publication_state_hash
                )

            now = utc_now()
            if reuse_latest_snapshot:
                snapshot_id = latest_snapshot["snapshot_id"]
            else:
                publication_nonce = uuid.uuid4().hex
                snapshot_digest = hashlib.sha256(
                    (
                        f"state={publication_state_hash}\n"
                        f"publication={publication_nonce}"
                    ).encode()
                ).hexdigest()
                snapshot_id = f"snap_{snapshot_digest[:32]}"
                connection.execute(
                    """
                    INSERT INTO snapshots(
                        snapshot_id, state, coverage_state, document_revision_set_hash,
                        extraction_projection_set_hash, document_count,
                        supported_document_count, extraction_schema_version,
                        created_at, completed_at
                    ) VALUES (?, 'complete', ?, ?, ?, ?, ?, ?, ?, ?)
                    """,
                    (
                        snapshot_id,
                        coverage_state,
                        revision_set_hash,
                        projection_set_hash,
                        len(mapping),
                        supported_count,
                        EXTRACTION_SCHEMA_VERSION,
                        now,
                        now,
                    ),
                )
                connection.executemany(
                    """
                    INSERT INTO snapshot_documents(
                        snapshot_id, document_id, revision_id, projection_id
                    ) VALUES (?, ?, ?, ?)
                    """,
                    [
                        (snapshot_id, document_id, revision_id, projection_id)
                        for document_id, revision_id, projection_id in mapping
                    ],
                )
                connection.execute(
                    """
                    INSERT INTO events(
                        event_id, event_type, payload_json, created_at
                    ) VALUES (?, 'snapshot_published', ?, ?)
                    """,
                    (
                        f"event_{uuid.uuid4().hex}",
                        encode_json(
                            {
                                "schema_version": 1,
                                "snapshot_id": snapshot_id,
                                "publication_state_hash": publication_state_hash,
                                "inventory_set_hash": inventory_set_hash,
                                "coverage_gaps": coverage_gaps,
                                "completeness_failure_scope": (
                                    completeness_failure_scope
                                ),
                                "completeness_failure_scope_fingerprint": (
                                    completeness_failure_scope_fingerprint
                                ),
                            }
                        ),
                        now,
                    ),
                )
        return {
            "snapshot_id": snapshot_id,
            "coverage_state": coverage_state,
            "document_count": len(mapping),
            "supported_document_count": supported_count,
            "document_revision_set_hash": revision_set_hash,
            "extraction_projection_set_hash": projection_set_hash,
            "inventory_set_hash": inventory_set_hash,
            "publication_state_hash": publication_state_hash,
            "coverage_gaps": coverage_gaps,
            "completeness_failure_scope": completeness_failure_scope,
            "completeness_failure_scope_fingerprint": (
                completeness_failure_scope_fingerprint
            ),
        }

    def _ingest_document(
        self,
        *,
        corpus: dict,
        document: dict,
        allow_hydration: bool,
        maximum_bytes: int,
        timeout_seconds: float,
    ) -> dict:
        source_root = Path(corpus["source_root"])
        source = Path(document["absolute_path"])
        paths = self._paths(corpus["corpus_id"])
        scanned_key = (
            document["logical_size"],
            document["modified_ns"],
            document["changed_ns"],
            document["device"],
            document["inode"],
        )
        try:
            captured = capture_to_staging(
                paths=paths,
                source_root=source_root,
                source=source,
                allow_hydration=allow_hydration,
                maximum_bytes=maximum_bytes,
                timeout_seconds=timeout_seconds,
                expected_source_identity=scanned_key,
            )
        except SourceBoundaryError as exc:
            if not _source_is_missing(exc):
                raise
            raise SourceUnavailableError(
                "source bytes are unavailable; refresh requires the registered original",
                details={
                    "document_id": document["document_id"],
                    "relative_path": document["relative_path"],
                    "existing_sqlite_index_unchanged": bool(
                        document.get("current_revision_id")
                    ),
                    "source_error": exc.details,
                },
            ) from exc

        result: dict | None = None
        primary_error: BaseException | None = None
        try:
            revision_id = _revision_id(document["document_id"], captured.sha256)
            capture_ref = _ephemeral_capture_ref(captured.sha256)
            adapter = self.adapter_registry.resolve(document["extension"])
            descriptor = adapter.descriptor

            with corpus_connection(self.data_root, corpus["corpus_id"]) as connection:
                existing = connection.execute(
                    """
                    SELECT p.projection_id, p.adapter_id, p.adapter_version, p.config_hash
                    FROM revisions r
                    LEFT JOIN extraction_projections p
                      ON p.revision_id = r.revision_id AND p.is_active = 1
                    WHERE r.revision_id = ?
                    """,
                    (revision_id,),
                ).fetchone()
                if (
                    existing
                    and existing["projection_id"] is not None
                    and existing["adapter_id"] == descriptor.adapter_id
                    and existing["adapter_version"] == descriptor.adapter_version
                    and existing["config_hash"] == descriptor.config_hash
                ):
                    self._reactivate_existing_projection(
                        connection,
                        document=document,
                        revision_id=revision_id,
                        projection_id=existing["projection_id"],
                        captured=captured,
                        blob_ref=capture_ref,
                    )
                    result = {
                        "document_id": document["document_id"],
                        "relative_path": document["relative_path"],
                        "revision_id": revision_id,
                        "projection_id": existing["projection_id"],
                        "sha256": captured.sha256,
                        "state": "already_indexed",
                        "hydrated": captured.hydration_was_required,
                        "source_copy_retention": "ephemeral",
                        "source_copy_cleanup": {"state": "deleted"},
                    }
                    return result

            try:
                extraction = adapter.extract(
                    captured.capture_path,
                    format_id=document["extension"],
                )
            except ExtractionError as exc:
                self._record_failed_extraction(
                    corpus_id=corpus["corpus_id"],
                    document=document,
                    captured=captured,
                    revision_id=revision_id,
                    blob_ref=capture_ref,
                    descriptor=descriptor,
                    error=exc,
                )
                raise

            committed = self._commit_extraction(
                corpus_id=corpus["corpus_id"],
                document=document,
                captured=captured,
                revision_id=revision_id,
                blob_ref=capture_ref,
                extraction=extraction,
            )
            result = {
                "document_id": document["document_id"],
                "relative_path": document["relative_path"],
                "revision_id": revision_id,
                "projection_id": committed["projection_id"],
                "sha256": captured.sha256,
                "state": "indexed",
                "source_units": committed["source_units"],
                "completeness_state": committed["completeness_state"],
                "hydrated": captured.hydration_was_required,
                "native_capture": captured.used_native_helper,
                "bytes_copied": captured.bytes_copied,
                "source_copy_retention": "ephemeral",
                "source_copy_cleanup": {"state": "deleted"},
                "extraction_issues": committed["extraction_issues"],
            }
            return result
        except BaseException as exc:
            primary_error = exc
            raise
        finally:
            try:
                discard_staged_capture(paths, captured)
            except Exception as exc:
                cleanup_failure = {
                    "state": "failed",
                    "error_code": getattr(exc, "code", "unexpected_error"),
                    "message": str(exc),
                    "details": getattr(exc, "details", {}),
                }
                if result is not None:
                    result["source_copy_retention"] = "cleanup_failed"
                    result["source_copy_cleanup"] = cleanup_failure
                elif primary_error is not None:
                    primary_error.source_copy_cleanup = cleanup_failure
                else:
                    raise

    def _set_document_current(
        self,
        connection,
        document: dict,
        revision_id: str,
        captured: CapturedSource,
    ) -> None:
        post_identity = captured.post_identity
        connection.execute(
            """
            UPDATE documents SET
                current_revision_id = ?, logical_size = ?, modified_ns = ?,
                changed_ns = ?, device = ?, inode = ?, mode = ?, flags = ?,
                is_dataless = ?, residency_state = ?, allocated_size = ?,
                last_seen_at = ?
            WHERE document_id = ?
            """,
            (
                revision_id,
                post_identity.size,
                post_identity.modified_ns,
                post_identity.changed_ns,
                post_identity.device,
                post_identity.inode,
                post_identity.mode,
                post_identity.flags,
                int(post_identity.dataless),
                "remote_only" if post_identity.dataless else "resident",
                post_identity.allocated_size,
                utc_now(),
                document["document_id"],
            ),
        )

    def _refresh_revision_observation(
        self,
        connection,
        *,
        revision_id: str,
        captured: CapturedSource,
        blob_ref: str,
    ) -> None:
        now = utc_now()
        connection.execute(
            """
            UPDATE revisions SET
                immutable_blob_ref = ?,
                source_size = ?,
                source_modified_ns = ?,
                source_changed_ns = ?,
                source_device = ?,
                source_inode = ?,
                observed_at = ?,
                captured_at = ?
            WHERE revision_id = ?
            """,
            (
                blob_ref,
                captured.post_identity.size,
                captured.post_identity.modified_ns,
                captured.post_identity.changed_ns,
                captured.post_identity.device,
                captured.post_identity.inode,
                now,
                now,
                revision_id,
            ),
        )

    def _revalidate_current_claim_dependencies(
        self,
        connection,
        *,
        projection_id: str,
    ) -> None:
        connection.execute(
            """
            UPDATE atomic_claims
            SET dependency_state = 'valid'
            WHERE dependency_state = 'stale'
              AND claim_id IN (
                  SELECT e.claim_id
                  FROM evidence_links e
                  JOIN source_units u ON u.unit_id = e.source_unit_id
                  WHERE u.projection_id = ?
              )
              AND NOT EXISTS (
                  SELECT 1
                  FROM evidence_links dependency
                  JOIN source_units dependency_unit
                    ON dependency_unit.unit_id = dependency.source_unit_id
                  JOIN revisions dependency_revision
                    ON dependency_revision.revision_id = dependency_unit.revision_id
                  JOIN documents dependency_document
                    ON dependency_document.document_id = dependency_revision.document_id
                  LEFT JOIN extraction_projections active_projection
                    ON active_projection.revision_id = dependency_unit.revision_id
                   AND active_projection.is_active = 1
                  WHERE dependency.claim_id = atomic_claims.claim_id
                    AND (
                        dependency_document.deleted_at IS NOT NULL
                        OR dependency_document.current_revision_id
                           != dependency_unit.revision_id
                        OR active_projection.projection_id IS NULL
                        OR active_projection.projection_id
                           != dependency_unit.projection_id
                        OR dependency_revision.source_size
                           != dependency_document.logical_size
                        OR dependency_revision.source_modified_ns
                           != dependency_document.modified_ns
                        OR dependency_revision.source_changed_ns
                           != dependency_document.changed_ns
                        OR dependency_revision.source_device
                           != dependency_document.device
                        OR dependency_revision.source_inode
                           != dependency_document.inode
                    )
              )
            """,
            (projection_id,),
        )

    @staticmethod
    def _checkpoint_completes_projection(
        connection,
        *,
        projection_id: str,
        checkpoint_json: str | None,
    ) -> bool:
        if not checkpoint_json:
            return False
        try:
            checkpoint = json.loads(checkpoint_json)
        except (TypeError, json.JSONDecodeError):
            return False
        if (
            not isinstance(checkpoint, dict)
            or type(checkpoint.get("schema_version")) is not int
            or checkpoint.get("schema_version") != 1
            or checkpoint.get("projection_id") != projection_id
            or checkpoint.get("next_ordinal") is not None
            or checkpoint.get("remaining_ordinal_ranges") != []
        ):
            return False
        inventory = connection.execute(
            """
            SELECT COUNT(*) AS unit_count, MIN(ordinal) AS min_ordinal,
                   MAX(ordinal) AS max_ordinal
            FROM source_units
            WHERE projection_id = ?
            """,
            (projection_id,),
        ).fetchone()
        if (
            inventory is None
            or inventory["unit_count"] < 1
            or type(checkpoint.get("total_units")) is not int
            or checkpoint.get("total_units") != inventory["unit_count"]
            or inventory["min_ordinal"] != 1
            or inventory["max_ordinal"] != inventory["unit_count"]
        ):
            return False
        ranges = checkpoint.get("processed_ordinal_ranges")
        if not isinstance(ranges, list) or not ranges:
            return False
        expected_start = 1
        for value in ranges:
            if (
                not isinstance(value, list)
                or len(value) != 2
                or type(value[0]) is not int
                or type(value[1]) is not int
                or value[0] != expected_start
                or value[1] < value[0]
            ):
                return False
            expected_start = value[1] + 1
        return expected_start == inventory["unit_count"] + 1

    def _reactivate_existing_projection(
        self,
        connection,
        *,
        document: dict,
        revision_id: str,
        projection_id: str,
        captured: CapturedSource,
        blob_ref: str,
    ) -> None:
        previous_revision_id = document.get("current_revision_id")
        now = utc_now()
        self._refresh_revision_observation(
            connection,
            revision_id=revision_id,
            captured=captured,
            blob_ref=blob_ref,
        )
        self._set_document_current(connection, document, revision_id, captured)
        if previous_revision_id and previous_revision_id != revision_id:
            connection.execute(
                """
                UPDATE interpretation_queue
                SET state = 'stale',
                    reason = 'superseded_source_revision',
                    updated_at = ?
                WHERE revision_id = ? AND state != 'stale'
                """,
                (now, previous_revision_id),
            )
            connection.execute(
                """
                UPDATE atomic_claims
                SET dependency_state = 'stale'
                WHERE claim_id IN (
                    SELECT claim_id FROM evidence_links
                    WHERE source_revision_id = ?
                )
                  AND dependency_state = 'valid'
                """,
                (previous_revision_id,),
            )
        reason = (
            "reactivated_source_revision"
            if previous_revision_id and previous_revision_id != revision_id
            else "source_observation_reconfirmed"
        )
        if self.maintain_legacy_semantic_cache:
            existing_queue = connection.execute(
                """
                SELECT checkpoint_json
                FROM interpretation_queue
                WHERE projection_id = ?
                """,
                (projection_id,),
            ).fetchone()
            queue_state = (
                "complete"
                if existing_queue is not None
                and self._checkpoint_completes_projection(
                    connection,
                    projection_id=projection_id,
                    checkpoint_json=existing_queue["checkpoint_json"],
                )
                else "pending"
            )
            connection.execute(
                """
                INSERT INTO interpretation_queue(
                    queue_id, document_id, revision_id, projection_id,
                    state, reason, created_at, updated_at
                ) VALUES (?, ?, ?, ?, ?, ?, ?, ?)
                ON CONFLICT(projection_id) DO UPDATE SET
                    state = excluded.state,
                    reason = excluded.reason,
                    updated_at = excluded.updated_at
                """,
                (
                    f"queue_{uuid.uuid4().hex}",
                    document["document_id"],
                    revision_id,
                    projection_id,
                    queue_state,
                    reason,
                    now,
                    now,
                ),
            )
        self._revalidate_current_claim_dependencies(
            connection,
            projection_id=projection_id,
        )
        connection.execute(
            """
            INSERT INTO events(
                event_id, event_type, document_id, revision_id, payload_json, created_at
            ) VALUES (?, 'revision_reactivated', ?, ?, ?, ?)
            """,
            (
                f"event_{uuid.uuid4().hex}",
                document["document_id"],
                revision_id,
                encode_json(
                    {
                        "projection_id": projection_id,
                        "previous_revision_id": previous_revision_id,
                        "reason": reason,
                    }
                ),
                now,
            ),
        )

    def _insert_revision(
        self,
        connection,
        *,
        document: dict,
        captured: CapturedSource,
        revision_id: str,
        blob_ref: str,
        extraction_state: str,
        extractor_version: str,
    ) -> str | None:
        predecessor = document.get("current_revision_id")
        now = utc_now()
        connection.execute(
            """
            INSERT INTO revisions(
                revision_id, document_id, sha256, immutable_blob_ref,
                source_size, source_modified_ns, source_changed_ns,
                source_device, source_inode,
                observed_at, captured_at, extraction_state, extractor_version,
                predecessor_revision_id
            ) VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
            ON CONFLICT(revision_id) DO UPDATE SET
                immutable_blob_ref = excluded.immutable_blob_ref,
                source_size = excluded.source_size,
                source_modified_ns = excluded.source_modified_ns,
                source_changed_ns = excluded.source_changed_ns,
                source_device = excluded.source_device,
                source_inode = excluded.source_inode,
                observed_at = excluded.observed_at,
                captured_at = excluded.captured_at,
                extraction_state = CASE
                    WHEN revisions.extraction_state = 'complete'
                         AND excluded.extraction_state = 'failed'
                    THEN revisions.extraction_state
                    ELSE excluded.extraction_state
                END,
                extractor_version = CASE
                    WHEN revisions.extraction_state = 'complete'
                         AND excluded.extraction_state = 'failed'
                    THEN revisions.extractor_version
                    ELSE excluded.extractor_version
                END
            """,
            (
                revision_id,
                document["document_id"],
                captured.sha256,
                blob_ref,
                captured.post_identity.size,
                captured.post_identity.modified_ns,
                captured.post_identity.changed_ns,
                captured.post_identity.device,
                captured.post_identity.inode,
                now,
                now,
                extraction_state,
                extractor_version,
                predecessor if predecessor != revision_id else None,
            ),
        )
        self._set_document_current(connection, document, revision_id, captured)
        return predecessor

    def _record_failed_extraction(
        self,
        *,
        corpus_id: str,
        document: dict,
        captured: CapturedSource,
        revision_id: str,
        blob_ref: str,
        descriptor: AdapterDescriptor,
        error: ExtractionError,
    ) -> None:
        with corpus_connection(self.data_root, corpus_id) as connection:
            self._insert_revision(
                connection,
                document=document,
                captured=captured,
                revision_id=revision_id,
                blob_ref=blob_ref,
                extraction_state="failed",
                extractor_version=descriptor.adapter_version,
            )
            now = utc_now()
            attempt_id = f"attempt_{uuid.uuid4().hex}"
            connection.execute(
                """
                INSERT INTO extraction_attempts(
                    attempt_id, revision_id, adapter_id, adapter_version,
                    config_hash, state, error_json, started_at, completed_at
                ) VALUES (?, ?, ?, ?, ?, 'failed', ?, ?, ?)
                """,
                (
                    attempt_id,
                    revision_id,
                    descriptor.adapter_id,
                    descriptor.adapter_version,
                    descriptor.config_hash,
                    encode_json(
                        {
                            "code": error.code,
                            "message": str(error),
                            "details": error.details,
                        }
                    ),
                    now,
                    now,
                ),
            )
            structural_locator, locator_key = _issue_locator(error.details)
            connection.execute(
                """
                INSERT INTO extraction_issues(
                    issue_id, document_id, revision_id, attempt_id,
                    stage, severity, code, message, details_json,
                    structural_locator_json, locator_key, lifecycle_state, created_at
                ) VALUES (?, ?, ?, ?, 'extract', 'error', ?, ?, ?, ?, ?, 'active', ?)
                """,
                (
                    f"issue_{uuid.uuid4().hex}",
                    document["document_id"],
                    revision_id,
                    attempt_id,
                    error.code,
                    str(error),
                    encode_json(error.details),
                    encode_json(structural_locator),
                    locator_key,
                    now,
                ),
            )

    def _commit_extraction(
        self,
        *,
        corpus_id: str,
        document: dict,
        captured: CapturedSource,
        revision_id: str,
        blob_ref: str,
        extraction: ExtractionEnvelope,
    ) -> dict:
        descriptor = extraction.descriptor
        adapter_id = descriptor.adapter_id
        result_manifest_hash = extraction.manifest_hash
        projection_id = _projection_id(
            revision_id,
            adapter_id,
            descriptor.adapter_version,
            descriptor.config_hash,
            result_manifest_hash,
        )
        completeness_state = extraction.completeness
        capability_manifest = descriptor.capabilities.to_dict()
        content_hashes = [
            hashlib.sha256(unit.content.encode("utf-8")).hexdigest() for unit in extraction.units
        ]
        unit_ids = [
            _unit_id(projection_id, ordinal, content_hash)
            for ordinal, content_hash in enumerate(content_hashes, start=1)
        ]
        with corpus_connection(self.data_root, corpus_id) as connection:
            predecessor = self._insert_revision(
                connection,
                document=document,
                captured=captured,
                revision_id=revision_id,
                blob_ref=blob_ref,
                extraction_state="complete",
                extractor_version=descriptor.adapter_version,
            )
            now = utc_now()
            attempt_id = f"attempt_{uuid.uuid4().hex}"
            connection.execute(
                """
                INSERT INTO extraction_attempts(
                    attempt_id, revision_id, adapter_id, adapter_version,
                    config_hash, state, started_at
                ) VALUES (?, ?, ?, ?, ?, 'running', ?)
                """,
                (
                    attempt_id,
                    revision_id,
                    adapter_id,
                    descriptor.adapter_version,
                    descriptor.config_hash,
                    now,
                ),
            )
            old_projection_ids = [
                row["projection_id"]
                for row in connection.execute(
                    """
                    SELECT projection_id
                    FROM extraction_projections
                    WHERE revision_id = ? AND is_active = 1 AND projection_id != ?
                    """,
                    (revision_id, projection_id),
                )
            ]
            existing_projection = connection.execute(
                """
                SELECT projection_id
                FROM extraction_projections
                WHERE projection_id = ?
                """,
                (projection_id,),
            ).fetchone()
            if existing_projection is None:
                connection.execute(
                    """
                    INSERT INTO extraction_projections(
                        projection_id, revision_id, adapter_id, adapter_version,
                        config_hash, result_manifest_hash, completeness_state,
                        capability_manifest_json, assurance_state, is_active, created_at
                    ) VALUES (?, ?, ?, ?, ?, ?, ?, ?, 'declared', 0, ?)
                    """,
                    (
                        projection_id,
                        revision_id,
                        adapter_id,
                        descriptor.adapter_version,
                        descriptor.config_hash,
                        result_manifest_hash,
                        completeness_state,
                        encode_json(capability_manifest),
                        now,
                    ),
                )

                for index, unit in enumerate(extraction.units):
                    ordinal = index + 1
                    unit_id = unit_ids[index]
                    unit_payload = unit.to_dict()
                    structure = unit_payload["structure_path"]
                    anchor = {
                        "schema_version": 2,
                        "extraction_schema_version": EXTRACTION_SCHEMA_VERSION,
                        "document_id": document["document_id"],
                        "revision_id": revision_id,
                        "projection_id": projection_id,
                        "content_hash": captured.sha256,
                        "canonical_locator": document["relative_path"],
                        "absolute_path": document["absolute_path"],
                        "structural_locator": structure,
                        "source_span": _source_span(structure),
                        "surface_open_target": document["absolute_path"],
                    }
                    previous_unit_id = unit_ids[index - 1] if index > 0 else None
                    next_unit_id = unit_ids[index + 1] if index + 1 < len(unit_ids) else None
                    connection.execute(
                        """
                        INSERT INTO source_units(
                            unit_id, revision_id, projection_id, ordinal, unit_type,
                            structure_path_json, source_anchor_json, normalized_content,
                            content_sha256, previous_unit_id, next_unit_id,
                            extraction_issues_json, derivation_method, geometry_json,
                            confidence, quality_flags_json
                        ) VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
                        """,
                        (
                            unit_id,
                            revision_id,
                            projection_id,
                            ordinal,
                            unit.unit_type,
                            encode_json(structure),
                            encode_json(anchor),
                            unit.content,
                            content_hashes[index],
                            previous_unit_id,
                            next_unit_id,
                            encode_json(unit_payload["issues"]),
                            unit.derivation_method,
                            encode_json(unit_payload["geometry"]),
                            unit.confidence,
                            encode_json(unit_payload["quality_flags"]),
                        ),
                    )
                    connection.execute(
                        """
                        INSERT INTO source_units_fts(
                            unit_id, document_id, relative_path,
                            structure_path, normalized_content
                        ) VALUES (?, ?, ?, ?, ?)
                        """,
                        (
                            unit_id,
                            document["document_id"],
                            document["relative_path"],
                            json.dumps(structure, ensure_ascii=False),
                            unit.content,
                        ),
                    )

            connection.execute(
                """
                UPDATE extraction_projections
                SET is_active = 0
                WHERE revision_id = ? AND projection_id != ?
                """,
                (revision_id, projection_id),
            )
            connection.execute(
                "UPDATE extraction_projections SET is_active = 1 WHERE projection_id = ?",
                (projection_id,),
            )
            if old_projection_ids:
                placeholders = ",".join("?" for _ in old_projection_ids)
                connection.execute(
                    f"""
                    UPDATE interpretation_queue
                    SET state = 'stale',
                        reason = 'superseded_extraction_projection',
                        updated_at = ?
                    WHERE projection_id IN ({placeholders}) AND state != 'stale'
                    """,
                    (now, *old_projection_ids),
                )
                connection.execute(
                    f"""
                    UPDATE extraction_issues
                    SET lifecycle_state = 'superseded'
                    WHERE projection_id IN ({placeholders})
                      AND lifecycle_state = 'active'
                    """,
                    old_projection_ids,
                )
                connection.execute(
                    f"""
                    UPDATE atomic_claims
                    SET dependency_state = 'stale'
                    WHERE claim_id IN (
                        SELECT e.claim_id
                        FROM evidence_links e
                        JOIN source_units u ON u.unit_id = e.source_unit_id
                        WHERE u.projection_id IN ({placeholders})
                    )
                      AND dependency_state = 'valid'
                    """,
                    old_projection_ids,
                )
            if predecessor and predecessor != revision_id:
                connection.execute(
                    """
                    UPDATE interpretation_queue
                    SET state = 'stale', reason = 'superseded_source_revision', updated_at = ?
                    WHERE revision_id = ? AND state != 'stale'
                    """,
                    (now, predecessor),
                )
                connection.execute(
                    """
                    UPDATE atomic_claims
                    SET dependency_state = 'stale'
                    WHERE claim_id IN (
                        SELECT claim_id FROM evidence_links
                        WHERE source_revision_id = ?
                    )
                      AND dependency_state = 'valid'
                    """,
                    (predecessor,),
                )
            if self.maintain_legacy_semantic_cache:
                connection.execute(
                    """
                    INSERT INTO interpretation_queue(
                        queue_id, document_id, revision_id, projection_id,
                        state, reason, created_at, updated_at
                    ) VALUES (?, ?, ?, ?, 'pending', ?, ?, ?)
                    ON CONFLICT(projection_id) DO UPDATE SET
                        state = 'pending',
                        reason = excluded.reason,
                        updated_at = excluded.updated_at
                    """,
                    (
                        f"queue_{uuid.uuid4().hex}",
                        document["document_id"],
                        revision_id,
                        projection_id,
                        (
                            "new_source_revision"
                            if predecessor and predecessor != revision_id
                            else "new_extraction_projection"
                        ),
                        now,
                        now,
                    ),
                )
            connection.execute(
                """
                UPDATE extraction_issues
                SET lifecycle_state = 'resolved'
                WHERE revision_id = ?
                  AND projection_id IS NULL
                  AND stage = 'extract'
                  AND lifecycle_state = 'active'
                """,
                (revision_id,),
            )
            registry_issues = [
                (issue, issue.to_dict()) for issue in extraction.issues
            ]
            for ordinal, unit in enumerate(extraction.units, start=1):
                structure = unit.to_dict()["structure_path"]
                for issue in unit.issues:
                    issue_payload = issue.to_dict()
                    issue_payload["structural_locator"] = structure
                    issue_payload["details"] = {
                        **issue_payload.get("details", {}),
                        "unit_ordinal": ordinal,
                        "unit_type": unit.unit_type,
                    }
                    registry_issues.append((issue, issue_payload))
            for issue, issue_payload in registry_issues:
                structural_locator, locator_key = _issue_locator(issue_payload)
                connection.execute(
                    """
                    INSERT INTO extraction_issues(
                        issue_id, document_id, revision_id, attempt_id, projection_id,
                        stage, severity, code, message, details_json,
                        structural_locator_json, locator_key, lifecycle_state, created_at
                    ) VALUES (?, ?, ?, ?, ?, 'extract', ?, ?, ?, ?, ?, ?, 'active', ?)
                    ON CONFLICT(projection_id, stage, code, locator_key)
                    WHERE projection_id IS NOT NULL AND lifecycle_state = 'active'
                    DO UPDATE SET
                        attempt_id = excluded.attempt_id,
                        severity = excluded.severity,
                        message = excluded.message,
                        details_json = excluded.details_json
                    """,
                    (
                        f"issue_{uuid.uuid4().hex}",
                        document["document_id"],
                        revision_id,
                        attempt_id,
                        projection_id,
                        issue.severity,
                        issue.code,
                        issue.message,
                        encode_json(issue_payload),
                        encode_json(structural_locator),
                        locator_key,
                        now,
                    ),
                )
            connection.execute(
                """
                UPDATE extraction_attempts
                SET state = 'succeeded', projection_id = ?, completed_at = ?
                WHERE attempt_id = ?
                """,
                (projection_id, now, attempt_id),
            )
            connection.execute(
                """
                INSERT INTO events(
                    event_id, event_type, document_id, revision_id, payload_json, created_at
                ) VALUES (?, 'revision_extracted', ?, ?, ?, ?)
                """,
                (
                    f"event_{uuid.uuid4().hex}",
                    document["document_id"],
                    revision_id,
                    encode_json(
                        {
                            "projection_id": projection_id,
                            "source_units": len(extraction.units),
                            "completeness_state": completeness_state,
                        }
                    ),
                    now,
                ),
            )
        return {
            "projection_id": projection_id,
            "source_units": len(extraction.units),
            "completeness_state": completeness_state,
            "extraction_issues": len(registry_issues),
        }

    def search(self, corpus_id: str, query: str, *, limit: int = 20) -> dict:
        if not 1 <= limit <= _MAX_SEARCH_RESULTS:
            raise BudgetExceededError(
                "search limit must be between 1 and 200",
                details={"limit": limit, "maximum": _MAX_SEARCH_RESULTS},
            )
        if len(query) > 2_000:
            raise BudgetExceededError(
                "search query must contain at most 2000 characters",
                details={"query_chars": len(query), "maximum": 2_000},
            )
        normalized = query.strip()
        if not normalized:
            return {
                "query": query,
                "strategy": "lexical_candidate_acquisition",
                "query_mode": "exact_phrase_fts",
                "zero_results_establish_absence": False,
                "candidates": [],
                "count": 0,
            }
        fts_query = '"' + normalized.replace('"', '""') + '"'
        with corpus_read_connection(self.data_root, corpus_id) as connection:
            connection.create_function(
                "corpus_projection_is_current",
                4,
                self._projection_uses_current_adapter,
                deterministic=True,
            )
            rows = connection.execute(
                """
                SELECT f.unit_id, f.document_id, f.relative_path, f.structure_path,
                       instr(u.normalized_content, ?) AS literal_position,
                       LENGTH(CAST(u.normalized_content AS BLOB))
                           AS source_content_bytes,
                       bm25(source_units_fts) AS lexical_score,
                       u.revision_id, u.projection_id, u.unit_type,
                       u.derivation_method, u.confidence, u.quality_flags_json,
                       u.source_anchor_json, u.trust_lineage,
                       p.completeness_state
                FROM source_units_fts f
                JOIN source_units u ON u.unit_id = f.unit_id
                JOIN extraction_projections p ON p.projection_id = u.projection_id
                JOIN revisions r ON r.revision_id = u.revision_id
                JOIN documents d ON d.document_id = f.document_id
                WHERE source_units_fts MATCH ?
                  AND d.current_revision_id = u.revision_id
                  AND p.is_active = 1
                  AND d.deleted_at IS NULL
                  AND d.eligibility_state = 'supported'
                  AND corpus_projection_is_current(
                          d.extension,
                          p.adapter_id,
                          p.adapter_version,
                          p.config_hash
                      ) = 1
                  AND r.source_size = d.logical_size
                  AND r.source_modified_ns = d.modified_ns
                  AND r.source_changed_ns = d.changed_ns
                  AND r.source_device = d.device
                  AND r.source_inode = d.inode
                ORDER BY bm25(source_units_fts)
                LIMIT ?
                """,
                (
                    normalized,
                    fts_query,
                    limit,
                ),
            ).fetchall()
            excerpt_details = {}
            for row in rows:
                literal_position = int(row["literal_position"])
                excerpt_start = (
                    max(
                        1,
                        literal_position
                        - CORPUS_SEARCH_EXCERPT_CONTEXT_BEFORE_CHARS,
                    )
                    if literal_position > 0
                    else 1
                )
                excerpt_probe = connection.execute(
                    """
                    SELECT substr(normalized_content, ?, ?) AS excerpt_probe
                    FROM source_units
                    WHERE unit_id = ?
                    """,
                    (
                        excerpt_start,
                        CORPUS_SEARCH_EXCERPT_MAX_CHARS + 1,
                        row["unit_id"],
                    ),
                ).fetchone()["excerpt_probe"]
                excerpt_details[row["unit_id"]] = {
                    "excerpt_probe": excerpt_probe,
                    "excerpt_start": excerpt_start,
                    "generation": (
                        "literal_query_window"
                        if literal_position > 0
                        else "bounded_content_prefix"
                    ),
                    "source_content_bytes": row["source_content_bytes"],
                }
        candidates = []
        truncated_excerpt_count = 0
        for row in rows:
            item = dict(row)
            item["structure_path"] = json.loads(item["structure_path"])
            item["source_anchor"] = json.loads(item.pop("source_anchor_json"))
            item["quality_flags"] = json.loads(item.pop("quality_flags_json"))
            item.pop("literal_position")
            item.pop("source_content_bytes")
            excerpt = excerpt_details[item["unit_id"]]
            excerpt_probe = excerpt["excerpt_probe"]
            excerpt_truncated = (
                excerpt["excerpt_start"] > 1
                or excerpt["source_content_bytes"] > len(excerpt_probe.encode())
                or len(excerpt_probe) > CORPUS_SEARCH_EXCERPT_MAX_CHARS
            )
            if excerpt_truncated:
                truncated_excerpt_count += 1
            item["untrusted_excerpt"] = excerpt_probe[
                :CORPUS_SEARCH_EXCERPT_MAX_CHARS
            ]
            item["excerpt_truncated"] = excerpt_truncated
            item["excerpt_max_characters"] = CORPUS_SEARCH_EXCERPT_MAX_CHARS
            item["excerpt_generation"] = excerpt["generation"]
            item["surfaced_by"] = "lexical_fts"
            item["ranking_is_evidence"] = False
            candidates.append(item)
        response = {
            "query": query,
            "strategy": "lexical_candidate_acquisition",
            "query_mode": "exact_phrase_fts",
            "zero_results_establish_absence": False,
            "count": len(candidates),
            "candidates": candidates,
            "truncated_excerpt_count": truncated_excerpt_count,
            "notice": (
                "Candidate excerpts may be truncated and require agent interpretation. "
                "Use corpus_read for exact source content."
            ),
        }
        serialized_bytes = len(encode_json(response).encode())
        if serialized_bytes > CORPUS_SEARCH_MAX_SERIALIZED_BYTES:
            raise BudgetExceededError(
                "search response exceeds the serialized response budget",
                details={
                    "candidate_count": len(candidates),
                    "serialized_bytes": serialized_bytes,
                    "maximum_serialized_bytes": (
                        CORPUS_SEARCH_MAX_SERIALIZED_BYTES
                    ),
                },
            )
        return response

    def read_units(
        self,
        corpus_id: str,
        unit_ids: list[str],
        *,
        neighbor_span: int = 0,
        max_chars: int = CORPUS_READ_DEFAULT_CHARS,
    ) -> dict:
        if len(unit_ids) > _MAX_READ_UNITS:
            raise BudgetExceededError(
                "source unit selection must contain at most 200 ids",
                details={"unit_id_count": len(unit_ids), "maximum": _MAX_READ_UNITS},
            )
        if any(not unit_id or len(unit_id) > 200 for unit_id in unit_ids):
            raise BudgetExceededError(
                "source unit ids must contain between 1 and 200 characters",
                details={"maximum_unit_id_chars": 200},
            )
        if not 0 <= neighbor_span <= _MAX_NEIGHBOR_SPAN:
            raise BudgetExceededError(
                "neighbor span must be between 0 and 10",
                details={"neighbor_span": neighbor_span, "maximum": _MAX_NEIGHBOR_SPAN},
            )
        if not CORPUS_READ_MIN_CHARS <= max_chars <= CORPUS_READ_MAX_CHARS:
            raise BudgetExceededError(
                "source unit max_chars must be between 1000 and 200000",
                details={
                    "max_chars": max_chars,
                    "minimum": CORPUS_READ_MIN_CHARS,
                    "maximum": CORPUS_READ_MAX_CHARS,
                },
            )
        if not unit_ids:
            return {"units": [], "count": 0}
        requested_ids = list(dict.fromkeys(unit_ids))
        requested = set(requested_ids)
        with corpus_read_connection(self.data_root, corpus_id) as connection:
            connection.create_function(
                "corpus_character_length",
                1,
                lambda value: len(value or ""),
                deterministic=True,
            )
            placeholders = ",".join("?" for _ in requested_ids)
            seed_rows = connection.execute(
                f"""
                SELECT u.unit_id, u.projection_id, u.ordinal
                FROM source_units u
                JOIN revisions r ON r.revision_id = u.revision_id
                JOIN documents d ON d.document_id = r.document_id
                WHERE u.unit_id IN ({placeholders})
                  AND d.deleted_at IS NULL
                """,
                requested_ids,
            ).fetchall()
            seed_by_id = {row["unit_id"]: row for row in seed_rows}
            seeds = [
                seed_by_id[unit_id]
                for unit_id in requested_ids
                if unit_id in seed_by_id
            ]
            selected_ids: list[str] = []
            selected: set[str] = set()
            for seed in seeds:
                inventory = connection.execute(
                    """
                    SELECT unit_id
                    FROM source_units u
                    WHERE u.projection_id = ? AND u.ordinal BETWEEN ? AND ?
                    ORDER BY ordinal
                    """,
                    (
                        seed["projection_id"],
                        max(1, seed["ordinal"] - neighbor_span),
                        seed["ordinal"] + neighbor_span,
                    ),
                ).fetchall()
                for row in inventory:
                    if row["unit_id"] in selected:
                        continue
                    selected.add(row["unit_id"])
                    selected_ids.append(row["unit_id"])

            if len(selected_ids) > CORPUS_READ_MAX_SELECTED_UNITS:
                raise BudgetExceededError(
                    "source unit response exceeds the aggregate read budget",
                    details={
                        "selected_unit_count": len(selected_ids),
                        "maximum_selected_units": CORPUS_READ_MAX_SELECTED_UNITS,
                    },
                )
            if not selected_ids:
                return {"units": [], "count": 0}

            selected_placeholders = ",".join("?" for _ in selected_ids)
            inventory_rows = connection.execute(
                f"""
                SELECT unit_id,
                       corpus_character_length(normalized_content)
                           AS content_chars,
                       (
                           LENGTH(CAST(unit_id AS BLOB))
                           + LENGTH(CAST(revision_id AS BLOB))
                           + LENGTH(CAST(projection_id AS BLOB))
                           + LENGTH(CAST(unit_type AS BLOB))
                           + LENGTH(CAST(structure_path_json AS BLOB))
                           + LENGTH(CAST(source_anchor_json AS BLOB))
                           + LENGTH(CAST(normalized_content AS BLOB))
                           + LENGTH(CAST(content_sha256 AS BLOB))
                           + LENGTH(CAST(COALESCE(previous_unit_id, '') AS BLOB))
                           + LENGTH(CAST(COALESCE(next_unit_id, '') AS BLOB))
                           + LENGTH(CAST(extraction_issues_json AS BLOB))
                           + LENGTH(CAST(derivation_method AS BLOB))
                           + LENGTH(CAST(geometry_json AS BLOB))
                           + LENGTH(CAST(quality_flags_json AS BLOB))
                           + LENGTH(CAST(trust_lineage AS BLOB))
                           + 1024
                       ) AS payload_bytes
                FROM source_units
                WHERE unit_id IN ({selected_placeholders})
                """,
                selected_ids,
            ).fetchall()
            selected_content_chars = sum(
                row["content_chars"] for row in inventory_rows
            )
            selected_payload_bytes = sum(
                row["payload_bytes"] for row in inventory_rows
            )
            if (
                selected_content_chars > max_chars
                or selected_payload_bytes > CORPUS_READ_MAX_SERIALIZED_BYTES
            ):
                raise BudgetExceededError(
                    "source unit response exceeds the aggregate read budget",
                    details={
                        "selected_unit_count": len(selected_ids),
                        "maximum_selected_units": CORPUS_READ_MAX_SELECTED_UNITS,
                        "selected_content_chars": selected_content_chars,
                        "max_chars": max_chars,
                        "selected_payload_bytes": selected_payload_bytes,
                        "maximum_serialized_bytes": (
                            CORPUS_READ_MAX_SERIALIZED_BYTES
                        ),
                    },
                )

            body_rows = connection.execute(
                f"""
                SELECT u.*, d.document_id, d.relative_path, d.current_revision_id,
                       active.projection_id AS active_projection_id,
                       projection.completeness_state,
                       projection.assurance_state,
                       projection.adapter_id,
                       projection.adapter_version,
                       projection.config_hash AS projection_config_hash,
                       d.extension AS document_extension,
                       CASE WHEN
                           r.source_size = d.logical_size
                           AND r.source_modified_ns = d.modified_ns
                           AND r.source_changed_ns = d.changed_ns
                           AND r.source_device = d.device
                           AND r.source_inode = d.inode
                       THEN 1 ELSE 0 END AS source_observation_current
                FROM source_units u
                JOIN revisions r ON r.revision_id = u.revision_id
                JOIN documents d ON d.document_id = r.document_id
                JOIN extraction_projections projection
                  ON projection.projection_id = u.projection_id
                LEFT JOIN extraction_projections active
                  ON active.revision_id = u.revision_id AND active.is_active = 1
                WHERE u.unit_id IN ({selected_placeholders})
                  AND d.deleted_at IS NULL
                """,
                selected_ids,
            ).fetchall()
            row_by_id = {row["unit_id"]: row for row in body_rows}
            rows = [
                row_by_id[unit_id]
                for unit_id in selected_ids
                if unit_id in row_by_id
            ]
        units = []
        for row in rows:
            item = dict(row)
            item["structure_path"] = json.loads(item.pop("structure_path_json"))
            item["source_anchor"] = json.loads(item.pop("source_anchor_json"))
            item["extraction_issues"] = json.loads(item.pop("extraction_issues_json"))
            item["geometry"] = json.loads(item.pop("geometry_json"))
            item["quality_flags"] = json.loads(item.pop("quality_flags_json"))
            item["requested"] = item["unit_id"] in requested
            source_observation_current = bool(item.pop("source_observation_current"))
            projection_current = self._projection_uses_current_adapter(
                item.pop("document_extension"),
                item["adapter_id"],
                item["adapter_version"],
                item.pop("projection_config_hash"),
            )
            if item["revision_id"] != item["current_revision_id"]:
                item["dependency_state"] = "stale_source_revision"
            elif not source_observation_current:
                item["dependency_state"] = "stale_source_observation"
            elif item["projection_id"] != item["active_projection_id"]:
                item["dependency_state"] = "stale_extraction_projection"
            elif not projection_current:
                item["dependency_state"] = "stale_extraction_adapter"
            else:
                item["dependency_state"] = "valid"
            item["untrusted_content"] = item.pop("normalized_content")
            units.append(item)
        response = {"count": len(units), "units": units}
        serialized_bytes = len(encode_json(response).encode())
        if serialized_bytes > CORPUS_READ_MAX_SERIALIZED_BYTES:
            raise BudgetExceededError(
                "source unit response exceeds the serialized response budget",
                details={
                    "serialized_bytes": serialized_bytes,
                    "maximum_serialized_bytes": CORPUS_READ_MAX_SERIALIZED_BYTES,
                },
            )
        return response

    def interpretation_queue(
        self,
        corpus_id: str,
        *,
        limit: int = 50,
        include_outdated: bool = False,
    ) -> dict:
        if not 1 <= limit <= 200:
            raise BudgetExceededError(
                "interpretation queue limit must be between 1 and 200",
                details={"limit": limit, "maximum": 200},
            )
        with corpus_read_connection(self.data_root, corpus_id) as connection:
            snapshot = connection.execute(
                """
                SELECT snapshot_id FROM snapshots
                WHERE state = 'complete'
                ORDER BY rowid DESC LIMIT 1
                """
            ).fetchone()
            if snapshot is None:
                return {
                    "count": 0,
                    "items": [],
                    "base_snapshot_id": None,
                    "semantic_state": _semantic_state(connection),
                    "excluded_outdated": 0,
                    "semantic_materialization_state": (
                        "optional_explicit_maintenance"
                    ),
                }
            rows = connection.execute(
                """
                SELECT q.*, d.relative_path, r.sha256, r.extraction_state,
                       d.extension,
                       p.completeness_state, p.assurance_state,
                       p.adapter_id, p.adapter_version, p.config_hash,
                       COUNT(u.unit_id) AS source_units
                FROM interpretation_queue q
                JOIN documents d ON d.document_id = q.document_id
                JOIN revisions r ON r.revision_id = q.revision_id
                JOIN extraction_projections p ON p.projection_id = q.projection_id
                JOIN snapshot_documents snapshot_document
                  ON snapshot_document.snapshot_id = ?
                 AND snapshot_document.document_id = q.document_id
                 AND snapshot_document.revision_id = q.revision_id
                 AND snapshot_document.projection_id = q.projection_id
                LEFT JOIN source_units u ON u.projection_id = q.projection_id
                WHERE q.state IN ('pending', 'in_progress', 'failed')
                  AND d.current_revision_id = q.revision_id
                  AND d.deleted_at IS NULL
                  AND p.is_active = 1
                  AND r.source_size = d.logical_size
                  AND r.source_modified_ns = d.modified_ns
                  AND r.source_changed_ns = d.changed_ns
                  AND r.source_device = d.device
                  AND r.source_inode = d.inode
                GROUP BY q.queue_id
                ORDER BY q.updated_at, d.relative_path_nfc
                """,
                (snapshot["snapshot_id"],),
            ).fetchall()
            items = []
            excluded_outdated = 0
            for row in rows:
                item = dict(row)
                try:
                    descriptor = self.adapter_registry.resolve(item["extension"]).descriptor
                    projection_current = (
                        item["adapter_id"] == descriptor.adapter_id
                        and item["adapter_version"] == descriptor.adapter_version
                        and item["config_hash"] == descriptor.config_hash
                    )
                except ExtractionError:
                    projection_current = False
                item["projection_current"] = projection_current
                checkpoint_json = item.pop("checkpoint_json")
                item["checkpoint"] = json.loads(checkpoint_json) if checkpoint_json else None
                if not projection_current and not include_outdated:
                    excluded_outdated += 1
                    continue
                if len(items) < limit:
                    items.append(item)
            semantic_state = _semantic_state(connection)
        return {
            "count": len(items),
            "items": items,
            "base_snapshot_id": snapshot["snapshot_id"],
            "semantic_state": semantic_state,
            "excluded_outdated": excluded_outdated,
            "semantic_materialization_state": "optional_explicit_maintenance",
        }

    def reconcile_completed_checkpoint(
        self,
        corpus_id: str,
        *,
        queue_id: str,
        expected_snapshot_id: str,
        expected_updated_at: str,
    ) -> dict:
        """Repair one historical pending queue row after verifying its full checkpoint."""

        _require_semantic_commit_string(
            corpus_id,
            field="corpus_id",
            maximum=64,
            nonempty=True,
        )
        corpus_id = normalize_corpus_id(corpus_id)
        for field, value in (
            ("queue_id", queue_id),
            ("expected_snapshot_id", expected_snapshot_id),
            ("expected_updated_at", expected_updated_at),
        ):
            _require_semantic_commit_string(
                value,
                field=field,
                maximum=SEMANTIC_COMMIT_MAX_IDENTIFIER_CHARS,
                nonempty=True,
            )

        corpus = get_corpus(self.data_root, corpus_id)
        source_root = Path(corpus["source_root"])
        paths = self._paths(corpus_id)
        paths.ensure()
        with (
            writer_lock(paths.corpus_root / "writer.lock"),
            corpus_connection(self.data_root, corpus_id) as connection,
        ):
            latest_snapshot = connection.execute(
                """
                SELECT snapshot_id
                FROM snapshots
                WHERE state = 'complete'
                ORDER BY rowid DESC
                LIMIT 1
                """
            ).fetchone()
            current_snapshot_id = (
                latest_snapshot["snapshot_id"] if latest_snapshot is not None else None
            )
            if current_snapshot_id != expected_snapshot_id:
                raise SnapshotConflictError(
                    "checkpoint reconciliation snapshot is not current",
                    details={
                        "expected_snapshot_id": expected_snapshot_id,
                        "current_snapshot_id": current_snapshot_id,
                    },
                )

            row = connection.execute(
                """
                SELECT q.*,
                       d.relative_path, d.absolute_path, d.extension,
                       d.deleted_at, d.current_revision_id,
                       d.logical_size, d.modified_ns, d.changed_ns,
                       d.device, d.inode, d.residency_state, d.is_dataless,
                       r.sha256, r.source_size, r.source_modified_ns,
                       r.source_changed_ns, r.source_device, r.source_inode,
                       p.adapter_id, p.adapter_version, p.config_hash,
                       p.completeness_state, p.is_active
                FROM interpretation_queue q
                JOIN documents d ON d.document_id = q.document_id
                JOIN revisions r
                  ON r.revision_id = q.revision_id
                 AND r.document_id = q.document_id
                JOIN extraction_projections p
                  ON p.projection_id = q.projection_id
                 AND p.revision_id = q.revision_id
                WHERE q.queue_id = ?
                """,
                (queue_id,),
            ).fetchone()
            if row is None:
                raise SemanticCommitError(
                    "checkpoint reconciliation queue item does not exist",
                    details={"queue_id": queue_id},
                )
            if row["updated_at"] != expected_updated_at:
                raise SnapshotConflictError(
                    "checkpoint reconciliation queue changed since inspection",
                    details={"queue_id": queue_id},
                )
            if (
                row["state"] != "pending"
                or row["reason"] != "source_observation_reconfirmed"
            ):
                raise SemanticCommitError(
                    "checkpoint reconciliation only accepts the historical "
                    "pending reconfirmation state",
                    details={
                        "queue_id": queue_id,
                        "state": row["state"],
                        "reason": row["reason"],
                    },
                )

            source_observation_current = (
                row["source_size"],
                row["source_modified_ns"],
                row["source_changed_ns"],
                row["source_device"],
                row["source_inode"],
            ) == (
                row["logical_size"],
                row["modified_ns"],
                row["changed_ns"],
                row["device"],
                row["inode"],
            )
            if (
                row["deleted_at"] is not None
                or row["current_revision_id"] != row["revision_id"]
                or row["is_active"] != 1
                or not source_observation_current
            ):
                raise SnapshotConflictError(
                    "checkpoint reconciliation target is no longer current",
                    details={
                        "queue_id": queue_id,
                        "revision_id": row["revision_id"],
                        "projection_id": row["projection_id"],
                    },
                )

            try:
                descriptor = self.adapter_registry.resolve(row["extension"]).descriptor
            except ExtractionError as exc:
                raise SnapshotConflictError(
                    "checkpoint reconciliation projection has no current adapter",
                    details={"queue_id": queue_id},
                ) from exc
            if (
                row["adapter_id"],
                row["adapter_version"],
                row["config_hash"],
            ) != (
                descriptor.adapter_id,
                descriptor.adapter_version,
                descriptor.config_hash,
            ):
                raise SnapshotConflictError(
                    "checkpoint reconciliation projection uses an outdated adapter",
                    details={
                        "queue_id": queue_id,
                        "projection_id": row["projection_id"],
                    },
                )

            current_membership = connection.execute(
                """
                SELECT 1
                FROM snapshot_documents
                WHERE snapshot_id = ?
                  AND document_id = ?
                  AND revision_id = ?
                  AND projection_id = ?
                """,
                (
                    expected_snapshot_id,
                    row["document_id"],
                    row["revision_id"],
                    row["projection_id"],
                ),
            ).fetchone()
            if current_membership is None:
                raise SnapshotConflictError(
                    "checkpoint reconciliation target is outside the current snapshot",
                    details={
                        "queue_id": queue_id,
                        "expected_snapshot_id": expected_snapshot_id,
                    },
                )

            checkpoint_json = row["checkpoint_json"]
            if not self._checkpoint_completes_projection(
                connection,
                projection_id=row["projection_id"],
                checkpoint_json=checkpoint_json,
            ):
                raise SemanticCommitError(
                    "checkpoint reconciliation requires a full stored checkpoint",
                    details={
                        "queue_id": queue_id,
                        "projection_id": row["projection_id"],
                    },
                )
            checkpoint = json.loads(checkpoint_json)
            updated_by_commit_id = checkpoint.get("updated_by_commit_id")
            if (
                not isinstance(updated_by_commit_id, str)
                or not updated_by_commit_id.strip()
                or len(updated_by_commit_id) > SEMANTIC_COMMIT_MAX_IDENTIFIER_CHARS
            ):
                raise SemanticCommitError(
                    "checkpoint reconciliation requires an updated_by semantic commit",
                    details={"queue_id": queue_id},
                )

            semantic_commit = connection.execute(
                """
                SELECT base_snapshot_id, result_json
                FROM semantic_commits
                WHERE commit_id = ?
                """,
                (updated_by_commit_id,),
            ).fetchone()
            if semantic_commit is None:
                raise SemanticCommitError(
                    "checkpoint reconciliation semantic commit does not exist",
                    details={
                        "queue_id": queue_id,
                        "commit_id": updated_by_commit_id,
                    },
                )
            base_membership = connection.execute(
                """
                SELECT 1
                FROM snapshots s
                JOIN snapshot_documents sd ON sd.snapshot_id = s.snapshot_id
                WHERE s.snapshot_id = ?
                  AND s.state = 'complete'
                  AND sd.document_id = ?
                  AND sd.revision_id = ?
                  AND sd.projection_id = ?
                """,
                (
                    semantic_commit["base_snapshot_id"],
                    row["document_id"],
                    row["revision_id"],
                    row["projection_id"],
                ),
            ).fetchone()
            if base_membership is None:
                raise SnapshotConflictError(
                    "checkpoint semantic commit does not depend on the target projection",
                    details={
                        "queue_id": queue_id,
                        "commit_id": updated_by_commit_id,
                        "base_snapshot_id": semantic_commit["base_snapshot_id"],
                    },
                )
            try:
                semantic_result = json.loads(semantic_commit["result_json"])
            except (TypeError, json.JSONDecodeError) as exc:
                raise SemanticCommitError(
                    "checkpoint semantic commit result is invalid",
                    details={
                        "queue_id": queue_id,
                        "commit_id": updated_by_commit_id,
                    },
                ) from exc
            queue_updates = (
                semantic_result.get("queue_updates")
                if isinstance(semantic_result, dict)
                else None
            )
            matching_updates = (
                [
                    update
                    for update in queue_updates
                    if isinstance(update, dict) and update.get("queue_id") == queue_id
                ]
                if isinstance(queue_updates, list)
                else []
            )
            if (
                not isinstance(semantic_result, dict)
                or semantic_result.get("commit_id") != updated_by_commit_id
                or semantic_result.get("base_snapshot_id")
                != semantic_commit["base_snapshot_id"]
                or len(matching_updates) != 1
                or matching_updates[0].get("revision_id") != row["revision_id"]
                or matching_updates[0].get("state") != "complete"
                or matching_updates[0].get("checkpoint") != checkpoint
            ):
                raise SemanticCommitError(
                    "checkpoint is not the completed result of its semantic commit",
                    details={
                        "queue_id": queue_id,
                        "commit_id": updated_by_commit_id,
                    },
                )

            source_path = Path(row["absolute_path"])
            expected_source_path = source_root / row["relative_path"]
            if (
                source_path != expected_source_path
                or row["residency_state"] != "resident"
                or row["is_dataless"] != 0
                or not _resident_source_matches_revision(
                    source_path=source_path,
                    source_root=source_root,
                    staging_root=paths.staging,
                    document=dict(row),
                    revision_sha256=row["sha256"],
                )
            ):
                raise SnapshotConflictError(
                    "checkpoint reconciliation source no longer exactly matches the revision",
                    details={
                        "queue_id": queue_id,
                        "revision_id": row["revision_id"],
                    },
                )

            now = utc_now()
            updated = connection.execute(
                """
                UPDATE interpretation_queue
                SET state = 'complete', updated_at = ?
                WHERE queue_id = ?
                  AND state = 'pending'
                  AND reason = 'source_observation_reconfirmed'
                  AND updated_at = ?
                  AND checkpoint_json = ?
                  AND EXISTS (
                      SELECT 1
                      FROM snapshots latest
                      WHERE latest.snapshot_id = ?
                        AND latest.state = 'complete'
                        AND latest.rowid = (
                            SELECT MAX(rowid)
                            FROM snapshots
                            WHERE state = 'complete'
                        )
                  )
                """,
                (
                    now,
                    queue_id,
                    expected_updated_at,
                    checkpoint_json,
                    expected_snapshot_id,
                ),
            )
            if updated.rowcount != 1:
                raise SnapshotConflictError(
                    "checkpoint reconciliation compare-and-set failed",
                    details={"queue_id": queue_id},
                )

            event_id = f"event_{uuid.uuid4().hex}"
            connection.execute(
                """
                INSERT INTO events(
                    event_id, event_type, document_id, revision_id,
                    payload_json, created_at
                ) VALUES (?, 'semantic_checkpoint_reconciled', ?, ?, ?, ?)
                """,
                (
                    event_id,
                    row["document_id"],
                    row["revision_id"],
                    encode_json(
                        {
                            "schema_version": 1,
                            "queue_id": queue_id,
                            "projection_id": row["projection_id"],
                            "expected_snapshot_id": expected_snapshot_id,
                            "updated_by_commit_id": updated_by_commit_id,
                            "previous_state": "pending",
                            "previous_reason": row["reason"],
                            "previous_updated_at": expected_updated_at,
                            "new_state": "complete",
                            "checkpoint_sha256": hashlib.sha256(
                                checkpoint_json.encode()
                            ).hexdigest(),
                            "checkpoint_preserved": True,
                            "resident_source_hash_verified": True,
                        }
                    ),
                    now,
                ),
            )
            semantic_state = _semantic_state(connection)

        return {
            "corpus_id": corpus_id,
            "queue_id": queue_id,
            "document_id": row["document_id"],
            "revision_id": row["revision_id"],
            "projection_id": row["projection_id"],
            "state": "complete",
            "reason": row["reason"],
            "snapshot_id": expected_snapshot_id,
            "updated_by_commit_id": updated_by_commit_id,
            "updated_at": now,
            "event_id": event_id,
            "checkpoint_preserved": True,
            "resident_source_hash_verified": True,
            "semantic_state": semantic_state,
        }

    def interpretation_material(
        self,
        corpus_id: str,
        *,
        queue_id: str,
        start_ordinal: int | None = None,
        max_units: int = 40,
        max_chars: int = 30_000,
    ) -> dict:
        if not queue_id.strip():
            raise SemanticCommitError("queue_id is required")
        if not 1 <= max_units <= 100:
            raise BudgetExceededError(
                "materialization max_units must be between 1 and 100",
                details={"max_units": max_units},
            )
        if not 1_000 <= max_chars <= 200_000:
            raise BudgetExceededError(
                "materialization max_chars must be between 1000 and 200000",
                details={"max_chars": max_chars},
            )
        with corpus_read_connection(self.data_root, corpus_id) as connection:
            snapshot = connection.execute(
                """
                SELECT snapshot_id, coverage_state, document_count, supported_document_count
                FROM snapshots
                WHERE state = 'complete'
                ORDER BY rowid DESC LIMIT 1
                """
            ).fetchone()
            if snapshot is None:
                raise SemanticCommitError("semantic materialization requires a complete snapshot")
            row = connection.execute(
                """
                SELECT q.*, d.relative_path, d.extension, d.current_revision_id,
                       r.sha256, r.source_size, r.source_modified_ns,
                       r.source_changed_ns,
                       p.adapter_id, p.adapter_version, p.config_hash,
                       p.completeness_state, p.assurance_state,
                       LENGTH(CAST(p.capability_manifest_json AS BLOB))
                         AS capability_manifest_bytes
                FROM interpretation_queue q
                JOIN documents d ON d.document_id = q.document_id
                JOIN revisions r ON r.revision_id = q.revision_id
                JOIN extraction_projections p ON p.projection_id = q.projection_id
                JOIN snapshot_documents snapshot_document
                  ON snapshot_document.snapshot_id = ?
                 AND snapshot_document.document_id = q.document_id
                 AND snapshot_document.revision_id = q.revision_id
                 AND snapshot_document.projection_id = q.projection_id
                WHERE q.queue_id = ?
                  AND q.state IN ('pending', 'in_progress', 'failed')
                  AND d.deleted_at IS NULL
                  AND d.current_revision_id = q.revision_id
                  AND p.is_active = 1
                  AND r.source_size = d.logical_size
                  AND r.source_modified_ns = d.modified_ns
                  AND r.source_changed_ns = d.changed_ns
                  AND r.source_device = d.device
                  AND r.source_inode = d.inode
                """,
                (snapshot["snapshot_id"], queue_id),
            ).fetchone()
            if row is None:
                raise SemanticCommitError(
                    "queue item is not current in the latest snapshot",
                    details={"queue_id": queue_id, "base_snapshot_id": snapshot["snapshot_id"]},
                )
            item = dict(row)
            descriptor = self.adapter_registry.resolve(item["extension"]).descriptor
            if (
                item["adapter_id"],
                item["adapter_version"],
                item["config_hash"],
            ) != (
                descriptor.adapter_id,
                descriptor.adapter_version,
                descriptor.config_hash,
            ):
                raise SemanticCommitError(
                    "queue item uses an outdated extraction projection",
                    details={
                        "queue_id": queue_id,
                        "projection_id": item["projection_id"],
                        "current_adapter_id": descriptor.adapter_id,
                        "current_adapter_version": descriptor.adapter_version,
                    },
                )

            bounds = connection.execute(
                """
                SELECT MIN(ordinal) AS min_ordinal, MAX(ordinal) AS max_ordinal,
                       COUNT(*) AS unit_count
                FROM source_units WHERE projection_id = ?
                """,
                (item["projection_id"],),
            ).fetchone()
            if not bounds["unit_count"]:
                raise SemanticCommitError(
                    "queue projection has no source units",
                    details={"projection_id": item["projection_id"]},
                )
            checkpoint = json.loads(item["checkpoint_json"]) if item["checkpoint_json"] else None
            if start_ordinal is None:
                start_ordinal = (
                    int(checkpoint["next_ordinal"])
                    if checkpoint and checkpoint.get("next_ordinal") is not None
                    else int(bounds["min_ordinal"])
                )
            if not bounds["min_ordinal"] <= start_ordinal <= bounds["max_ordinal"]:
                raise SemanticCommitError(
                    "start_ordinal is outside the queued projection",
                    details={
                        "start_ordinal": start_ordinal,
                        "min_ordinal": bounds["min_ordinal"],
                        "max_ordinal": bounds["max_ordinal"],
                    },
                )
            connection.create_function(
                "corpus_character_length",
                1,
                lambda value: len(value or ""),
            )
            unit_inventory_rows = connection.execute(
                """
                SELECT unit_id, ordinal,
                       corpus_character_length(normalized_content)
                         AS content_chars,
                       LENGTH(CAST(COALESCE(unit_id, '') AS BLOB))
                         + LENGTH(CAST(COALESCE(revision_id, '') AS BLOB))
                         + LENGTH(CAST(COALESCE(projection_id, '') AS BLOB))
                         + LENGTH(CAST(COALESCE(unit_type, '') AS BLOB))
                         + LENGTH(CAST(COALESCE(structure_path_json, '') AS BLOB))
                         + LENGTH(CAST(COALESCE(source_anchor_json, '') AS BLOB))
                         + LENGTH(CAST(COALESCE(normalized_content, '') AS BLOB))
                         + LENGTH(CAST(COALESCE(content_sha256, '') AS BLOB))
                         + LENGTH(CAST(COALESCE(previous_unit_id, '') AS BLOB))
                         + LENGTH(CAST(COALESCE(next_unit_id, '') AS BLOB))
                         + LENGTH(CAST(COALESCE(extraction_issues_json, '') AS BLOB))
                         + LENGTH(CAST(COALESCE(derivation_method, '') AS BLOB))
                         + LENGTH(CAST(COALESCE(geometry_json, '') AS BLOB))
                         + LENGTH(CAST(COALESCE(quality_flags_json, '') AS BLOB))
                         + LENGTH(CAST(COALESCE(trust_lineage, '') AS BLOB))
                         + 1024 AS stored_payload_bytes
                FROM source_units
                WHERE projection_id = ? AND ordinal >= ?
                ORDER BY ordinal
                LIMIT ?
                """,
                (item["projection_id"], start_ordinal, max_units),
            ).fetchall()
            selected_inventory = []
            character_count = 0
            for inventory_row in unit_inventory_rows:
                content_chars = inventory_row["content_chars"]
                if character_count + content_chars > max_chars:
                    if not selected_inventory:
                        raise BudgetExceededError(
                            "one source unit exceeds the materialization character budget",
                            details={
                                "unit_id": inventory_row["unit_id"],
                                "content_chars": content_chars,
                                "max_chars": max_chars,
                            },
                        )
                    break
                selected_inventory.append(inventory_row)
                character_count += content_chars
            if not selected_inventory:
                raise BudgetExceededError("materialization budget returned no source units")

            issue_inventory = connection.execute(
                """
                SELECT COUNT(*) AS issue_count,
                       COALESCE(
                           SUM(
                               LENGTH(CAST(COALESCE(severity, '') AS BLOB))
                               + LENGTH(CAST(COALESCE(code, '') AS BLOB))
                               + LENGTH(CAST(COALESCE(message, '') AS BLOB))
                               + LENGTH(CAST(COALESCE(details_json, '') AS BLOB))
                               + LENGTH(
                                   CAST(COALESCE(structural_locator_json, '') AS BLOB)
                                 )
                               + 512
                           ),
                           0
                       ) AS stored_payload_bytes
                FROM extraction_issues
                WHERE projection_id = ? AND lifecycle_state = 'active'
                """,
                (item["projection_id"],),
            ).fetchone()
            existing_claim_inventory = connection.execute(
                """
                SELECT COUNT(*) AS claim_count,
                       COALESCE(
                           SUM(
                               LENGTH(CAST(COALESCE(c.claim_id, '') AS BLOB))
                               + LENGTH(CAST(COALESCE(c.body, '') AS BLOB))
                               + LENGTH(CAST(COALESCE(c.subject, '') AS BLOB))
                               + LENGTH(CAST(COALESCE(c.modality, '') AS BLOB))
                               + LENGTH(
                                   CAST(COALESCE(c.claim_assessment, '') AS BLOB)
                                 )
                               + LENGTH(CAST(COALESCE(c.validation_state, '') AS BLOB))
                               + LENGTH(
                                   CAST(COALESCE(c.temporal_applicability, '') AS BLOB)
                                 )
                               + LENGTH(CAST(COALESCE(c.contest_state, '') AS BLOB))
                               + LENGTH(CAST(COALESCE(c.apparent_status, '') AS BLOB))
                               + LENGTH(
                                   CAST(COALESCE(c.materializer_version, '') AS BLOB)
                                 )
                               + 512
                           ),
                           0
                       ) AS stored_payload_bytes
                FROM (
                    SELECT c.claim_id, c.body, c.subject, c.modality,
                           c.claim_assessment, c.validation_state,
                           c.temporal_applicability, c.contest_state,
                           c.apparent_status, c.materializer_version,
                           c.created_at
                    FROM atomic_claims c
                    WHERE c.dependency_state = 'valid'
                      AND EXISTS (
                          SELECT 1
                          FROM evidence_links e
                          WHERE e.claim_id = c.claim_id
                            AND e.source_revision_id = ?
                      )
                    ORDER BY c.created_at DESC, c.claim_id
                    LIMIT 50
                ) c
                """,
                (item["revision_id"],),
            ).fetchone()
            selected_unit_payload_bytes = sum(
                row["stored_payload_bytes"] for row in selected_inventory
            )
            stored_payload_bytes = (
                selected_unit_payload_bytes
                + issue_inventory["stored_payload_bytes"]
                + existing_claim_inventory["stored_payload_bytes"]
                + item["capability_manifest_bytes"]
                + 8_192
            )
            if stored_payload_bytes > INTERPRETATION_MATERIAL_MAX_SERIALIZED_BYTES:
                raise BudgetExceededError(
                    "interpretation material exceeds the serialized response budget",
                    details={
                        "requested_max_units": max_units,
                        "requested_max_chars": max_chars,
                        "selected_unit_count": len(selected_inventory),
                        "selected_unit_payload_bytes": selected_unit_payload_bytes,
                        "extraction_issue_count": issue_inventory["issue_count"],
                        "extraction_issue_payload_bytes": issue_inventory[
                            "stored_payload_bytes"
                        ],
                        "stored_payload_bytes": stored_payload_bytes,
                        "maximum_serialized_bytes": (
                            INTERPRETATION_MATERIAL_MAX_SERIALIZED_BYTES
                        ),
                        "retry_with_lower": ["max_units", "max_chars"],
                    },
                )

            unit_rows = connection.execute(
                """
                SELECT * FROM source_units
                WHERE projection_id = ? AND ordinal >= ?
                ORDER BY ordinal
                LIMIT ?
                """,
                (
                    item["projection_id"],
                    start_ordinal,
                    len(selected_inventory),
                ),
            ).fetchall()
            units = []
            for unit_row in unit_rows:
                unit = dict(unit_row)
                unit["structure_path"] = json.loads(unit.pop("structure_path_json"))
                unit["source_anchor"] = json.loads(unit.pop("source_anchor_json"))
                unit["extraction_issues"] = json.loads(unit.pop("extraction_issues_json"))
                unit["geometry"] = json.loads(unit.pop("geometry_json"))
                unit["quality_flags"] = json.loads(unit.pop("quality_flags_json"))
                unit["untrusted_content"] = unit.pop("normalized_content")
                unit["dependency_state"] = "valid"
                units.append(unit)

            end_ordinal = units[-1]["ordinal"]
            next_ordinal = end_ordinal + 1
            has_more = end_ordinal < bounds["max_ordinal"]
            issues = [
                {
                    **dict(issue),
                    "details": json.loads(issue["details_json"]),
                    "structural_locator": json.loads(issue["structural_locator_json"]),
                }
                for issue in connection.execute(
                    """
                    SELECT severity, code, message, details_json, structural_locator_json
                    FROM extraction_issues
                    WHERE projection_id = ? AND lifecycle_state = 'active'
                    ORDER BY severity, code, locator_key
                    """,
                    (item["projection_id"],),
                )
            ]
            for issue in issues:
                issue.pop("details_json")
                issue.pop("structural_locator_json")
            capability_manifest_json = connection.execute(
                """
                SELECT capability_manifest_json
                FROM extraction_projections
                WHERE projection_id = ?
                """,
                (item["projection_id"],),
            ).fetchone()["capability_manifest_json"]
            claim_rows = connection.execute(
                """
                SELECT DISTINCT c.claim_id, c.body, c.subject, c.modality,
                       c.claim_assessment, c.validation_state,
                       c.temporal_applicability, c.contest_state,
                       c.apparent_status, c.materializer_version
                FROM atomic_claims c
                JOIN evidence_links e ON e.claim_id = c.claim_id
                WHERE e.source_revision_id = ? AND c.dependency_state = 'valid'
                ORDER BY c.created_at DESC, c.claim_id
                LIMIT 50
                """,
                (item["revision_id"],),
            ).fetchall()
            existing_claims = [dict(claim) for claim in claim_rows]
            existing_claim_count = connection.execute(
                """
                SELECT COUNT(DISTINCT c.claim_id)
                FROM atomic_claims c
                JOIN evidence_links e ON e.claim_id = c.claim_id
                WHERE e.source_revision_id = ? AND c.dependency_state = 'valid'
                """,
                (item["revision_id"],),
            ).fetchone()[0]
            semantic_state = _semantic_state(connection)
            checkpoint_payload = {
                "revision_id": item["revision_id"],
                "processed_from_ordinal": units[0]["ordinal"],
                "next_ordinal": next_ordinal,
                "processed_through_ordinal": end_ordinal,
            }
            receipt_key = _materialization_receipt_key(
                self.data_root,
                create=True,
            )
            checkpoint_payload["batch_receipt"] = _materialization_batch_receipt(
                receipt_key=receipt_key,
                base_snapshot_id=snapshot["snapshot_id"],
                base_semantic_state_hash=semantic_state["semantic_state_hash"],
                queue_id=queue_id,
                revision_id=item["revision_id"],
                projection_id=item["projection_id"],
                processed_from_ordinal=units[0]["ordinal"],
                processed_through_ordinal=end_ordinal,
                units=units,
            )
            response = {
                "base_snapshot_id": snapshot["snapshot_id"],
                "semantic_state": semantic_state,
                "queue": {
                    "queue_id": queue_id,
                    "state": item["state"],
                    "reason": item["reason"],
                    "checkpoint": checkpoint,
                },
                "document": {
                    "document_id": item["document_id"],
                    "relative_path": item["relative_path"],
                    "revision_id": item["revision_id"],
                    "revision_sha256": item["sha256"],
                },
                "extraction": {
                    "projection_id": item["projection_id"],
                    "adapter_id": item["adapter_id"],
                    "adapter_version": item["adapter_version"],
                    "config_hash": item["config_hash"],
                    "completeness_state": item["completeness_state"],
                    "assurance_state": item["assurance_state"],
                    "capabilities": json.loads(capability_manifest_json),
                    "issues": issues,
                },
                "batch": {
                    "start_ordinal": units[0]["ordinal"],
                    "end_ordinal": end_ordinal,
                    "unit_count": len(units),
                    "character_count": character_count,
                    "has_more": has_more,
                    "next_ordinal": next_ordinal if has_more else None,
                    "total_units": bounds["unit_count"],
                    "units": units,
                },
                "existing_claims": existing_claims,
                "existing_claim_count": existing_claim_count,
                "existing_claims_truncated": existing_claim_count > len(existing_claims),
                "progress_update": checkpoint_payload,
                "materialization_contract": {
                    "source_content": "untrusted",
                    "created_claims": "model_authored_unreviewed",
                    "completion": "agent_declared_after_coherent_pass",
                    "authority": "none",
                },
            }
            serialized_bytes = len(encode_json(response).encode())
            if serialized_bytes > INTERPRETATION_MATERIAL_MAX_SERIALIZED_BYTES:
                raise BudgetExceededError(
                    "interpretation material exceeds the serialized response budget",
                    details={
                        "requested_max_units": max_units,
                        "requested_max_chars": max_chars,
                        "returned_unit_count": len(units),
                        "extraction_issue_count": len(issues),
                        "serialized_bytes": serialized_bytes,
                        "maximum_serialized_bytes": (
                            INTERPRETATION_MATERIAL_MAX_SERIALIZED_BYTES
                        ),
                        "retry_with_lower": ["max_units", "max_chars"],
                    },
                )
            return response

    def interpretation_commit(
        self,
        corpus_id: str,
        *,
        base_snapshot_id: str,
        base_semantic_state_hash: str,
        idempotency_key: str,
        claims: list[dict],
        completed_revision_ids: list[str] | None = None,
        progress_updates: list[dict] | None = None,
        materializer_version: str = "corpus-agent-v1",
    ) -> dict:
        _require_semantic_commit_string(
            corpus_id,
            field="corpus_id",
            maximum=64,
            nonempty=True,
        )
        corpus_id = normalize_corpus_id(corpus_id)
        (
            claims,
            completed_revision_ids,
            progress_updates,
            canonical_input_bytes,
        ) = _bounded_semantic_commit_input(
            base_snapshot_id=base_snapshot_id,
            base_semantic_state_hash=base_semantic_state_hash,
            idempotency_key=idempotency_key,
            claims=claims,
            completed_revision_ids=completed_revision_ids,
            progress_updates=progress_updates,
            materializer_version=materializer_version,
        )
        input_sha256 = hashlib.sha256(canonical_input_bytes).hexdigest()
        commit_digest = hashlib.sha256(
            f"{corpus_id}\0{base_snapshot_id}\0{idempotency_key}".encode()
        ).hexdigest()
        commit_id = f"semcommit_{commit_digest[:32]}"
        paths = self._paths(corpus_id)
        paths.ensure()

        with (
            writer_lock(paths.corpus_root / "writer.lock"),
            corpus_connection(self.data_root, corpus_id) as connection,
        ):
            existing = connection.execute(
                """
                    SELECT input_sha256, result_json FROM semantic_commits
                    WHERE idempotency_key = ?
                    """,
                (idempotency_key,),
            ).fetchone()
            if existing:
                if existing["input_sha256"] != input_sha256:
                    raise SemanticCommitError(
                        "idempotency key was already used with different content",
                        details={"idempotency_key": idempotency_key},
                    )
                return json.loads(existing["result_json"])

            current_snapshot = connection.execute(
                """
                    SELECT snapshot_id FROM snapshots
                    WHERE state = 'complete'
                    ORDER BY rowid DESC LIMIT 1
                    """
            ).fetchone()
            current_snapshot_id = current_snapshot["snapshot_id"] if current_snapshot else None
            if current_snapshot_id != base_snapshot_id:
                raise SnapshotConflictError(
                    "semantic commit base snapshot is not current",
                    details={
                        "base_snapshot_id": base_snapshot_id,
                        "current_snapshot_id": current_snapshot_id,
                    },
                )
            current_semantic_state = _semantic_state(connection)
            if current_semantic_state["semantic_state_hash"] != base_semantic_state_hash:
                raise SnapshotConflictError(
                    "semantic commit base semantic state is not current",
                    details={
                        "base_semantic_state_hash": base_semantic_state_hash,
                        "current_semantic_state_hash": current_semantic_state[
                            "semantic_state_hash"
                        ],
                    },
                )
            receipt_key = (
                _materialization_receipt_key(self.data_root, create=False)
                if progress_updates
                else None
            )

            snapshot_projection_by_revision = {
                row["revision_id"]: row["projection_id"]
                for row in connection.execute(
                    """
                        SELECT revision_id, projection_id FROM snapshot_documents
                        WHERE snapshot_id = ?
                        """,
                    (base_snapshot_id,),
                )
            }
            snapshot_revisions = set(snapshot_projection_by_revision)
            invalid_completed = sorted(set(completed_revision_ids) - snapshot_revisions)
            if invalid_completed:
                raise SemanticCommitError(
                    "completed revisions are not part of the base snapshot",
                    details={"revision_ids": invalid_completed},
                )
            progress_revision_ids = [
                str(update.get("revision_id", "")).strip() for update in progress_updates
            ]
            if len(progress_revision_ids) != len(set(progress_revision_ids)) or any(
                not revision_id for revision_id in progress_revision_ids
            ):
                raise SemanticCommitError(
                    "progress revision_id values must be non-empty and unique"
                )
            invalid_progress = sorted(set(progress_revision_ids) - snapshot_revisions)
            if invalid_progress:
                raise SemanticCommitError(
                    "progress revisions are not part of the base snapshot",
                    details={"revision_ids": invalid_progress},
                )

            prepared_claims: list[dict] = []
            prepared_links: list[dict] = []
            client_refs: set[str] = set()
            forbidden_fields = {
                "created_by",
                "validation_state",
                "trust_lineage",
                "review_record",
                "review_record_ids",
                "authority_record",
                "authority_record_ids",
            }
            assessment_values = {
                "supported",
                "qualified",
                "conflicting",
                "unresolved",
            }
            temporal_values = {"current", "expired", "future", "unknown"}
            contest_values = {"uncontested", "disputed", "unknown"}
            stance_values = {"supports", "qualifies", "contradicts", "mentions"}

            for claim_index, claim in enumerate(claims, start=1):
                forbidden = sorted(forbidden_fields.intersection(claim))
                if forbidden:
                    raise SemanticCommitError(
                        "model commits cannot write review, authority, or lineage fields",
                        details={"claim_index": claim_index, "fields": forbidden},
                    )
                client_ref = str(claim.get("client_ref", f"claim-{claim_index}")).strip()
                body = str(claim.get("body", "")).strip()
                if not client_ref or client_ref in client_refs:
                    raise SemanticCommitError(
                        "claim client_ref must be non-empty and unique",
                        details={"claim_index": claim_index, "client_ref": client_ref},
                    )
                client_refs.add(client_ref)
                if not body or len(body) > SEMANTIC_COMMIT_MAX_BODY_CHARS:
                    raise SemanticCommitError(
                        "claim body must contain "
                        f"1-{SEMANTIC_COMMIT_MAX_BODY_CHARS} characters",
                        details={"claim_index": claim_index},
                    )
                assessment = claim.get("claim_assessment", "unresolved")
                temporal = claim.get("temporal_applicability", "unknown")
                contest = claim.get("contest_state", "unknown")
                if assessment not in assessment_values:
                    raise SemanticCommitError(
                        "invalid claim assessment",
                        details={"claim_index": claim_index, "value": assessment},
                    )
                if temporal not in temporal_values or contest not in contest_values:
                    raise SemanticCommitError(
                        "invalid temporal or contest state",
                        details={"claim_index": claim_index},
                    )
                evidence = claim.get("evidence", [])
                if (
                    not isinstance(evidence, list)
                    or not 1
                    <= len(evidence)
                    <= SEMANTIC_COMMIT_MAX_EVIDENCE_PER_CLAIM
                ):
                    raise SemanticCommitError(
                        "each claim requires "
                        f"1-{SEMANTIC_COMMIT_MAX_EVIDENCE_PER_CLAIM} evidence links",
                        details={"claim_index": claim_index},
                    )
                claim_digest = hashlib.sha256(
                    f"{commit_id}\0{client_ref}\0{body}".encode()
                ).hexdigest()
                claim_id = f"claim_{claim_digest[:32]}"
                prepared_claims.append(
                    {
                        "claim_id": claim_id,
                        "body": body,
                        "subject": claim.get("subject"),
                        "modality": str(claim.get("modality", "asserted")),
                        "scope_and_conditions": claim.get("scope_and_conditions", {}),
                        "time_window": claim.get("time_window", {}),
                        "claim_assessment": assessment,
                        "temporal_applicability": temporal,
                        "contest_state": contest,
                        "apparent_status": claim.get("apparent_status"),
                    }
                )
                for evidence_index, link in enumerate(evidence, start=1):
                    source_unit_id = str(link.get("source_unit_id", "")).strip()
                    stance = link.get("stance", "supports")
                    if not source_unit_id or stance not in stance_values:
                        raise SemanticCommitError(
                            "evidence requires a source_unit_id and valid stance",
                            details={
                                "claim_index": claim_index,
                                "evidence_index": evidence_index,
                            },
                        )
                    prepared_links.append(
                        {
                            "claim_id": claim_id,
                            "source_unit_id": source_unit_id,
                            "stance": stance,
                            "qualifier": link.get("qualifier"),
                            "applicability": link.get("applicability", {}),
                            "claim_index": claim_index,
                            "evidence_index": evidence_index,
                        }
                    )

            requested_unit_ids = list(
                dict.fromkeys(link["source_unit_id"] for link in prepared_links)
            )
            source_units = {}
            if requested_unit_ids:
                placeholders = ",".join("?" for _ in requested_unit_ids)
                source_units = {
                    row["unit_id"]: dict(row)
                    for row in connection.execute(
                        f"""
                            SELECT u.unit_id, u.revision_id, u.projection_id,
                                   u.source_anchor_json, d.current_revision_id,
                                   active.projection_id AS active_projection_id
                            FROM source_units u
                            JOIN revisions r ON r.revision_id = u.revision_id
                            JOIN documents d ON d.document_id = r.document_id
                            LEFT JOIN extraction_projections active
                              ON active.revision_id = u.revision_id AND active.is_active = 1
                            WHERE u.unit_id IN ({placeholders})
                            """,
                        requested_unit_ids,
                    )
                }
            missing_units = sorted(set(requested_unit_ids) - set(source_units))
            if missing_units:
                raise SemanticCommitError(
                    "evidence references unknown source units",
                    details={"source_unit_ids": missing_units},
                )

            touched_revision_ids = (
                set(progress_revision_ids)
                | set(completed_revision_ids)
                | {unit["revision_id"] for unit in source_units.values()}
            )
            for revision_id in sorted(touched_revision_ids):
                if revision_id not in snapshot_projection_by_revision:
                    raise SnapshotConflictError(
                        "semantic input revision is outside the base snapshot",
                        details={"revision_id": revision_id},
                    )
                projection_id = snapshot_projection_by_revision[revision_id]
                projection_row = connection.execute(
                    """
                    SELECT d.extension, p.adapter_id, p.adapter_version, p.config_hash
                    FROM revisions r
                    JOIN documents d ON d.document_id = r.document_id
                    JOIN extraction_projections p ON p.projection_id = ?
                    WHERE r.revision_id = ?
                    """,
                    (projection_id, revision_id),
                ).fetchone()
                if projection_row is None:
                    raise SnapshotConflictError(
                        "semantic input projection is no longer available",
                        details={"revision_id": revision_id},
                    )
                descriptor = self.adapter_registry.resolve(
                    projection_row["extension"]
                ).descriptor
                if (
                    projection_row["adapter_id"],
                    projection_row["adapter_version"],
                    projection_row["config_hash"],
                ) != (
                    descriptor.adapter_id,
                    descriptor.adapter_version,
                    descriptor.config_hash,
                ):
                    raise SnapshotConflictError(
                        "semantic input uses an outdated extraction adapter",
                        details={"revision_id": revision_id},
                    )

            evidence_link_ids: list[str] = []
            for link in prepared_links:
                unit = source_units[link["source_unit_id"]]
                if (
                    unit["revision_id"] not in snapshot_revisions
                    or unit["revision_id"] != unit["current_revision_id"]
                    or unit["projection_id"]
                    != snapshot_projection_by_revision[unit["revision_id"]]
                    or unit["projection_id"] != unit["active_projection_id"]
                ):
                    raise SnapshotConflictError(
                        "evidence unit is stale or outside the base snapshot",
                        details={"source_unit_id": link["source_unit_id"]},
                    )
                anchor = json.loads(unit["source_anchor_json"])
                link_digest = hashlib.sha256(
                    (
                        f"{link['claim_id']}\0{link['source_unit_id']}\0"
                        f"{link['stance']}\0{link['evidence_index']}"
                    ).encode()
                ).hexdigest()
                link["evidence_link_id"] = f"evidence_{link_digest[:32]}"
                link["source_revision_id"] = unit["revision_id"]
                link["source_span"] = anchor.get("source_span")
                evidence_link_ids.append(link["evidence_link_id"])

            prepared_progress: list[dict] = []
            for progress_index, update in enumerate(progress_updates, start=1):
                revision_id = str(update.get("revision_id", "")).strip()
                projection_id = snapshot_projection_by_revision[revision_id]
                queue_row = connection.execute(
                    """
                    SELECT q.queue_id, q.checkpoint_json,
                           d.extension, p.adapter_id, p.adapter_version, p.config_hash,
                           MIN(u.ordinal) AS min_ordinal,
                           MAX(u.ordinal) AS max_ordinal,
                           COUNT(u.unit_id) AS total_units
                    FROM interpretation_queue q
                    JOIN documents d ON d.document_id = q.document_id
                    JOIN extraction_projections p ON p.projection_id = q.projection_id
                    JOIN source_units u ON u.projection_id = q.projection_id
                    WHERE q.revision_id = ? AND q.projection_id = ?
                      AND q.state IN ('pending', 'in_progress', 'failed')
                    GROUP BY q.queue_id
                    """,
                    (revision_id, projection_id),
                ).fetchone()
                if queue_row is None:
                    raise SemanticCommitError(
                        "progress update does not target a current queue item",
                        details={
                            "progress_index": progress_index,
                            "revision_id": revision_id,
                            "projection_id": projection_id,
                        },
                    )
                descriptor = self.adapter_registry.resolve(queue_row["extension"]).descriptor
                if (
                    queue_row["adapter_id"],
                    queue_row["adapter_version"],
                    queue_row["config_hash"],
                ) != (
                    descriptor.adapter_id,
                    descriptor.adapter_version,
                    descriptor.config_hash,
                ):
                    raise SnapshotConflictError(
                        "progress targets an outdated extraction adapter",
                        details={
                            "progress_index": progress_index,
                            "revision_id": revision_id,
                        },
                    )
                try:
                    processed_from = int(update["processed_from_ordinal"])
                    processed_through = int(update["processed_through_ordinal"])
                    next_ordinal = int(update["next_ordinal"])
                except (KeyError, TypeError, ValueError) as exc:
                    raise SemanticCommitError(
                        "progress requires integer processed_from_ordinal, "
                        "processed_through_ordinal, and next_ordinal",
                        details={"progress_index": progress_index},
                    ) from exc
                checkpoint = (
                    json.loads(queue_row["checkpoint_json"])
                    if queue_row["checkpoint_json"]
                    else None
                )
                expected_start = (
                    int(checkpoint["next_ordinal"])
                    if checkpoint and checkpoint.get("next_ordinal") is not None
                    else int(queue_row["min_ordinal"])
                )
                if (
                    processed_from != expected_start
                    or processed_through < processed_from
                    or next_ordinal != processed_through + 1
                    or processed_through > queue_row["max_ordinal"]
                ):
                    raise SemanticCommitError(
                        "progress must advance the queued projection contiguously",
                        details={
                            "progress_index": progress_index,
                            "expected_start_ordinal": expected_start,
                            "processed_from_ordinal": processed_from,
                            "processed_through_ordinal": processed_through,
                            "next_ordinal": next_ordinal,
                            "max_ordinal": queue_row["max_ordinal"],
                        },
                    )
                batch_receipt = str(update.get("batch_receipt", ""))
                if (
                    len(batch_receipt) != 64
                    or any(
                        character not in "0123456789abcdef"
                        for character in batch_receipt
                    )
                ):
                    raise SemanticCommitError(
                        "progress requires a valid server-issued batch_receipt",
                        details={"progress_index": progress_index},
                    )
                receipt_units = [
                    dict(unit)
                    for unit in connection.execute(
                        """
                        SELECT unit_id, ordinal, content_sha256
                        FROM source_units
                        WHERE projection_id = ?
                          AND ordinal BETWEEN ? AND ?
                        ORDER BY ordinal
                        """,
                        (
                            projection_id,
                            processed_from,
                            processed_through,
                        ),
                    )
                ]
                expected_receipt = _materialization_batch_receipt(
                    receipt_key=receipt_key,
                    base_snapshot_id=base_snapshot_id,
                    base_semantic_state_hash=base_semantic_state_hash,
                    queue_id=queue_row["queue_id"],
                    revision_id=revision_id,
                    projection_id=projection_id,
                    processed_from_ordinal=processed_from,
                    processed_through_ordinal=processed_through,
                    units=receipt_units,
                )
                if (
                    len(receipt_units) != processed_through - processed_from + 1
                    or not hmac.compare_digest(batch_receipt, expected_receipt)
                ):
                    raise SemanticCommitError(
                        "progress does not match the server-issued material batch",
                        details={"progress_index": progress_index},
                    )
                note = update.get("note")
                if (
                    note is not None
                    and len(str(note)) > SEMANTIC_COMMIT_MAX_STATUS_CHARS
                ):
                    raise SemanticCommitError(
                        "progress note may contain at most "
                        f"{SEMANTIC_COMMIT_MAX_STATUS_CHARS} characters",
                        details={"progress_index": progress_index},
                    )
                prepared_progress.append(
                    {
                        "queue_id": queue_row["queue_id"],
                        "revision_id": revision_id,
                        "projection_id": projection_id,
                        "processed_from_ordinal": processed_from,
                        "processed_through_ordinal": processed_through,
                        "next_ordinal": next_ordinal,
                        "min_ordinal": queue_row["min_ordinal"],
                        "max_ordinal": queue_row["max_ordinal"],
                        "total_units": queue_row["total_units"],
                        "has_more": processed_through < queue_row["max_ordinal"],
                        "batch_receipt": batch_receipt,
                        "note": str(note) if note is not None else None,
                    }
                )

            fully_progressed_revision_ids = {
                progress["revision_id"]
                for progress in prepared_progress
                if not progress["has_more"]
            }
            unauthorized_completed = sorted(
                set(completed_revision_ids) - fully_progressed_revision_ids
            )
            if unauthorized_completed:
                raise SemanticCommitError(
                    "completed revisions require a full server-issued checkpoint pass "
                    "in the same commit",
                    details={"revision_ids": unauthorized_completed},
                )
            effective_completed_revision_ids = sorted(fully_progressed_revision_ids)

            claim_ids = [claim["claim_id"] for claim in prepared_claims]
            result = {
                "commit_id": commit_id,
                "base_snapshot_id": base_snapshot_id,
                "base_semantic_state_hash": current_semantic_state["semantic_state_hash"],
                "claim_ids": claim_ids,
                "evidence_link_ids": evidence_link_ids,
                "completed_revision_ids": effective_completed_revision_ids,
                "requested_completed_revision_ids": completed_revision_ids,
                "progress_updates": prepared_progress,
                "created_by": "model",
                "validation_state": "unchecked",
            }
            now = utc_now()
            connection.execute(
                """
                    INSERT INTO semantic_commits(
                        commit_id, base_snapshot_id, idempotency_key, input_sha256,
                        created_by, materializer_version, result_json, created_at
                    ) VALUES (?, ?, ?, ?, 'model', ?, ?, ?)
                    """,
                (
                    commit_id,
                    base_snapshot_id,
                    idempotency_key,
                    input_sha256,
                    materializer_version,
                    encode_json(result),
                    now,
                ),
            )
            for claim in prepared_claims:
                connection.execute(
                    """
                        INSERT INTO atomic_claims(
                            claim_id, commit_id, body, subject, modality,
                            scope_and_conditions_json, time_window_json,
                            claim_assessment, validation_state, trust_lineage,
                            dependency_state, temporal_applicability, contest_state,
                            apparent_status, created_by, materializer_version, created_at
                        ) VALUES (?, ?, ?, ?, ?, ?, ?, ?, 'unchecked',
                                  'untrusted_source_derived', 'valid', ?, ?, ?,
                                  'model', ?, ?)
                        """,
                    (
                        claim["claim_id"],
                        commit_id,
                        claim["body"],
                        claim["subject"],
                        claim["modality"],
                        encode_json(claim["scope_and_conditions"]),
                        encode_json(claim["time_window"]),
                        claim["claim_assessment"],
                        claim["temporal_applicability"],
                        claim["contest_state"],
                        claim["apparent_status"],
                        materializer_version,
                        now,
                    ),
                )
                connection.execute(
                    """
                        INSERT INTO atomic_claims_fts(
                            claim_id, subject, body, scope_and_conditions
                        ) VALUES (?, ?, ?, ?)
                        """,
                    (
                        claim["claim_id"],
                        claim["subject"] or "",
                        claim["body"],
                        json.dumps(claim["scope_and_conditions"], ensure_ascii=False),
                    ),
                )
            for link in prepared_links:
                connection.execute(
                    """
                        INSERT INTO evidence_links(
                            evidence_link_id, claim_id, source_unit_id,
                            source_revision_id, source_span_json, stance,
                            qualifier, applicability_json
                        ) VALUES (?, ?, ?, ?, ?, ?, ?, ?)
                        """,
                    (
                        link["evidence_link_id"],
                        link["claim_id"],
                        link["source_unit_id"],
                        link["source_revision_id"],
                        encode_json(link["source_span"]),
                        link["stance"],
                        link["qualifier"],
                        encode_json(link["applicability"]),
                    ),
                )
            queue_updates = []
            for progress in prepared_progress:
                previous_row = connection.execute(
                    """
                    SELECT checkpoint_json FROM interpretation_queue
                    WHERE queue_id = ?
                    """,
                    (progress["queue_id"],),
                ).fetchone()
                previous_checkpoint = (
                    json.loads(previous_row["checkpoint_json"])
                    if previous_row and previous_row["checkpoint_json"]
                    else {}
                )
                ranges = [
                    list(value)
                    for value in previous_checkpoint.get("processed_ordinal_ranges", [])
                    if isinstance(value, list) and len(value) == 2
                ]
                if (
                    ranges
                    and ranges[-1][1] + 1 == progress["processed_from_ordinal"]
                ):
                    ranges[-1][1] = progress["processed_through_ordinal"]
                else:
                    ranges.append(
                        [
                            progress["processed_from_ordinal"],
                            progress["processed_through_ordinal"],
                        ]
                    )
                remaining = (
                    [[progress["next_ordinal"], progress["max_ordinal"]]]
                    if progress["has_more"]
                    else []
                )
                checkpoint_payload = {
                    "schema_version": 1,
                    "projection_id": progress["projection_id"],
                    "processed_ordinal_ranges": ranges,
                    "remaining_ordinal_ranges": remaining,
                    "next_ordinal": (
                        progress["next_ordinal"] if progress["has_more"] else None
                    ),
                    "total_units": progress["total_units"],
                    "updated_by_commit_id": commit_id,
                }
                if progress["note"] is not None:
                    checkpoint_payload["note"] = progress["note"]
                is_complete = progress["revision_id"] in effective_completed_revision_ids
                connection.execute(
                    """
                    UPDATE interpretation_queue
                    SET state = ?, reason = ?, checkpoint_json = ?, updated_at = ?
                    WHERE queue_id = ?
                    """,
                    (
                        "complete" if is_complete else "in_progress",
                        "semantic_commit" if is_complete else "semantic_partial_commit",
                        encode_json(checkpoint_payload),
                        now,
                        progress["queue_id"],
                    ),
                )
                queue_updates.append(
                    {
                        "queue_id": progress["queue_id"],
                        "revision_id": progress["revision_id"],
                        "state": "complete" if is_complete else "in_progress",
                        "checkpoint": checkpoint_payload,
                    }
                )
            result["queue_updates"] = queue_updates
            result["semantic_state"] = _semantic_state(connection)
            connection.execute(
                """
                UPDATE semantic_commits SET result_json = ?
                WHERE commit_id = ?
                """,
                (encode_json(result), commit_id),
            )
        return result

    def semantic_context(
        self,
        corpus_id: str,
        *,
        query: str | None = None,
        limit: int = 50,
    ) -> dict:
        if not 1 <= limit <= 200:
            raise BudgetExceededError(
                "semantic context limit must be between 1 and 200",
                details={"limit": limit, "maximum": 200},
            )
        if query is not None and len(query) > 2_000:
            raise BudgetExceededError(
                "semantic context query must contain at most 2000 characters",
                details={"query_chars": len(query), "maximum": 2_000},
            )
        with corpus_read_connection(self.data_root, corpus_id) as connection:
            snapshot = connection.execute(
                """
                SELECT snapshot_id, coverage_state, document_count, supported_document_count
                FROM snapshots
                WHERE state = 'complete'
                ORDER BY rowid DESC LIMIT 1
                """
            ).fetchone()
            materialization_by_projection: dict[str, dict] = {}
            queue_coverage = {
                "complete": 0,
                "in_progress": 0,
                "pending": 0,
                "failed": 0,
                "stale": 0,
                "outdated": 0,
                "untracked": 0,
            }
            if snapshot is not None:
                queue_rows = connection.execute(
                    """
                    SELECT sd.projection_id, d.extension,
                           p.adapter_id, p.adapter_version, p.config_hash,
                           q.state
                    FROM snapshot_documents sd
                    JOIN documents d ON d.document_id = sd.document_id
                    JOIN extraction_projections p
                      ON p.projection_id = sd.projection_id
                    LEFT JOIN interpretation_queue q
                      ON q.revision_id = sd.revision_id
                     AND q.projection_id = sd.projection_id
                    WHERE sd.snapshot_id = ?
                    """,
                    (snapshot["snapshot_id"],),
                )
                for queue_row in queue_rows:
                    try:
                        descriptor = self.adapter_registry.resolve(
                            queue_row["extension"]
                        ).descriptor
                        current = (
                            queue_row["adapter_id"],
                            queue_row["adapter_version"],
                            queue_row["config_hash"],
                        ) == (
                            descriptor.adapter_id,
                            descriptor.adapter_version,
                            descriptor.config_hash,
                        )
                    except ExtractionError:
                        current = False
                    state = queue_row["state"] or "untracked"
                    coverage_state = state if current else "outdated"
                    queue_coverage[coverage_state] = (
                        queue_coverage.get(coverage_state, 0) + 1
                    )
            claims = []
            orientation_by_document: dict[str, dict] = {}
            context_payload_bytes = 4_096
            if snapshot is not None:
                has_query = bool(query and query.strip())
                fts_query = (
                    '"' + query.strip().replace('"', '""') + '"' if has_query else None
                )
                surfaced_by = (
                    "semantic_claim_lexical_acquisition"
                    if has_query
                    else "recent_semantic_claims"
                )
                page_size = max(50, min(400, limit * 2))
                offset = 0
                while len(claims) < limit:
                    if has_query:
                        claim_rows = connection.execute(
                            """
                            SELECT c.claim_id,
                                   bm25(atomic_claims_fts) AS acquisition_score,
                                   LENGTH(CAST(COALESCE(c.claim_id, '') AS BLOB))
                                   + LENGTH(CAST(COALESCE(c.commit_id, '') AS BLOB))
                                   + LENGTH(CAST(COALESCE(c.body, '') AS BLOB))
                                   + LENGTH(CAST(COALESCE(c.subject, '') AS BLOB))
                                   + LENGTH(CAST(COALESCE(c.modality, '') AS BLOB))
                                   + LENGTH(
                                       CAST(
                                           COALESCE(c.scope_and_conditions_json, '')
                                           AS BLOB
                                       )
                                     )
                                   + LENGTH(
                                       CAST(COALESCE(c.time_window_json, '') AS BLOB)
                                     )
                                   + LENGTH(
                                       CAST(COALESCE(c.claim_assessment, '') AS BLOB)
                                     )
                                   + LENGTH(
                                       CAST(COALESCE(c.validation_state, '') AS BLOB)
                                     )
                                   + LENGTH(
                                       CAST(COALESCE(c.trust_lineage, '') AS BLOB)
                                     )
                                   + LENGTH(
                                       CAST(COALESCE(c.dependency_state, '') AS BLOB)
                                     )
                                   + LENGTH(
                                       CAST(
                                           COALESCE(c.temporal_applicability, '')
                                           AS BLOB
                                       )
                                     )
                                   + LENGTH(
                                       CAST(COALESCE(c.contest_state, '') AS BLOB)
                                     )
                                   + LENGTH(
                                       CAST(COALESCE(c.apparent_status, '') AS BLOB)
                                     )
                                   + LENGTH(CAST(COALESCE(c.created_by, '') AS BLOB))
                                   + LENGTH(
                                       CAST(COALESCE(c.materializer_version, '') AS BLOB)
                                     )
                                   + LENGTH(CAST(COALESCE(c.created_at, '') AS BLOB))
                                   + 1024 AS stored_payload_bytes
                            FROM atomic_claims_fts f
                            JOIN atomic_claims c ON c.claim_id = f.claim_id
                            WHERE atomic_claims_fts MATCH ?
                              AND c.dependency_state = 'valid'
                              AND NOT EXISTS (
                                  SELECT 1
                                  FROM evidence_links dependency
                                  JOIN source_units dependency_unit
                                    ON dependency_unit.unit_id = dependency.source_unit_id
                                  JOIN revisions dependency_revision
                                    ON dependency_revision.revision_id
                                     = dependency_unit.revision_id
                                  JOIN documents dependency_document
                                    ON dependency_document.document_id
                                     = dependency_revision.document_id
                                  LEFT JOIN extraction_projections active_projection
                                    ON active_projection.revision_id
                                     = dependency_unit.revision_id
                                   AND active_projection.is_active = 1
                                  LEFT JOIN snapshot_documents dependency_snapshot
                                    ON dependency_snapshot.snapshot_id = ?
                                   AND dependency_snapshot.document_id
                                     = dependency_document.document_id
                                   AND dependency_snapshot.revision_id
                                     = dependency_unit.revision_id
                                   AND dependency_snapshot.projection_id
                                     = dependency_unit.projection_id
                                  WHERE dependency.claim_id = c.claim_id
                                    AND (
                                        dependency_snapshot.snapshot_id IS NULL
                                        OR dependency_document.deleted_at IS NOT NULL
                                        OR dependency_document.current_revision_id
                                           != dependency_unit.revision_id
                                        OR active_projection.projection_id IS NULL
                                        OR active_projection.projection_id
                                           != dependency_unit.projection_id
                                        OR dependency_revision.source_size
                                           != dependency_document.logical_size
                                        OR dependency_revision.source_modified_ns
                                           != dependency_document.modified_ns
                                        OR dependency_revision.source_changed_ns
                                           != dependency_document.changed_ns
                                        OR dependency_revision.source_device
                                           != dependency_document.device
                                        OR dependency_revision.source_inode
                                           != dependency_document.inode
                                    )
                              )
                            ORDER BY bm25(atomic_claims_fts), c.claim_id
                            LIMIT ? OFFSET ?
                            """,
                            (
                                fts_query,
                                snapshot["snapshot_id"],
                                page_size,
                                offset,
                            ),
                        ).fetchall()
                    else:
                        claim_rows = connection.execute(
                            """
                            SELECT c.claim_id, NULL AS acquisition_score,
                                   LENGTH(CAST(COALESCE(c.claim_id, '') AS BLOB))
                                   + LENGTH(CAST(COALESCE(c.commit_id, '') AS BLOB))
                                   + LENGTH(CAST(COALESCE(c.body, '') AS BLOB))
                                   + LENGTH(CAST(COALESCE(c.subject, '') AS BLOB))
                                   + LENGTH(CAST(COALESCE(c.modality, '') AS BLOB))
                                   + LENGTH(
                                       CAST(
                                           COALESCE(c.scope_and_conditions_json, '')
                                           AS BLOB
                                       )
                                     )
                                   + LENGTH(
                                       CAST(COALESCE(c.time_window_json, '') AS BLOB)
                                     )
                                   + LENGTH(
                                       CAST(COALESCE(c.claim_assessment, '') AS BLOB)
                                     )
                                   + LENGTH(
                                       CAST(COALESCE(c.validation_state, '') AS BLOB)
                                     )
                                   + LENGTH(
                                       CAST(COALESCE(c.trust_lineage, '') AS BLOB)
                                     )
                                   + LENGTH(
                                       CAST(COALESCE(c.dependency_state, '') AS BLOB)
                                     )
                                   + LENGTH(
                                       CAST(
                                           COALESCE(c.temporal_applicability, '')
                                           AS BLOB
                                       )
                                     )
                                   + LENGTH(
                                       CAST(COALESCE(c.contest_state, '') AS BLOB)
                                     )
                                   + LENGTH(
                                       CAST(COALESCE(c.apparent_status, '') AS BLOB)
                                     )
                                   + LENGTH(CAST(COALESCE(c.created_by, '') AS BLOB))
                                   + LENGTH(
                                       CAST(COALESCE(c.materializer_version, '') AS BLOB)
                                     )
                                   + LENGTH(CAST(COALESCE(c.created_at, '') AS BLOB))
                                   + 1024 AS stored_payload_bytes
                            FROM atomic_claims c
                            WHERE c.dependency_state = 'valid'
                              AND NOT EXISTS (
                                  SELECT 1
                                  FROM evidence_links dependency
                                  JOIN source_units dependency_unit
                                    ON dependency_unit.unit_id = dependency.source_unit_id
                                  JOIN revisions dependency_revision
                                    ON dependency_revision.revision_id
                                     = dependency_unit.revision_id
                                  JOIN documents dependency_document
                                    ON dependency_document.document_id
                                     = dependency_revision.document_id
                                  LEFT JOIN extraction_projections active_projection
                                    ON active_projection.revision_id
                                     = dependency_unit.revision_id
                                   AND active_projection.is_active = 1
                                  LEFT JOIN snapshot_documents dependency_snapshot
                                    ON dependency_snapshot.snapshot_id = ?
                                   AND dependency_snapshot.document_id
                                     = dependency_document.document_id
                                   AND dependency_snapshot.revision_id
                                     = dependency_unit.revision_id
                                   AND dependency_snapshot.projection_id
                                     = dependency_unit.projection_id
                                  WHERE dependency.claim_id = c.claim_id
                                    AND (
                                        dependency_snapshot.snapshot_id IS NULL
                                        OR dependency_document.deleted_at IS NOT NULL
                                        OR dependency_document.current_revision_id
                                           != dependency_unit.revision_id
                                        OR active_projection.projection_id IS NULL
                                        OR active_projection.projection_id
                                           != dependency_unit.projection_id
                                        OR dependency_revision.source_size
                                           != dependency_document.logical_size
                                        OR dependency_revision.source_modified_ns
                                           != dependency_document.modified_ns
                                        OR dependency_revision.source_changed_ns
                                           != dependency_document.changed_ns
                                        OR dependency_revision.source_device
                                           != dependency_document.device
                                        OR dependency_revision.source_inode
                                           != dependency_document.inode
                                    )
                              )
                            ORDER BY c.created_at DESC, c.claim_id
                            LIMIT ? OFFSET ?
                            """,
                            (snapshot["snapshot_id"], page_size, offset),
                        ).fetchall()
                    if not claim_rows:
                        break
                    offset += len(claim_rows)
                    for claim_inventory in claim_rows:
                        projection_rows = connection.execute(
                            """
                            SELECT DISTINCT d.extension, p.adapter_id,
                                   p.adapter_version, p.config_hash
                            FROM evidence_links e
                            JOIN source_units u ON u.unit_id = e.source_unit_id
                            JOIN revisions r ON r.revision_id = u.revision_id
                            JOIN documents d ON d.document_id = r.document_id
                            JOIN extraction_projections p
                              ON p.projection_id = u.projection_id
                            JOIN snapshot_documents dependency_snapshot
                              ON dependency_snapshot.snapshot_id = ?
                             AND dependency_snapshot.document_id = d.document_id
                             AND dependency_snapshot.revision_id = u.revision_id
                             AND dependency_snapshot.projection_id = u.projection_id
                            WHERE e.claim_id = ?
                            """,
                            (
                                snapshot["snapshot_id"],
                                claim_inventory["claim_id"],
                            ),
                        ).fetchall()
                        projection_current = bool(projection_rows)
                        for projection_row in projection_rows:
                            try:
                                descriptor = self.adapter_registry.resolve(
                                    projection_row["extension"]
                                ).descriptor
                            except ExtractionError:
                                projection_current = False
                                break
                            if (
                                projection_row["adapter_id"],
                                projection_row["adapter_version"],
                                projection_row["config_hash"],
                            ) != (
                                descriptor.adapter_id,
                                descriptor.adapter_version,
                                descriptor.config_hash,
                            ):
                                projection_current = False
                                break
                        if not projection_current:
                            continue

                        link_inventory = connection.execute(
                            """
                            SELECT COUNT(*) AS link_count,
                                   COALESCE(
                                       SUM(
                                           LENGTH(
                                               CAST(
                                                   COALESCE(e.evidence_link_id, '')
                                                   AS BLOB
                                               )
                                             )
                                           + LENGTH(
                                               CAST(COALESCE(e.claim_id, '') AS BLOB)
                                             )
                                           + LENGTH(
                                               CAST(
                                                   COALESCE(e.source_unit_id, '')
                                                   AS BLOB
                                               )
                                             )
                                           + LENGTH(
                                               CAST(
                                                   COALESCE(e.source_revision_id, '')
                                                   AS BLOB
                                               )
                                             )
                                           + LENGTH(
                                               CAST(
                                                   COALESCE(e.source_span_json, '')
                                                   AS BLOB
                                               )
                                             )
                                           + LENGTH(
                                               CAST(COALESCE(e.stance, '') AS BLOB)
                                             )
                                           + LENGTH(
                                               CAST(COALESCE(e.qualifier, '') AS BLOB)
                                             )
                                           + LENGTH(
                                               CAST(
                                                   COALESCE(e.applicability_json, '')
                                                   AS BLOB
                                               )
                                             )
                                           + LENGTH(
                                               CAST(
                                                   COALESCE(u.source_anchor_json, '')
                                                   AS BLOB
                                               )
                                             )
                                           + LENGTH(
                                               CAST(
                                                   COALESCE(u.projection_id, '')
                                                   AS BLOB
                                               )
                                             )
                                           + LENGTH(
                                               CAST(COALESCE(d.document_id, '') AS BLOB)
                                             )
                                           + LENGTH(
                                               CAST(COALESCE(d.relative_path, '') AS BLOB)
                                             )
                                           + LENGTH(
                                               CAST(COALESCE(q.state, '') AS BLOB)
                                             )
                                           + LENGTH(
                                               CAST(
                                                   COALESCE(q.checkpoint_json, '')
                                                   AS BLOB
                                               )
                                             )
                                           + 1024
                                       ),
                                       0
                                   ) AS stored_payload_bytes
                            FROM evidence_links e
                            JOIN source_units u ON u.unit_id = e.source_unit_id
                            JOIN revisions r ON r.revision_id = u.revision_id
                            JOIN documents d ON d.document_id = r.document_id
                            LEFT JOIN interpretation_queue q
                              ON q.projection_id = u.projection_id
                            JOIN snapshot_documents dependency_snapshot
                              ON dependency_snapshot.snapshot_id = ?
                             AND dependency_snapshot.document_id = d.document_id
                             AND dependency_snapshot.revision_id = u.revision_id
                             AND dependency_snapshot.projection_id = u.projection_id
                            WHERE e.claim_id = ?
                            """,
                            (
                                snapshot["snapshot_id"],
                                claim_inventory["claim_id"],
                            ),
                        ).fetchone()
                        candidate_stored_payload_bytes = (
                            claim_inventory["stored_payload_bytes"]
                            + link_inventory["stored_payload_bytes"]
                        )
                        projected_payload_bytes = (
                            context_payload_bytes
                            + candidate_stored_payload_bytes
                        )
                        if (
                            projected_payload_bytes
                            > SEMANTIC_CONTEXT_MAX_SERIALIZED_BYTES
                        ):
                            raise BudgetExceededError(
                                "semantic context exceeds the serialized response budget",
                                details={
                                    "requested_limit": limit,
                                    "returned_claim_count": len(claims),
                                    "candidate_claim_id": claim_inventory[
                                        "claim_id"
                                    ],
                                    "candidate_evidence_count": link_inventory[
                                        "link_count"
                                    ],
                                    "response_bytes_before_candidate": (
                                        context_payload_bytes
                                    ),
                                    "candidate_stored_payload_bytes": (
                                        candidate_stored_payload_bytes
                                    ),
                                    "projected_payload_bytes": (
                                        projected_payload_bytes
                                    ),
                                    "maximum_serialized_bytes": (
                                        SEMANTIC_CONTEXT_MAX_SERIALIZED_BYTES
                                    ),
                                    "retry_with_lower": ["limit"],
                                    "retry_with_narrower_query": True,
                                },
                            )

                        claim_row = connection.execute(
                            """
                            SELECT c.*, ? AS acquisition_score
                            FROM atomic_claims c
                            WHERE c.claim_id = ?
                            """,
                            (
                                claim_inventory["acquisition_score"],
                                claim_inventory["claim_id"],
                            ),
                        ).fetchone()
                        claim = dict(claim_row)
                        claim["scope_and_conditions"] = json.loads(
                            claim.pop("scope_and_conditions_json")
                        )
                        claim["time_window"] = json.loads(
                            claim.pop("time_window_json")
                        )
                        links = connection.execute(
                            """
                            SELECT e.*, u.source_anchor_json, u.projection_id,
                                   d.document_id, d.relative_path, d.extension,
                                   p.adapter_id, p.adapter_version, p.config_hash
                            FROM evidence_links e
                            JOIN source_units u ON u.unit_id = e.source_unit_id
                            JOIN revisions r ON r.revision_id = u.revision_id
                            JOIN documents d ON d.document_id = r.document_id
                            JOIN extraction_projections p
                              ON p.projection_id = u.projection_id
                            JOIN snapshot_documents dependency_snapshot
                              ON dependency_snapshot.snapshot_id = ?
                             AND dependency_snapshot.document_id = d.document_id
                             AND dependency_snapshot.revision_id = u.revision_id
                             AND dependency_snapshot.projection_id = u.projection_id
                            WHERE e.claim_id = ?
                            ORDER BY e.evidence_link_id
                            """,
                            (snapshot["snapshot_id"], claim["claim_id"]),
                        ).fetchall()
                        claim["evidence"] = []
                        for link_row in links:
                            link = dict(link_row)
                            for internal_field in (
                                "extension",
                                "adapter_id",
                                "adapter_version",
                                "config_hash",
                            ):
                                link.pop(internal_field)
                            link["source_span"] = json.loads(
                                link.pop("source_span_json")
                            )
                            link["applicability"] = json.loads(
                                link.pop("applicability_json")
                            )
                            link["source_anchor"] = json.loads(
                                link.pop("source_anchor_json")
                            )
                            claim["evidence"].append(link)
                            projection_materialization = (
                                materialization_by_projection.get(
                                    link["projection_id"]
                                )
                            )
                            if projection_materialization is None:
                                queue_row = connection.execute(
                                    """
                                    SELECT state, checkpoint_json
                                    FROM interpretation_queue
                                    WHERE projection_id = ?
                                    """,
                                    (link["projection_id"],),
                                ).fetchone()
                                projection_materialization = {
                                    "state": (
                                        queue_row["state"]
                                        if queue_row is not None
                                        else "untracked"
                                    ),
                                    "checkpoint": (
                                        json.loads(queue_row["checkpoint_json"])
                                        if queue_row is not None
                                        and queue_row["checkpoint_json"]
                                        else None
                                    ),
                                }
                                materialization_by_projection[
                                    link["projection_id"]
                                ] = projection_materialization
                            orientation = orientation_by_document.setdefault(
                                link["document_id"],
                                {
                                    "document_id": link["document_id"],
                                    "relative_path": link["relative_path"],
                                    "revision_id": link["source_revision_id"],
                                    "projection_id": link["projection_id"],
                                    "materialization": projection_materialization,
                                    "claims": [],
                                },
                            )
                            if not any(
                                existing["claim_id"] == claim["claim_id"]
                                for existing in orientation["claims"]
                            ):
                                orientation["claims"].append(
                                    {
                                        "claim_id": claim["claim_id"],
                                        "body": claim["body"],
                                        "subject": claim["subject"],
                                        "claim_assessment": claim[
                                            "claim_assessment"
                                        ],
                                        "temporal_applicability": claim[
                                            "temporal_applicability"
                                        ],
                                        "contest_state": claim["contest_state"],
                                    }
                                )
                        claim["surfaced_by"] = surfaced_by
                        claim["acquisition_is_evidence"] = False
                        claims.append(claim)
                        context_payload_bytes = (
                            len(
                                encode_json(
                                    {
                                        "query": query,
                                        "count": len(claims),
                                        "claims": claims,
                                        "document_orientations": list(
                                            orientation_by_document.values()
                                        ),
                                    }
                                ).encode()
                            )
                            + 4_096
                        )
                        if (
                            context_payload_bytes
                            > SEMANTIC_CONTEXT_MAX_SERIALIZED_BYTES
                        ):
                            raise BudgetExceededError(
                                "semantic context exceeds the serialized response budget",
                                details={
                                    "requested_limit": limit,
                                    "returned_claim_count": len(claims),
                                    "document_orientation_count": len(
                                        orientation_by_document
                                    ),
                                    "serialized_bytes_with_reserved_overhead": (
                                        context_payload_bytes
                                    ),
                                    "maximum_serialized_bytes": (
                                        SEMANTIC_CONTEXT_MAX_SERIALIZED_BYTES
                                    ),
                                    "retry_with_lower": ["limit"],
                                    "retry_with_narrower_query": True,
                                },
                            )
                        if len(claims) >= limit:
                            break
                    if len(claim_rows) < page_size:
                        break
            claim_revision_count = 0
            if snapshot is not None:
                def projection_is_current(
                    extension: str,
                    adapter_id: str,
                    adapter_version: str,
                    config_hash: str,
                ) -> int:
                    try:
                        descriptor = self.adapter_registry.resolve(extension).descriptor
                    except ExtractionError:
                        return 0
                    return int(
                        (
                            adapter_id,
                            adapter_version,
                            config_hash,
                        )
                        == (
                            descriptor.adapter_id,
                            descriptor.adapter_version,
                            descriptor.config_hash,
                        )
                    )

                connection.create_function(
                    "corpus_projection_is_current",
                    4,
                    projection_is_current,
                )
                coverage = connection.execute(
                    """
                    SELECT COUNT(DISTINCT e.source_revision_id) AS revision_count
                    FROM evidence_links e
                    JOIN atomic_claims c ON c.claim_id = e.claim_id
                    JOIN source_units u ON u.unit_id = e.source_unit_id
                    JOIN revisions r ON r.revision_id = u.revision_id
                    JOIN documents d ON d.document_id = r.document_id
                    JOIN extraction_projections p
                      ON p.projection_id = u.projection_id AND p.is_active = 1
                    JOIN snapshot_documents sd
                      ON sd.snapshot_id = ?
                     AND sd.revision_id = e.source_revision_id
                     AND sd.projection_id = u.projection_id
                    WHERE c.dependency_state = 'valid'
                      AND d.deleted_at IS NULL
                      AND d.current_revision_id = u.revision_id
                      AND r.source_size = d.logical_size
                      AND r.source_modified_ns = d.modified_ns
                      AND r.source_changed_ns = d.changed_ns
                      AND r.source_device = d.device
                      AND r.source_inode = d.inode
                      AND corpus_projection_is_current(
                              d.extension,
                              p.adapter_id,
                              p.adapter_version,
                              p.config_hash
                          ) = 1
                      AND NOT EXISTS (
                          SELECT 1
                          FROM evidence_links dependency
                          JOIN source_units dependency_unit
                            ON dependency_unit.unit_id = dependency.source_unit_id
                          JOIN revisions dependency_revision
                            ON dependency_revision.revision_id
                             = dependency_unit.revision_id
                          JOIN documents dependency_document
                            ON dependency_document.document_id
                             = dependency_revision.document_id
                          LEFT JOIN extraction_projections active_projection
                            ON active_projection.revision_id
                             = dependency_unit.revision_id
                           AND active_projection.is_active = 1
                          LEFT JOIN snapshot_documents dependency_snapshot
                            ON dependency_snapshot.snapshot_id = ?
                           AND dependency_snapshot.document_id
                             = dependency_document.document_id
                           AND dependency_snapshot.revision_id
                             = dependency_unit.revision_id
                           AND dependency_snapshot.projection_id
                             = dependency_unit.projection_id
                          WHERE dependency.claim_id = c.claim_id
                            AND (
                                dependency_snapshot.snapshot_id IS NULL
                                OR dependency_document.deleted_at IS NOT NULL
                                OR dependency_document.current_revision_id
                                   != dependency_unit.revision_id
                                OR active_projection.projection_id IS NULL
                                OR active_projection.projection_id
                                   != dependency_unit.projection_id
                                OR dependency_revision.source_size
                                   != dependency_document.logical_size
                                OR dependency_revision.source_modified_ns
                                   != dependency_document.modified_ns
                                OR dependency_revision.source_changed_ns
                                   != dependency_document.changed_ns
                                OR dependency_revision.source_device
                                   != dependency_document.device
                                OR dependency_revision.source_inode
                                   != dependency_document.inode
                            )
                      )
                      AND NOT EXISTS (
                          SELECT 1
                          FROM evidence_links adapter_dependency
                          JOIN source_units adapter_unit
                            ON adapter_unit.unit_id
                             = adapter_dependency.source_unit_id
                          JOIN revisions adapter_revision
                            ON adapter_revision.revision_id
                             = adapter_unit.revision_id
                          JOIN documents adapter_document
                            ON adapter_document.document_id
                             = adapter_revision.document_id
                          JOIN extraction_projections adapter_projection
                            ON adapter_projection.projection_id
                             = adapter_unit.projection_id
                          WHERE adapter_dependency.claim_id = c.claim_id
                            AND corpus_projection_is_current(
                                    adapter_document.extension,
                                    adapter_projection.adapter_id,
                                    adapter_projection.adapter_version,
                                    adapter_projection.config_hash
                                ) = 0
                      )
                    """,
                    (snapshot["snapshot_id"], snapshot["snapshot_id"]),
                ).fetchone()
                claim_revision_count = coverage["revision_count"]
            semantic_state = _semantic_state(connection)
        response = {
            "query": query,
            "count": len(claims),
            "claims": claims,
            "document_orientations": list(orientation_by_document.values()),
            "source_snapshot": dict(snapshot) if snapshot is not None else None,
            "semantic_state": semantic_state,
            "semantic_coverage": {
                "enrichment_queue": queue_coverage,
                "snapshot_revisions_with_valid_claims": claim_revision_count,
                "snapshot_document_count": snapshot["document_count"] if snapshot else 0,
                "snapshot_revisions_without_valid_claims": (
                    snapshot["document_count"] - claim_revision_count
                    if snapshot
                    else 0
                ),
                "claims_returned": len(claims),
                "corpus_wide_absence_supported": False,
            },
            "notice": (
                "Claims are model-created derived interpretations. "
                "Follow evidence links to source units before reuse."
            ),
        }
        serialized_bytes = len(encode_json(response).encode())
        if serialized_bytes > SEMANTIC_CONTEXT_MAX_SERIALIZED_BYTES:
            raise BudgetExceededError(
                "semantic context exceeds the serialized response budget",
                details={
                    "requested_limit": limit,
                    "returned_claim_count": len(claims),
                    "document_orientation_count": len(orientation_by_document),
                    "serialized_bytes": serialized_bytes,
                    "maximum_serialized_bytes": (
                        SEMANTIC_CONTEXT_MAX_SERIALIZED_BYTES
                    ),
                    "retry_with_lower": ["limit"],
                    "retry_with_narrower_query": True,
                },
            )
        return response

    def evaluate_extraction_golden(self, annotation: dict) -> dict:
        annotation = validate_golden_annotation(annotation)
        annotation_id = annotation["annotation_id"]
        corpus_id = annotation["subject"]["corpus_id"]
        document_id = annotation["subject"]["document_id"]
        try:
            corpus = get_corpus(self.data_root, corpus_id)
        except CorpusError as exc:
            raise GoldenEvaluationError(
                "golden annotation subject could not be resolved",
                details={
                    "annotation_id": annotation_id,
                    "subject_match": False,
                },
            ) from exc
        paths = self._paths(corpus_id)
        with corpus_read_connection(self.data_root, corpus_id) as connection:
            row = connection.execute(
                """
                SELECT d.document_id, d.current_revision_id, d.logical_size,
                       d.modified_ns, d.changed_ns, d.device, d.inode, d.deleted_at,
                       d.relative_path, d.residency_state, d.is_dataless,
                       r.revision_id, r.sha256, r.source_size,
                       r.source_modified_ns, r.source_changed_ns,
                       r.source_device, r.source_inode,
                       p.projection_id, p.adapter_id, p.adapter_version,
                       p.config_hash, p.completeness_state
                FROM documents d
                LEFT JOIN revisions r ON r.revision_id = d.current_revision_id
                LEFT JOIN extraction_projections p
                  ON p.revision_id = r.revision_id AND p.is_active = 1
                WHERE d.document_id = ?
                """,
                (document_id,),
            ).fetchone()
            if row is None:
                raise GoldenEvaluationError(
                    "golden annotation subject could not be resolved",
                    details={
                        "annotation_id": annotation_id,
                        "subject_match": False,
                    },
                )
            observation = {
                "corpus_id": corpus_id,
                "document_id": document_id,
                "revision_sha256": row["sha256"],
                "projection_id": row["projection_id"],
                "adapter_id": row["adapter_id"],
                "adapter_version": row["adapter_version"],
                "config_hash": row["config_hash"],
                "completeness_state": row["completeness_state"],
                "source_observation_current": False,
                "unit_count": 0,
                "unit_type_counts": {},
                "derivation_method_counts": {},
                "geometry_coverage": 0.0,
                "confidence_coverage": 0.0,
                "issue_codes": [],
                "projection_observation_sha256": None,
                "unit_content_hashes_current": True,
            }
            source_root = Path(corpus["source_root"])
            source_path = source_root / row["relative_path"]
            observation["source_observation_current"] = bool(
                row["deleted_at"] is None
                and row["revision_id"] is not None
                and not row["is_dataless"]
                and row["source_size"] == row["logical_size"]
                and row["source_modified_ns"] == row["modified_ns"]
                and row["source_changed_ns"] == row["changed_ns"]
                and row["source_device"] == row["device"]
                and row["source_inode"] == row["inode"]
                and _resident_source_matches_revision(
                    source_path=source_path,
                    source_root=source_root,
                    staging_root=paths.staging,
                    document=dict(row),
                    revision_sha256=row["sha256"],
                )
            )
            if row["projection_id"] is not None:
                unit_rows = connection.execute(
                    """
                    SELECT ordinal, unit_type, normalized_content, content_sha256,
                           extraction_issues_json, derivation_method,
                           structure_path_json, geometry_json, confidence,
                           quality_flags_json,
                           geometry_json <> '{}' AS has_geometry,
                           confidence IS NOT NULL AS has_confidence
                    FROM source_units WHERE projection_id = ?
                    ORDER BY ordinal
                    """,
                    (row["projection_id"],),
                ).fetchall()
                observation["unit_count"] = len(unit_rows)
                for unit in unit_rows:
                    observation["unit_type_counts"][unit["unit_type"]] = (
                        observation["unit_type_counts"].get(unit["unit_type"], 0) + 1
                    )
                    observation["derivation_method_counts"][unit["derivation_method"]] = (
                        observation["derivation_method_counts"].get(
                            unit["derivation_method"], 0
                        )
                        + 1
                    )
                if unit_rows:
                    observation["geometry_coverage"] = sum(
                        bool(unit["has_geometry"]) for unit in unit_rows
                    ) / len(unit_rows)
                    observation["confidence_coverage"] = sum(
                        bool(unit["has_confidence"]) for unit in unit_rows
                    ) / len(unit_rows)
                recomputed_content_hashes = [
                    hashlib.sha256(unit["normalized_content"].encode("utf-8")).hexdigest()
                    for unit in unit_rows
                ]
                observation["unit_content_hashes_current"] = all(
                    actual == unit["content_sha256"]
                    for actual, unit in zip(
                        recomputed_content_hashes,
                        unit_rows,
                        strict=True,
                    )
                )
                projection_issues = [
                    {
                        "stage": issue["stage"],
                        "severity": issue["severity"],
                        "code": issue["code"],
                        "details": json.loads(issue["details_json"]),
                        "structural_locator": json.loads(
                            issue["structural_locator_json"]
                        ),
                    }
                    for issue in connection.execute(
                        """
                        SELECT stage, severity, code, details_json,
                               structural_locator_json
                        FROM extraction_issues
                        WHERE projection_id = ? AND lifecycle_state = 'active'
                        ORDER BY stage, code, locator_key
                        """,
                        (row["projection_id"],),
                    )
                ]
                observation["projection_observation_sha256"] = (
                    projection_observation_sha256(
                        [
                            {
                                "ordinal": unit["ordinal"],
                                "unit_type": unit["unit_type"],
                                "content_sha256": recomputed_content_hashes[index],
                                "derivation_method": unit["derivation_method"],
                                "structure_path": json.loads(
                                    unit["structure_path_json"]
                                ),
                                "geometry": json.loads(unit["geometry_json"]),
                                "confidence": unit["confidence"],
                                "quality_flags": json.loads(
                                    unit["quality_flags_json"]
                                ),
                                "issues": json.loads(
                                    unit["extraction_issues_json"]
                                ),
                            }
                            for index, unit in enumerate(unit_rows)
                        ],
                        projection_issues=projection_issues,
                    )
                )
                observation["issue_codes"] = sorted(
                    {issue["code"] for issue in projection_issues}
                )
        return evaluate_projection_observation(annotation, observation)

    def doctor(self, corpus_id: str | None = None) -> dict:
        checks: list[dict] = []
        checks.append(
            {
                "name": "data_root_outside_cloud_source",
                "status": "pass",
                "data_root": str(self.data_root),
            }
        )
        try:
            import docx  # noqa: F401
            import openpyxl  # noqa: F401
            import pptx  # noqa: F401
            import pypdf  # noqa: F401

            checks.append({"name": "document_extractors", "status": "pass"})
        except ImportError as exc:
            checks.append({"name": "document_extractors", "status": "fail", "message": str(exc)})
        try:
            import olefile  # noqa: F401

            checks.append({"name": "binary_hwp_reader", "status": "pass"})
        except ImportError as exc:
            checks.append(
                {
                    "name": "binary_hwp_reader",
                    "status": "fail",
                    "message": str(exc),
                }
            )
        native_source = native_source_path()
        checks.append(
            {
                "name": "native_helper_source",
                "status": "pass"
                if native_source is not None and native_source.is_file()
                else "pending",
            }
        )
        pdf_descriptor = self.adapter_registry.resolve("pdf").descriptor
        checks.append(
            {
                "name": "pdf_extraction_adapter",
                "status": "pass",
                "adapter_id": pdf_descriptor.adapter_id,
                "supports_ocr": pdf_descriptor.capabilities.supports_ocr,
                "supports_geometry": pdf_descriptor.capabilities.supports_geometry,
            }
        )
        if corpus_id:
            corpus = get_corpus(self.data_root, corpus_id)
            source_root = Path(corpus["source_root"])
            checks.append(
                {
                    "name": "source_root",
                    "status": "pass" if source_root.is_dir() else "fail",
                    "path": str(source_root),
                }
            )
        return {
            "ok": all(check["status"] != "fail" for check in checks),
            "checks": checks,
        }
