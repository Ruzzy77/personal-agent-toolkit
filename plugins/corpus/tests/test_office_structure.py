from __future__ import annotations

import sys
import tempfile
import unittest
from pathlib import Path
from unittest import mock

from corpus.adapters import run_builtin_extraction


class OfficeStructureTest(unittest.TestCase):
    @unittest.skipUnless(sys.platform == "darwin", "Vision requires macOS")
    def test_local_image_ocr_has_source_object_location_and_keeps_native_text(self):
        from corpus.adapter_registry import build_default_registry
        from PIL import Image, ImageDraw, ImageFont
        from pptx import Presentation
        from pptx.util import Inches

        with tempfile.TemporaryDirectory() as temporary:
            base = Path(temporary)
            picture = base / "image.png"
            image = Image.new("RGB", (1000, 250), "white")
            font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 70)
            ImageDraw.Draw(image).text(
                (50, 50), "Corpus local OCR", font=font, fill="black"
            )
            image.save(picture)
            # Reusing a large package part must not spend its byte budget again.
            with picture.open("ab") as stream:
                stream.write(b"\0" * (2 * 1024 * 1024))
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            slide.shapes.add_picture(str(picture), 0, 0, width=Inches(6))
            second = presentation.slides.add_slide(presentation.slide_layouts[6])
            second.shapes.add_textbox(
                0, 0, Inches(6), Inches(1)
            ).text = "Existing native content remains authoritative and is not replaced"
            second.shapes.add_picture(str(picture), 0, 0, width=Inches(6))
            for _ in range(17):
                repeated = presentation.slides.add_slide(presentation.slide_layouts[6])
                repeated.shapes.add_picture(str(picture), 0, 0, width=Inches(6))
            path = base / "images.pptx"
            presentation.save(path)
            original = path.read_bytes()
            adapter = build_default_registry(base / "runtime").resolve("pptx")
            with mock.patch.object(
                adapter, "_image_text", wraps=adapter._image_text
            ) as recognize:
                result = adapter.extract(path, format_id="pptx")
                self.assertEqual(recognize.call_count, 1)
            self.assertEqual(path.read_bytes(), original)
        recognized = [u for u in result.units if u.derivation_method == "ocr"]
        self.assertTrue(recognized)
        self.assertIn("Corpus local OCR", " ".join(u.content for u in recognized))
        self.assertEqual(
            {u.structure_path["slide"] for u in recognized}, {1, *range(3, 20)}
        )
        first_ocr = next(
            i for i, u in enumerate(result.units) if u.derivation_method == "ocr"
        )
        native_position = next(
            i
            for i, u in enumerate(result.units)
            if u.content.startswith("Existing native content")
        )
        self.assertLess(first_ocr, native_position)
        self.assertNotIn("office_image_size_limit", {i.code for i in result.issues})
        self.assertTrue(
            all(u.geometry and u.confidence is not None for u in recognized)
        )
        self.assertTrue(
            all(
                u.structure_path["image_part"].startswith("ppt/media/")
                for u in recognized
            )
        )
        self.assertEqual(
            sum(u.content.startswith("Existing native content") for u in result.units),
            1,
        )

    def test_word_body_order_nested_tables_and_merged_cells(self):
        from docx import Document

        with tempfile.TemporaryDirectory() as temporary:
            path = Path(temporary) / "source.docx"
            document = Document()
            document.add_heading("Title", level=1)
            document.add_paragraph("Before the table")
            table = document.add_table(rows=3, cols=2)
            table.cell(0, 0).merge(table.cell(0, 1)).text = "Horizontal merge"
            table.cell(1, 0).merge(table.cell(2, 0)).text = "Vertical merge"
            table.cell(1, 1).text = "Before nested"
            nested = table.cell(1, 1).add_table(rows=1, cols=1)
            nested.cell(0, 0).text = "Nested cell"
            table.cell(1, 1).add_paragraph("After nested")
            table.cell(2, 1).text = "Last cell"
            document.add_paragraph("After the table")
            document.sections[0].header.paragraphs[0].text = "Header text"
            document.sections[0].footer.paragraphs[0].text = "Footer text"
            document.save(path)
            original = path.read_bytes()
            result = run_builtin_extraction(path, "docx")
            self.assertEqual(original, path.read_bytes())
        texts = [u.content for u in result.units if u.content]
        self.assertEqual(
            texts,
            [
                "Title",
                "Before the table",
                "Horizontal merge",
                "Vertical merge",
                "Before nested",
                "Nested cell",
                "After nested",
                "Last cell",
                "After the table",
                "Header text",
                "Footer text",
            ],
        )
        vertical = next(u for u in result.units if u.content == "Vertical merge")
        self.assertEqual(vertical.structure_path["row_span"], 2)
        horizontal = next(u for u in result.units if u.content == "Horizontal merge")
        self.assertEqual(horizontal.structure_path["col_span"], 2)
        nested = next(u for u in result.units if u.content == "Nested cell")
        self.assertTrue(nested.structure_path["container_path"])
        self.assertEqual(result.units[0].unit_type, "heading")

    def test_word_textbox_and_field_do_not_duplicate_or_evaluate(self):
        from docx import Document
        from docx.oxml import parse_xml

        with tempfile.TemporaryDirectory() as temporary:
            path = Path(temporary) / "source.docx"
            document = Document()
            paragraph = document.add_paragraph("Before ")
            paragraph._p.append(
                parse_xml("""
              <w:r xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"
                   xmlns:mc="http://schemas.openxmlformats.org/markup-compatibility/2006">
                <mc:AlternateContent>
                  <mc:Choice Requires="w"><w:drawing><w:txbxContent><w:p><w:r>
                    <w:t>Textbox</w:t></w:r></w:p></w:txbxContent></w:drawing></mc:Choice>
                  <mc:Fallback><w:pict><w:txbxContent><w:p><w:r>
                    <w:t>Textbox</w:t></w:r></w:p></w:txbxContent></w:pict></mc:Fallback>
                </mc:AlternateContent>
              </w:r>""")
            )
            paragraph.add_run("After ")
            paragraph._p.append(
                parse_xml("""
              <w:fldSimple xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"
                           w:instr="DATE"><w:r><w:t>Stored date</w:t></w:r></w:fldSimple>""")
            )
            document.save(path)
            result = run_builtin_extraction(path, "docx")
        self.assertEqual(
            [u.content for u in result.units if u.content],
            ["Before", "Textbox", "After", "Stored date"],
        )
        textbox = next(u for u in result.units if u.content == "Textbox")
        self.assertIn("owner_paragraph", textbox.structure_path)
        field = next(u for u in result.units if u.unit_type == "field")
        self.assertEqual(field.structure_path["instruction"], "DATE")
        self.assertEqual(field.structure_path["evaluation"], "stored_result_only")

    def test_slides_recursive_shapes_tables_chart_and_notes(self):
        from pptx import Presentation
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE
        from pptx.util import Inches

        with tempfile.TemporaryDirectory() as temporary:
            path = Path(temporary) / "source.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            slide.shapes.add_textbox(0, 0, Inches(2), Inches(1)).text = "Top-level"
            group = slide.shapes.add_group_shape()
            group.shapes.add_textbox(0, 0, Inches(2), Inches(1)).text = "Group text"
            nested = group.shapes.add_group_shape()
            nested.shapes.add_textbox(0, 0, Inches(2), Inches(1)).text = "Nested group"
            table = slide.shapes.add_table(2, 2, 0, 0, Inches(4), Inches(2)).table
            table.cell(0, 0).merge(table.cell(0, 1))
            table.cell(0, 0).text = "Merged heading"
            table.cell(1, 0).text = "Cell A"
            table.cell(1, 1).text = "Cell B"
            data = CategoryChartData()
            data.categories = ["North", "South"]
            data.add_series("Revenue", [12, 18])
            slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED, 0, 0, Inches(4), Inches(2), data
            )
            slide.notes_slide.notes_text_frame.text = "Speaker notes"
            presentation.save(path)
            original = path.read_bytes()
            first = run_builtin_extraction(path, "pptx")
            second = run_builtin_extraction(path, "pptx")
            self.assertEqual(path.read_bytes(), original)
        self.assertEqual(first.manifest_hash, second.manifest_hash)
        texts = [u.content for u in first.units if u.content]
        self.assertEqual(
            texts[:6],
            [
                "Top-level",
                "Group text",
                "Nested group",
                "Merged heading",
                "Cell A",
                "Cell B",
            ],
        )
        self.assertEqual(texts.count("Merged heading"), 1)
        self.assertEqual(texts[-1], "Speaker notes")
        self.assertTrue({"Revenue", "North", "South", "12", "18"}.issubset(texts))
        group_text = next(u for u in first.units if u.content == "Nested group")
        self.assertEqual(len(group_text.structure_path["group_path"]), 2)
        cell = next(u for u in first.units if u.content == "Merged heading")
        self.assertEqual(cell.structure_path["col_span"], 2)
        locators = [dict(u.structure_path) for u in first.units]
        self.assertTrue(all("part" in loc and "element" in loc for loc in locators))


if __name__ == "__main__":
    unittest.main()
